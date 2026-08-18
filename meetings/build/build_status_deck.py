"""Build the GENERAL STATUS deck for the 2026-08-16 supervision meeting.

The whole project, not one paper. `build_paper_deck.py` already writes into the same dated folder
and covers the CLPsych draft in depth; this deck is the wider one — where every arm stands, what
the main comparison says, how much of it survives a second grader, the NEW look-ahead result, and
the decisions that follow. It should stand on its own if nobody in the room has read the drafts.

Written for a ROOM: no "since the email", no "as discussed". Every slide states its own claim.

What is new since the 2026-08-03 deck, and why this is not just that deck re-dated:

- **The look-ahead comparison resolved, and not the way the earlier reading had it.** The PTO K=5
  arm is now trained 1–10 and scored 0–10 on both graders, so the matched K comparison runs over
  eleven points. K=5 closes the over-praise channel for the whole run (large, replicated), and the
  MI-inconsistent TOTAL falls only under the grader that was the training reward (dz 0.446) — the
  held-out judge sees no aggregate change (dz 0.099, n.s.). The cross-judge claim is SUBSTITUTION,
  not reduction. Three slides: the channel, the aggregate, and where the intervention acts.
- **A correction slide.** Two prose readings in the repo — one in STATUS.md, one written earlier
  the same day — were right about a component and wrong about the behaviour. Naming that on a
  slide is the point; a room that hears only the corrected version learns the wrong lesson about
  how much to trust the next single-channel result.
- **Two drafts now exist**, on disjoint claims, and the scope question is live.

Numbers are owned by the artifacts they are read from — `eda/results/{L0,L5}/` tables and each
paper's `NUMBERS.md`. This script only restates them. If a number here disagrees with those, they
are right.

Exports .pptx, then convert to .pdf via export_pdf.ps1 (PowerPoint COM).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))       # meetings/build/
REPO = os.path.dirname(os.path.dirname(HERE))           # repo root (meetings/ lives here)
ROOT = os.path.join(REPO, "Exp3_PTO_GRPO")             # the experiment the artifacts come from
L0F  = os.path.join(ROOT, "eda", "results", "L0", "figures")
L0T  = os.path.join(ROOT, "eda", "results", "L0", "tables")
L5F  = os.path.join(ROOT, "eda", "results", "L5", "figures")
L5T  = os.path.join(ROOT, "eda", "results", "L5", "tables")
METH = os.path.join(ROOT, "figures")   # hand-authored method schematics (build_method_figures.py)
OUT  = os.path.join(REPO, "meetings", "2026-08-16", "status_2026-08-16.pptx")

JUDGE = "gpt-4o-mini"
JUDGE_INVARIANT_FAMILIES = {"8_measurement"}

def _jp(base, p):
    *parts, name = p.split("/")
    if parts and parts[0] in JUDGE_INVARIANT_FAMILIES:
        return os.path.join(base, *parts, name)
    return os.path.join(base, *parts, JUDGE, name)

def f0(p): return _jp(L0F, p)
def t0(p): return _jp(L0T, p)
def f5(p): return _jp(L5F, p)
def t5(p): return _jp(L5T, p)
def fm(p): return os.path.join(METH, p)

NAVY  = RGBColor(0x1F,0x3A,0x5F); PTO = RGBColor(0x00,0x72,0xB2)
GRPO  = RGBColor(0xE6,0x9F,0x00); GREY= RGBColor(0x5A,0x5A,0x5A)
LIGHT = RGBColor(0xEF,0xF2,0xF6); GREEN=RGBColor(0x00,0x8A,0x63)
RED   = RGBColor(0xC0,0x4A,0x1A); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x22,0x28,0x30); ROWALT=RGBColor(0xF4,0xF7,0xFA)
AMBER = RGBColor(0xB0,0x77,0x00)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_N = 0

def slide(fnum=True):
    global _N
    s = prs.slides.add_slide(BLANK)
    if fnum:
        _N += 1
        tb = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"Exp3 · {_N}"; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xB0,0xB8,0xC0); r.font.name = FONT
    return s

def box(s,l,tp,w,h):
    tb = s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tb.text_frame.word_wrap = True; return tb

def rect(s,l,tp,w,h,color,line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp

def set_runs(para, segs, size, align=PP_ALIGN.LEFT):
    para.alignment = align
    for txt,bold,color in segs:
        r = para.add_run(); r.text = txt; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

def title_bar(s, text, kicker=None):
    rect(s,0,0,13.333,1.15,NAVY); rect(s,0,1.15,13.333,0.06,PTO)
    tb = box(s,0.55,0.18,12.2,0.9); tf = tb.text_frame
    p = tf.paragraphs[0]
    if kicker:
        set_runs(p,[(kicker,True,RGBColor(0x8F,0xC7,0xEC))],13)
        set_runs(tf.add_paragraph(),[(text,True,WHITE)],26)
    else:
        set_runs(p,[(text,True,WHITE)],27)

def bullets(s, items, left=0.6, top=1.5, width=12.1, size=16, gap=8):
    tb = box(s,left,top,width,5.6); tf = tb.text_frame
    first = True
    for segs,level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level = level; p.space_before = Pt(gap); p.space_after = Pt(2)
        mk = "•  " if level==0 else "–  "
        set_runs(p,[("    "*level+mk, False, PTO if level==0 else GREY)]+segs, size)

def fit(img,mw,mh):
    with Image.open(img) as im: iw,ih = im.size
    ar = iw/ih; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    return w,h

def figure(s,img,mw,mh,left,top,anchor="center"):
    w,h = fit(img,mw,mh)
    dy = 0.0 if anchor=="top" else (mh-h)/2
    s.shapes.add_picture(img,Inches(left+(mw-w)/2),Inches(top+dy),Inches(w),Inches(h))
    return top+dy+h

def caption(s,text,left,top,width,color=GREY,size=11.5,align=PP_ALIGN.CENTER):
    p = box(s,left,top,width,0.5).text_frame.paragraphs[0]
    set_runs(p,[(text,False,color)],size,align)

def side_notes(s, heading, rows, left, top, width, size=12, gap=10):
    p = box(s,left,top,width,0.45).text_frame.paragraphs[0]
    set_runs(p,[(heading,True,NAVY)],14.5)
    tb = box(s,left,top+0.45,width,5.2); tf = tb.text_frame
    first = True
    for segs in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(gap)
        set_runs(p,[("•  ",False,PTO)]+segs,size)

def cards(s, items, top, left=0.6, width=12.1, height=1.2, gap=0.12):
    y = top
    for name,desc,col,tag in items:
        rect(s,left,y,width,height,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,left,y,0.14,height,col)
        tb = box(s,left+0.35,y+0.10,width-0.55,height); tf = tb.text_frame
        seg = [(name,True,col)]
        if tag: seg.append(("    ("+tag+")",True,GREEN))
        set_runs(tf.paragraphs[0],seg,14.5)
        set_runs(tf.add_paragraph(),[(desc,False,DARK)],12)
        y += height + gap
    return y

def grid_table(s, rows, left, top, width, colw, fontsize=12.5, rowh=0.42):
    y = top
    for ri,(cells, color, is_hdr) in enumerate(rows):
        x = left
        for i,c in enumerate(cells):
            w = width*colw[i]
            rect(s,x,y,w,rowh, NAVY if is_hdr else (LIGHT if ri%2==0 else WHITE))
            p = box(s,x+0.1,y+0.04,w-0.15,rowh).text_frame.paragraphs[0]
            col = WHITE if is_hdr else (color if (color and i==len(cells)-1) else DARK)
            set_runs(p,[(c,is_hdr or i==0,col)],fontsize)
            x += w
        y += rowh + 0.04
    return y

def callout(s, left, top, width, height, segs, accent=PTO, fill=LIGHT, size=13):
    rect(s,left,top,width,height,fill); rect(s,left,top,0.06,height,accent)
    set_runs(box(s,left+0.25,top+0.10,width-0.45,height-0.1).text_frame.paragraphs[0], segs, size)

# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.5,13.333,0.08,PTO)
tb = box(s,0.9,1.5,11.5,2.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Looking Ahead in Goal-Oriented Dialogue",True,WHITE)],34)
set_runs(tf.add_paragraph(),[("Preference-Tree (PTO) vs Group-Relative (GRPO) optimization "
                             "of a small Motivational-Interviewing therapist",False,
                             RGBColor(0xBD,0xD6,0xEA))],18)
tb2 = box(s,0.9,4.75,11.5,2.2); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("Project status — all four arms, both graders, and a new "
                             "look-ahead result",True,WHITE)],20)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Both comparisons are finished. The look-ahead arm reached iteration 10 this week "
             "and resolved — not the way we last read it.",False,RGBColor(0x9F,0xB4,0xC8))],14)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_runs(p,[("Lior Baruch · Reichman University · 16 August 2026",
             False,RGBColor(0x9F,0xB4,0xC8))],13)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Llama-3.2-1B therapist (bf16)  ·  gpt-4o-mini simulated patient + oracle  ·  "
             "Claude Haiku 4.5 held-out judge  ·  96 personas, persona-paired",
             False,RGBColor(0x7E,0x93,0xA8))],12)

# =====================================================================
# 2 · AGENDA
# =====================================================================
s = slide(); title_bar(s,"Where the project stands, and what I'd like to decide","AGENDA")

p = box(s,0.6,1.45,6.0,0.4).text_frame.paragraphs[0]
set_runs(p,[("Settled",True,GREEN)],16.5)
bullets(s,[
  ([("The main comparison is finished. ",True,NAVY),("PTO and GRPO, 10 iterations each at matched "
    "settings, every conversation scored on eight instruments by two independent graders.",
    False,DARK)],0),
  ([("The instrument is measured, not assumed. ",True,NAVY),("Oracle repeatability, a full "
    "re-scoring by a judge from another model family that never played the patient, and an "
    "explicit list of where the two graders disagree.",False,DARK)],0),
  ([("The look-ahead arm is finished too. ",True,NAVY),("PTO K=5 now trains to iteration 10 and is "
    "scored 0–10 on both graders. It closes the reward-hack channel it targets for the whole run, "
    "and the therapist substitutes another violation for the one it lost. That is new this week.",
    False,DARK)],0),
  ([("Two paper drafts exist",True,NAVY),(", on deliberately disjoint claims, both compiling.",
    False,DARK)],0),
], left=0.6, top=1.95, width=6.0, size=13.5, gap=11)

rect(s,6.85,1.45,0.02,4.4,RGBColor(0xD5,0xDD,0xE5))
p = box(s,7.15,1.45,5.5,0.4).text_frame.paragraphs[0]
set_runs(p,[("Open — the agenda",True,RED)],16.5)
bullets(s,[
  ([("Scope. ",True,NAVY),("Two drafts, one thesis. Which is the chapter, which is the paper, and "
    "in what order do they get finished?",False,DARK)],0),
  ([("Budget. ",True,NAVY),("Two things are left to buy: GRPO at K=5, or nothing. Only the first "
    "changes a claim.",False,DARK)],0),
  ([("Human MI coders. ",True,NAVY),("The one validity gap that money alone does not close, and "
    "the one a reviewer will ask about first.",False,DARK)],0),
  ([("Authorship and venue",True,NAVY),(" for both drafts.",False,DARK)],0),
], left=7.15, top=1.95, width=5.5, size=14.5, gap=18)

callout(s,0.6,6.1,12.1,0.9,
  [("In one line: ",True,NAVY),("the therapist learns to flatter; a principled fix stops the "
    "flattery and the therapist starts lecturing instead, with no fall in total MI violations that "
    "a grader outside the loop can see. The reward function cannot tell any of these apart.",
    False,DARK)],size=13.5)

# =====================================================================
# 3 · BACKGROUND — the published paper
# =====================================================================
s = slide(); title_bar(s,"The starting point — the ICLR 2025 paper","BACKGROUND · 1")

rect(s,0.6,1.45,12.1,1.12,LIGHT)
tb = box(s,0.85,1.55,11.6,1.0); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Preference Tree Optimization: Enhancing Goal-Oriented Dialogue with "
                           "Look-Ahead Simulations",True,NAVY)],17)
p = tf.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Baruch, Butman, Bar, Friedman · ICLR 2025 (SSI-FM workshop poster)",False,GREY)],13)

bullets(s,[
  ([("The method. ",True,NAVY),("At each therapist turn, branch several candidate replies; for "
    "each, simulate ",False,DARK),("K further turns",True,DARK),("; let the oracle score the "
    "resulting trajectory; keep best and worst as a preference pair; update with DPO. Repeat, "
    "regenerating from the improved model.",False,DARK)],0),
  ([("Why look ahead. ",True,NAVY),("Scoring a reply on its own rewards openings that look good in "
    "isolation. Scoring the reply plus K simulated turns rewards openings that lead somewhere.",
    False,DARK)],0),
  ([("Headline finding. ",True,NAVY),("Every PTO model beat the untrained baseline, and ",
    False,DARK),("K = 5 scored higher and more stably than K = 0",True,DARK),
   (" — on a 7B therapist over 7 iterations.",False,DARK)],0),
  ([("Why that matters today. ",True,AMBER),("This experiment does not reproduce that on a 1B "
    "model over 10 iterations, and the reason turns out to be interesting rather than a "
    "contradiction. Slides 11–14.",False,DARK)],0),
], top=2.85, size=15, gap=15)

# =====================================================================
# 4 · WHAT WAS RUN — the four arms
# =====================================================================
s = slide(); title_bar(s,"What was run, and what is scored","SETUP")

grid_table(s,[
  (["Arm","Trained","Scored on BOTH graders","Status"],None,True),
  (["PTO  K=0","iterations 1–10","base + 1–10","complete"],GREEN,False),
  (["GRPO K=0","iterations 1–10","base + 1–10","complete"],GREEN,False),
  (["PTO  K=5","iterations 1–10","base + 1–10","complete"],GREEN,False),
  (["GRPO K=5","iteration 1","base + 1","thin — stopped on budget"],RED,False),
],0.5,1.45,12.3,[0.13,0.17,0.24,0.46],fontsize=12.5,rowh=0.40)

bullets(s,[
  ([("35 scored model states",True,NAVY),(", each 96 conversations, each conversation scored on "
    "eight instruments by two graders.",False,DARK)],0),
  ([("Matched by construction ",True,NAVY),("across every arm: sampling temperature 1.2, 8 "
    "candidates per branch point, minimum context length 12, the same τ filter, the same DPO β, "
    "the same oracle. Verified against each run's stored metadata, not assumed.",False,DARK)],0),
  ([("The same 96 patient personas recur in every iteration",True,NAVY),(" — reshuffled, so "
    "identity has to be recovered rather than read off file order. That is what makes every "
    "comparison in this deck ",False,DARK),("persona-paired",True,DARK),
   (", which is far stronger than comparing group means.",False,DARK)],0),
  ([("Training reward = two rubrics ",True,NAVY),("(session satisfaction + working alliance). "
    "The other six instruments were never optimised against, which is what makes them evidence.",
    False,DARK)],0),
], top=3.65, size=13.5, gap=13)

# =====================================================================
# 5 · METHOD — the two loops
# =====================================================================
s = slide(); title_bar(s,"The two methods — same loop, one step different","METHOD")
figure(s,fm("pto_framework.png"),6.1,4.5,0.35,1.45,anchor="top")
figure(s,fm("grpo_framework.png"),6.1,4.5,6.85,1.45,anchor="top")
caption(s,"PTO — branch, score, keep best vs worst, DPO",0.35,6.15,6.1,color=PTO,size=12.5)
caption(s,"GRPO — sample a group, standardize rewards, policy gradient",6.85,6.15,6.1,
        color=GRPO,size=12.5)
callout(s,0.6,6.6,12.1,0.72,
  [("Both regenerate their own training data from the current policy each iteration. ",True,NAVY),
   ("Look-ahead (K) changes only what CONTEXT the oracle grades — never the loss, the "
    "hyperparameters, or the sampling.",False,DARK)],size=12.5)

# =====================================================================
# 6 · MAIN RESULT — trajectories
# =====================================================================
s = slide(); title_bar(s,"All eight instruments across 10 iterations (K = 0)","RESULTS · 1")
bot = figure(s,f0("0_headline/trajectories_all_metrics.png"),12.3,5.2,0.5,1.35,anchor="top")
caption(s,"Mean ± 95% CI over the 96 personas, primary oracle. MICI is lower-is-better. "
          "PTO climbs steadily to iteration 10; GRPO peaks at iteration 8 and regresses.",
        0.5,bot+0.05,12.3,size=11.5)

# =====================================================================
# 7 · ENDPOINT NUMBERS
# =====================================================================
s = slide(); title_bar(s,"Endpoint numbers (K = 0)","RESULTS · 2")
caption(s,"Primary oracle, final and best endpoint. Read off the generated scorecard "
          "(results/L0/tables/0_headline/gpt-4o-mini/) — every cell traces to the score lake.",
        0.6,1.3,12.1,size=12,align=PP_ALIGN.LEFT)
grid_table(s,[
  (["","Q1+Q2","MITI","MICI ↓","R:Q","% MI-consistent"],None,True),
  (["PTO  K=0 @ iter 10","4.26","4.27","0.49","0.75","0.70"],None,False),
  (["GRPO K=0 @ iter 10","3.75","3.92","0.84","1.44","0.83"],None,False),
  (["GRPO K=0 @ iter 8 (its peak)","4.08","4.23","0.54","1.04","0.71"],None,False),
],0.6,1.85,12.1,[0.30,0.14,0.14,0.14,0.14,0.14],fontsize=13,rowh=0.44)

bullets(s,[
  ([("PTO beats GRPO at the matched endpoint",True,NAVY),(" — Q1+Q2 4.26 vs 3.75, persona-paired "
    "+0.51, dz 0.73.",False,DARK)],0),
  ([("Credit GRPO at its peak instead",True,NAVY),(" and PTO still leads (+0.18, dz 0.30, "
    "p = .010) — but the MITI and MICI gaps stop being significant. ",False,DARK),
   ("We claim the weaker version.",True,GREEN)],0),
  ([("GRPO's higher R:Q is not a win. ",True,AMBER),("Its reflection-to-question ratio crosses the "
    "MITI competency bar because it stopped asking questions — the regex question rate falls "
    "0.83 → 0.15 per turn. A ratio is gameable through its denominator.",False,DARK)],0),
  ([("The held-out judge orders them the same way",True,GREEN),(" on Q1+Q2 at iteration 10 — "
    "2.87 (PTO) vs 2.26 (GRPO) — on a different absolute scale. ",False,DARK),
   ("Never average the two graders' raw scores",True,RED),(": the primary one WAS the training "
    "reward, so this is train-vs-test, and the level offset is model-dependent. Only contrasts "
    "combine.",False,DARK)],0),
], top=3.9, size=13.5, gap=13)

# =====================================================================
# 8 · WHAT ELSE MOVED — the reward hack
# =====================================================================
s = slide(); title_bar(s,"What the model actually learned (K = 0)","RESULTS · 3")
bot = figure(s,f0("0_headline/mici_detail_grid.png"),7.9,4.6,0.35,1.4,anchor="top")
side_notes(s,"One channel carries the whole thing",[
  [("Every ",False,DARK),("coercive",True,DARK),(" MI violation FALLS with training — confronting, "
    "warning, directing and judging all go to ~0.",False,DARK)],
  [("Over-praise rises from ",False,DARK),("0.02 to 0.70 per turn",True,RED),(" in GRPO and 0.01 "
    "to 0.30 in PTO, and accounts for more than 100% of the increase in total MI-inconsistency.",
    False,DARK)],
  [("The rubrics measure how the session ",False,DARK),("felt to the client",True,DARK),
   (". Flattery is the cheapest way to move that.",False,DARK)],
  [("So the model did not get broadly worse at MI. It swapped coercion for flattery.",False,DARK)],
],8.55,1.35,4.4,size=12,gap=11)
caption(s,"MI-inconsistent behaviour, decomposed, per therapist turn. Higher is worse.",
        0.35,bot+0.05,7.9,size=11)

# =====================================================================
# 9 · MEASUREMENT VALIDITY
# =====================================================================
s = slide(); title_bar(s,"How much of this survives a grader that took no part in training",
                       "MEASUREMENT VALIDITY")
bot = figure(s,f0("8_measurement/multijudge_gain_retention.png"),7.7,4.5,0.35,1.4,anchor="top")
side_notes(s,"Gain retention = a train/test split on the reward",[
  [("The primary oracle ",False,DARK),("was the training signal",True,DARK),("; Claude Haiku 4.5 "
    "is held out, from another family, and never played the patient. Agreement between them is a "
    "generalisation check, not inter-rater reliability.",False,DARK)],
  [("Oracle repeatability ICC(2,1) ",False,DARK),("0.86–0.99",True,GREEN),
   ("; 1,632 of 1,848 arm × metric contrasts keep their sign (98.9% at |Δ| ≥ 0.50).",False,DARK)],
  [("Q1 retention: PTO@10 ",False,DARK),("0.80",True,GREEN),(" vs GRPO@10 ",False,DARK),
   ("0.28",True,RED),(", non-overlapping — GRPO's late gains are increasingly grader-specific.",
    False,DARK)],
  [("⚠ Two standing caveats: MITI dependability is ",False,DARK),("0.65",True,RED),
   (" off one judge, and MICI agreement is weak at the conversation level. Both are limitations "
    "in the write-up, not footnotes.",False,DARK)],
],8.35,1.35,4.6,size=11.5,gap=9)
caption(s,"Fraction of each arm's gain over base that survives the judge swap. ~1.0 = a real "
          "behaviour change; ~0 = a gain that existed only in the optimised grader.",
        0.35,bot+0.05,7.7,size=11)

# =====================================================================
# 10 · SECTION BREAK — the new result
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,3.95,13.333,0.08,PTO)
tb = box(s,1.1,2.5,11.2,1.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("New this week",True,RGBColor(0x8F,0xC7,0xEC))],16)
set_runs(tf.add_paragraph(),[("Does look-ahead fix the reward hack?",True,WHITE)],36)
tb2 = box(s,1.1,4.3,11.2,1.6); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("It closes the channel it targets. The therapist substitutes another "
                             "violation for the one it lost.",False,RGBColor(0xBD,0xD6,0xEA))],19)
p = tf2.add_paragraph(); p.space_before = Pt(12)
set_runs(p,[("PTO only, K ∈ {0,5}, eleven matched iterations (0–10), both graders, persona-paired.",
             False,RGBColor(0x7E,0x93,0xA8))],13)

# =====================================================================
# 11 · LOOK-AHEAD 1 — the channel closes
# =====================================================================
s = slide(); title_bar(s,"Look-ahead closes the channel it targets, and keeps it closed",
                       "LOOK-AHEAD · 1")
bot = figure(s,f5("7_stats/k_overpraise_trajectory.png"),7.4,4.5,0.35,1.4,anchor="top")
side_notes(s,"Exactly what it was supposed to do",[
  [("The two PTO arms are ",False,DARK),("indistinguishable for six iterations",True,DARK),
   (" and then separate for good. Over-praise per session at iteration 10: ",False,DARK),
   ("3.042 (K=0) vs 0.625 (K=5)",True,DARK),(".",False,DARK)],
  [("Persona-paired, n = 96: ",False,DARK),("dz 0.887, p < 10⁻⁴",True,GREEN),
   (". The held-out judge reports a larger gap in the same direction — 4.750 vs 1.177, dz 0.999.",
    False,DARK)],
  [("Prevention, not delay. ",True,GREEN),("K=5's over-praise stays LOW to the end of the run "
    "— it rises 0.115 → 0.625 across the eleven points, while K=0 goes 0.167 → 3.042. Low and "
    "still creeping, not static; the draft's open question is closed in the direction of "
    "prevention.",False,DARK)],
  [("Not a denominator artifact — these are ",False,DARK),("raw per-session counts",True,DARK),
   (". K=5 actually takes MORE therapist turns (14.4 vs 10.2), which would push a rate the other "
    "way.",False,DARK)],
  [("Not a data-starvation artifact — over the full ten iterations K=5 trained on ",False,DARK),
   ("more",True,DARK),(" preference groups (6,416 vs 4,935).",False,DARK)],
  [("And it is visible ",False,DARK),("upstream",True,DARK),(", in what the reward itself selects "
    "for — not only in the output. Slide 13.",False,DARK)],
],8.0,1.35,4.95,size=11.5,gap=7)
caption(s,"Over-praise per therapist turn, all four arms. K=0 solid, K=5 dashed. Shading marks "
          "the measured divergence onset. Higher is worse.",0.35,bot+0.03,7.4,size=11)

# =====================================================================
# 12 · LOOK-AHEAD 2 — the substitution replicates, the reduction does not  (THE SLIDE)
# =====================================================================
s = slide(); title_bar(s,"…and the substitution replicates — the reduction does not",
                       "LOOK-AHEAD · 2")

grid_table(s,[
  (["per session, iteration 10","Δ (K=0 − K=5)","dz","p (Holm)"],None,True),
  (["Over-praise","+2.417","+0.887","< 10⁻⁴"],RED,False),
  (["Advice without permission","−0.313","−0.239",".057  n.s."],GREEN,False),
  (["Directing the client","−0.375","−0.389",".0016"],GREEN,False),
  (["ALL MI-inconsistent acts","+1.615","+0.446",".0003"],AMBER,False),
],0.5,1.4,6.6,[0.44,0.22,0.15,0.19],fontsize=12,rowh=0.40)

side_notes(s,"The therapist changed how it fails; whether it fails less depends on who grades",[
  [("Table above is the ",False,DARK),("primary oracle",True,DARK),(", which WAS the training "
    "reward: totals 4.958 (K=0) vs 3.344 (K=5) acts per session.",False,DARK)],
  [("The held-out judge reproduces every component and ",False,DARK),("larger",True,DARK),
   (" — over-praise dz 0.999, advice −0.709, directing −0.468 — but its total is ",False,DARK),
   ("null: +0.531, dz 0.099, p = .167",True,RED),(" (8.510 vs 7.979).",False,DARK)],
  [("So the claim that travels across graders is ",False,DARK),
   ("substitution, not reduction",True,GREEN),(". The significant aggregate is a "
    "primary-oracle-only result and is labelled as one wherever it appears.",False,DARK)],
  [("Composition flips: over-praise is ",False,DARK),("61.3%",True,RED),
   (" of K=0's violations and ",False,DARK),("18.7%",True,DARK),(" of K=5's; "
    "advice-without-permission is 32.1% and ",False,DARK),("57.0%",True,RED),(".",False,DARK)],
  [("Counts, not rates. ",True,NAVY),("The per-TURN MI-inconsistency rate is large under both "
    "judges (dz 0.708 / 0.655) only because K=5 takes 14.4 therapist turns to K=0's 10.2. "
    "Per-session counts are the primary reading.",False,DARK)],
  [("And the reward the run optimised still cannot separate the arms: Q1+Q2 dz −0.096, p_holm .695 — "
    "K=5 nominally ahead for the first time, not significantly; the held-out judge has K=0 ahead "
    "(dz 0.308, p = .032).",False,DARK)],
],7.35,1.35,5.55,size=11,gap=6)

# The table above runs 1.40 -> ~3.62 (5 rows x 0.44). Start the figure clear of it, or the
# "ALL MI-inconsistent acts" row — the one row this slide exists for — gets covered.
bot = figure(s,f5("7_stats/k_mici_composition_grid.png"),6.6,3.05,0.5,3.80,anchor="top")
caption(s,"Per-session counts — no denominator involved.",0.5,bot+0.02,6.6,size=10.5)

# =====================================================================
# 13 · LOOK-AHEAD 3 — where it acts
# =====================================================================
s = slide(); title_bar(s,"Where the intervention acts — and why the pressure relocates",
                       "LOOK-AHEAD · 3")
figure(s,f5("6_preference/k_mechanism_overpraise.png"),4.5,5.4,0.4,1.35,anchor="top")
# Selection-pressure figures RE-READ 2026-08-17 from 6_preference.xlsx sheet
# `k_mechanism_overpraise_chain` (policy axis). PTO_LA0 w_overpraise peaks at 0.083 (policy
# iter 8) and is NEGATIVE at iter 9 (-0.039); PTO_LA5's max is 0.025. The earlier 0.086 on this
# slide came from L0's update_lexical_push, which is a GRPO AFFIRMATION weight — a different
# table and a different feature. Do not reintroduce it.
side_notes(s,"A behaviour has to clear three gates",[
  [("Every training group logs all its candidates and their oracle scores, so we can measure what "
    "the reward ",False,DARK),("selected for",True,DARK),(" — not just what the policy produced.",
    False,DARK)],
  [("Under K=0 the selection weight on over-praise is ~0 for ",False,DARK),("six",True,DARK),
   (" iterations, then jumps — 0.063, 0.034, ",False,DARK),("0.083",True,RED),
   (" at policy iterations 6–8 (SEs 0.017–0.025), before falling back. The candidate pool and "
    "the evaluation follow with a lag.",False,DARK)],
  [("Under K=5 it ",False,DARK),("never clears +0.025",True,GREEN),(" at any iteration, and "
    "downstream stays an order smaller rather than static — the pool moves 0.004 → 0.065 and "
    "the evaluated rate 0.008 → 0.043. Look-ahead acted on the ",False,DARK),("reward",True,DARK),
   (", not on the policy.",False,DARK)],
  [("Which is also why the substitution is unsurprising. The objective did not change; only the "
    "price of one route went up. A trajectory-extended oracle still rewards a turn that makes the "
    "next five look productive — and ",False,DARK),("giving advice does that too",True,AMBER),
   (". Nothing penalises advice; the intervention penalises praise that does not compound.",
    False,DARK)],
  [("Scale worth holding together: the largest per-iteration selection pressure anywhere in "
    "this chain is ",False,DARK),("0.083",True,DARK),(", while the generated over-praise rate "
    "moves ",False,DARK),("0.002 → 0.318",True,RED),(". A bias too small to see in one update "
    "compounds through the on-policy loop.",False,DARK)],
],5.35,1.35,7.55,size=11.5,gap=8)

# =====================================================================
# 14 · THE CORRECTION SLIDE
# =====================================================================
s = slide(); title_bar(s,"Two readings this corrected — and why I am putting them on a slide",
                       "HOW WE GOT IT WRONG TWICE")

cards(s,[
  ("STATUS.md: \"drop the look-ahead MI-consistency claim\"",
   "Correct that TOTAL MICI flips sign across iterations 7–8 — because the total mixes channels "
   "that move in opposite directions. The over-praise component does not flip; it is monotone. "
   "The claim was dropped for a property of the summary, not of the data.",
   AMBER,"a summary hid its components"),
  ("Earlier today: \"look-ahead cuts the reward hack ~4×\"",
   "True of over-praise, replicated by a second judge, mechanistically explained at three levels "
   "— and wrong about the therapist's behaviour. It was scored on the channel that motivated the "
   "fix. The aggregate had not been looked at.",
   RED,"a channel scored as if it were the outcome"),
],1.45,left=0.6,width=12.1,height=1.12,gap=0.18)

bullets(s,[
  ([("Both errors have the same shape: ",True,NAVY),("a claim about a summary or a component, "
    "reported as a claim about the behaviour. Neither table was ever wrong.",False,DARK)],0),
  ([("The fix is cheap and now in the analysis code: ",True,NAVY),("per-session counts are tested "
    "as their own family beside the per-turn rates, and the aggregate is plotted on the same axes "
    "as the channel, so a channel result cannot be shown without its total.",False,DARK)],0),
  ([("The methodological point is the paper's contribution, not an aside. ",True,GREEN),
   ("Any intervention will look successful when it is scored on the channel that motivated it. "
    "The aggregate has to be nominated BEFORE the intervention.",False,DARK)],0),
], top=5.0, size=13.5, gap=12)

# =====================================================================
# 15 · TWO DRAFTS
# =====================================================================
s = slide(); title_bar(s,"Two drafts, deliberately disjoint","WRITE-UP")

cards(s,[
  ("Affirmation Without Inquiry  —  K = 0 only",
   "What the reward hack IS. PTO vs GRPO at matched look-ahead; the loop works on its own terms, "
   "what it actually taught (turns 2.3–3.4× longer, questions collapse), and how much survives the "
   "held-out judge. 8-page body, fits, 0 overfull boxes.",
   PTO,"complete draft"),
  ("The Hack Moves  —  PTO only, K ∈ {0,5}",
   "Whether a principled fix removes it. The channel closes for the whole run; the failure "
   "substitutes rather than falls; the mechanism is measured at the reward. 8 pages, 0 overfull "
   "boxes — related work is still a scaffold of citation slots.",
   GREEN,"drafted this week"),
],1.45,left=0.6,width=12.1,height=1.5,gap=0.18)

bullets(s,[
  ([("They share no claims. ",True,NAVY),("The split is the K axis: the first holds the optimiser "
    "as the variable and K fixed at 0; the second holds the optimiser fixed and varies K. Each "
    "carries a claim→artifact ledger, and the over-praise numbers appear in both at different "
    "iterations — the one place they could cross-contaminate.",False,DARK)],0),
  ([("Neither is a K × method result. ",True,AMBER),("GRPO at K=5 has one iteration. Both drafts "
    "say so in Limitations rather than implying otherwise.",False,DARK)],0),
  ([("Every number in both traces to a generated artifact",True,NAVY),(", not to a summary "
    "document — which is what caught the two errors on the previous slide.",False,DARK)],0),
], top=4.85, size=13.5, gap=13)

# =====================================================================
# 16 · DECISION 1 — scope
# =====================================================================
s = slide(); title_bar(s,"Decision 1 — scope and order","DECISION")
caption(s,"Two drafts exist and the thesis needs a spine. These are not mutually exclusive; the "
          "question is what gets finished first and what the chapter is built around.",
        0.6,1.3,12.1,size=13,align=PP_ALIGN.LEFT)
cards(s,[
  ("Submit the K=0 paper first, keep look-ahead for the thesis",
   "The reward-hacking result is the more complete and the more self-contained. Look-ahead becomes "
   "a thesis chapter where the GRPO gap is a stated limitation rather than a reviewer's first "
   "question.",
   PTO,"lowest risk"),
  ("Submit the look-ahead paper first",
   "It is the more novel claim and the one with a methodological contribution beyond MI. But it "
   "rests on one optimiser and one run per arm, and a reviewer will ask for GRPO at K=5.",
   AMBER,"higher ceiling, higher risk"),
  ("Merge into one longer paper",
   "The hack, then the attempted fix, then the substitution — one arc. Does not fit 8 pages, so it "
   "means a different venue and a later date.",
   GREY,None),
],1.95,left=0.6,width=12.1,height=1.42,gap=0.16)
callout(s,0.6,6.5,12.1,0.8,
  [("My read: ",True,NAVY),("the first option. The K=0 draft is closer to done, and the look-ahead "
    "result gets materially stronger for ~one more spend decision — which is Decision 2.",
    False,DARK)],size=13)

# =====================================================================
# 17 · DECISION 2 — budget
# =====================================================================
s = slide(); title_bar(s,"Decision 2 — what, if anything, is worth buying","DECISION · BUDGET")
caption(s,"Total API spend to date is roughly $312 and is the binding constraint on everything "
          "below. Two candidates are left; only one of them changes a claim.",
        0.6,1.3,12.1,size=13,align=PP_ALIGN.LEFT)
cards(s,[
  ("GRPO at K = 5, five or six iterations",
   "The only purchase that turns the look-ahead result from a within-PTO finding into a K × method "
   "result. It is also the one a reviewer will ask for. Requires training, not just scoring, so it "
   "is by far the most expensive thing on this slide.",
   GREEN,"changes a claim"),
  ("Nothing further",
   "Both drafts are defensible as they stand, with the gaps named in Limitations. Spend the "
   "remaining budget on nothing, and write.",
   GREY,None),
  ("PTO K = 5 out to iterations 9 and 10",
   "Was the third candidate on this slide. It has been run: scoring both graders cost ≈ $2 and "
   "closed the prevention-vs-delay question — K=5's over-praise stays low through iteration 10 "
   "under both. It did not change the substitution result, which now stands over eleven matched "
   "iterations.",
   PTO,"done — off the list"),
],1.95,left=0.6,width=12.1,height=1.42,gap=0.16)
callout(s,0.6,6.5,12.1,0.8,
  [("Note on cost asymmetry: ",True,AMBER),("K=5 spent MORE than K=0 over the full ten iterations "
    "(7,548 preference groups built vs 6,240, each paying five extra simulated turns). Matching "
    "iterations is the right control, but it is generous to K=0 on budget.",False,DARK)],
  accent=AMBER,size=12.5)

# =====================================================================
# 18 · ASKS
# =====================================================================
s = slide(); title_bar(s,"What I'd like to leave with","ASKS")
bullets(s,[
  ([("A decision on order. ",True,NAVY),("Which draft is finished first, and is the other a thesis "
    "chapter or a second submission?",False,DARK)],0),
  ([("A go / no-go on GRPO at K = 5. ",True,NAVY),("It is the only remaining purchase that changes "
    "what either paper can claim.",False,DARK)],0),
  ([("A view on human MI coders. ",True,NAVY),("Every behaviour count in both drafts is produced "
    "by a language model applying a coding scheme. Two graders from different families agree, but "
    "agreement is not validity, and no MI-trained human has read a single transcript. Even 40–50 "
    "coded sessions would change what we can claim — do we have access to a coder?",False,DARK)],0),
  ([("Co-author list and order",True,NAVY),(" for both drafts, plus venue and target date.",
    False,DARK)],0),
  ([("A sanity check on the clinical reading. ",True,NAVY),("Over-praise and unsolicited advice "
    "are both MI-inconsistent, but they fail differently, and our two graders disagree about which "
    "mix is more severe. I have declined to rank them — is that the right call?",False,DARK)],0),
], top=1.6, size=15, gap=17)

callout(s,0.6,6.35,12.1,0.85,
  [("Everything in this deck regenerates from one command. ",True,NAVY),
   ("Figures, tables and summaries rebuild under either grader with seeded intervals; the analysis "
    "package carries a 20-check self-test that runs after every change.",False,DARK)],size=13)

# =====================================================================
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
