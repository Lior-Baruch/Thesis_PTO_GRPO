#!/usr/bin/env python
"""
build_rq_deck_2026-08-21.py — the 2026-08-21 supervisor deck.

A NEW deck, not a fork of any previous builder. Every earlier deck in ``meetings/build/`` uses the
same project-status spine (status -> setup -> ICLR background -> arms -> results -> drafts ->
decisions). This one is organised by the **three research questions the thesis actually asks**, and
every answer slide carries the same three bands:

    VERDICT   — the one-line answer, stated as a claim
    EVIDENCE  — the table or figure that carries it, with the owning artifact path on the slide
    NOT THIS  — the nearest stronger claim the evidence does NOT support

The "NOT THIS" band is the point of the deck. It exists so nobody leaves the room repeating a
stronger version of a result than the tables defend — the failure mode CLAUDE.md's "Epistemic
status" section was written to prevent.

NUMBERS ARE HARD-CODED, like every builder here: a deck is a snapshot of what was presented on its
date, not a live view. Every figure below was verified against the rendered tables on 2026-08-21,
after the full re-render that day (6 units / 21 notebooks, no failures). Owning artifacts, all
under ``Exp3_PTO_GRPO/eda/results/``:

    arms/outcomes/tables/<judge>/leaderboard_scorecard.md   endpoint levels, both graders
    lookahead/reward/tables/k_table1.md                     K contrast, persona-paired
    lookahead/reward/tables/k_did.md                        K x method difference-in-differences
    lookahead/transfer/tables/k_retention_summary.md        gain retention by K
    method/contrast/tables/method_paired_by_K.md            PTO - GRPO at matched K + iteration
    method/contrast/tables/method_paired_best.md            PTO - GRPO best-vs-best
    compute/cost/tables/compute_by_{arm,iteration}.md       GPU-hours
    compute/cost/tables/budget_sweep_GRPO_K_<judge>.md      K contrast at matched budget
    measurement/validity/tables/multijudge_*.md             sign preservation, dependability

Build:
    & ..\\..\\.venv\\Scripts\\python.exe build_rq_deck_2026-08-21.py
    .\\export_pdf.ps1 ..\\2026-08-21\\rq_2026-08-21.pptx
"""

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
REPO = os.path.dirname(os.path.dirname(HERE))
ROOT = os.path.join(REPO, "Exp3_PTO_GRPO")
RES = os.path.join(ROOT, "eda", "results")
OUTDIR = os.path.join(REPO, "meetings", "2026-08-21")
OUT = os.path.join(OUTDIR, "rq_2026-08-21.pptx")

PRIMARY = "gpt-4o-mini"
HELDOUT = "claude-haiku-4-5"

# ── palette ───────────────────────────────────────────────────────────────────
INK = RGBColor(0x12, 0x20, 0x3A)   # deep navy — headings, dividers
BODY = RGBColor(0x2B, 0x33, 0x45)  # body text
MUTED = RGBColor(0x69, 0x72, 0x82)  # captions, provenance
RULE = RGBColor(0xD8, 0xDE, 0xE7)  # hairlines
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF4, 0xF6, 0xF9)  # band fill
VERDICT = RGBColor(0x1E, 0x5B, 0x45)  # green — the answer band
VWASH = RGBColor(0xE8, 0xF1, 0xED)
CAVEAT = RGBColor(0x8A, 0x5A, 0x12)  # amber — the "not this" band
CWASH = RGBColor(0xFB, 0xF3, 0xE3)
PTO_C = RGBColor(0x2F, 0x6F, 0x8F)   # teal — PTO
GRPO_C = RGBColor(0xC2, 0x70, 0x3D)  # burnt orange — GRPO

FONT = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
BOTTOM = Inches(0.66)   # clearance kept under a bottom-anchored band for the source line
ML, MR = Inches(0.85), Inches(0.85)
CW = W - ML - MR


# ── primitives ────────────────────────────────────────────────────────────────
def _txbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _run(p, text, *, size, bold=False, color=BODY, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _para(tf, *, space_before=0, space_after=0, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.alignment = align
    return p


def _rect(slide, x, y, w, h, fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _pic(slide, path, x, y, max_w, max_h):
    """Insert a PNG scaled to fit (max_w, max_h), centred in that box.

    Fails loudly if the artifact moved — a deck is never written half-updated.
    """
    if not os.path.exists(path):
        raise SystemExit(
            "MISSING ARTIFACT: %s\n"
            "The EDA tree moved or the family was not rendered. Re-run\n"
            "  python tools/render_results.py\n"
            "from Exp3_PTO_GRPO/eda/ and rebuild." % path
        )
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(path, int(x + (max_w - w) / 2), int(y + (max_h - h) / 2),
                                    width=w, height=h)


def _rel(path):
    """results-relative path, for the provenance line printed on the slide."""
    return os.path.relpath(path, RES).replace(os.sep, "/")


# ── slide furniture ───────────────────────────────────────────────────────────
def _running_head(slide, kicker):
    tf = _txbox(slide, ML, Inches(0.34), CW, Inches(0.26))
    p = _para(tf, first=True)
    _run(p, kicker.upper(), size=10.5, bold=True, color=MUTED)
    _rect(slide, ML, Inches(0.66), CW, Emu(9525), RULE)


def _heading(slide, text, y=Inches(0.86), size=27):
    tf = _txbox(slide, ML, y, CW, Inches(0.62))
    p = _para(tf, first=True)
    _run(p, text, size=size, bold=True, color=INK)
    return y + Inches(0.66)


def _band_h(text, size):
    """Height a band needs for `text` at `size`, from an empirically calibrated wrap width.

    Measured off the rendered deck: a 97-char run at 15.5 pt sits on one line inside CW, a
    152-char run wraps to two. 1600/size reproduces both, and the 12.5 pt caveat bands.
    """
    cpl = max(40, int(1600 / size))
    lines = max(1, -(-len(text) // cpl))
    return max(Inches(0.80), Inches(0.32) + Inches(0.26) * lines)


def _band(slide, y, label, text, *, fill, edge, label_color, height=None,
          size=15.5, bold_text=True):
    """A labelled band — the deck's repeating device (VERDICT / NOT THIS).

    `height=None` sizes the band to its text, so a long caveat can never spill over the
    provenance line (it did, on the first build).
    """
    height = _band_h(text, size) if height is None else height
    _rect(slide, ML, y, CW, height, fill)
    _rect(slide, ML, y, Inches(0.055), height, edge)
    tf = _txbox(slide, ML + Inches(0.28), y + Inches(0.11), CW - Inches(0.5),
                height - Inches(0.22))
    p = _para(tf, first=True)
    _run(p, label, size=9.5, bold=True, color=label_color)
    p2 = _para(tf, space_before=3)
    _run(p2, text, size=size, bold=bold_text, color=INK)
    return y + height + Inches(0.2)


def _bandbot(slide, label, text, **kw):
    """Bottom-anchored band: its BOTTOM edge sits just above the source line, whatever its height."""
    h = _band_h(text, kw.get("size", 15.5))
    return _band(slide, H - BOTTOM - h, label, text, height=h, **kw)


def _figband(slide, y, png, label, text, **kw):
    """Figure filling the space above a bottom-anchored caption band (band sized first)."""
    h = _band_h(text, kw.get("size", 12.5))
    band_y = H - BOTTOM - h
    _pic(slide, png, ML, y + Inches(0.02), CW, band_y - y - Inches(0.14))
    _band(slide, band_y, label, text, height=h, **kw)


def _provenance(slide, paths, y=None):
    y = y if y is not None else H - Inches(0.52)
    tf = _txbox(slide, ML, y, CW, Inches(0.3))
    p = _para(tf, first=True)
    _run(p, "source:  ", size=8.5, color=MUTED)
    _run(p, "   ·   ".join(paths), size=8.5, color=MUTED, font=MONO)


def _footer(slide, n):
    tf = _txbox(slide, W - MR - Inches(1.2), H - Inches(0.5), Inches(1.2), Inches(0.3))
    p = _para(tf, first=True, align=PP_ALIGN.RIGHT)
    _run(p, str(n), size=9.5, color=MUTED)


def _bullets(slide, y, items, *, size=13.5, gap=9, width=None, x=None):
    width = width or CW
    x = ML if x is None else x
    tf = _txbox(slide, x, y, width, H - y - Inches(0.7))
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
        else:
            lead, rest = None, it
        p = _para(tf, first=(i == 0), space_after=gap)
        _run(p, "—   ", size=size, color=RULE)
        if lead:
            _run(p, lead, size=size, bold=True, color=INK)
            _run(p, "  ", size=size)
        _run(p, rest, size=size, color=BODY)
    return tf


def _table(slide, y, headers, rows, *, col_w, size=11.5, head_size=10, row_h=Inches(0.34),
           emphasis=None, prose_cols=()):
    """Native pptx table with the deck's flat styling (no banding, hairline rules)."""
    total = sum(col_w)
    left = ML + (CW - total) / 2
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), int(left), int(y),
                                   int(total), int(row_h * (len(rows) + 1)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = int(cw)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = INK
        c.margin_left = c.margin_right = Inches(0.09)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        _run(p, htxt, size=head_size, bold=True, color=PAPER)
    for i, row in enumerate(rows, start=1):
        tbl.rows[i].height = int(row_h)
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = PAPER if i % 2 else WASH
            c.margin_left = c.margin_right = Inches(0.09)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0 or j in prose_cols) else PP_ALIGN.CENTER
            colr = BODY
            bold = False
            if emphasis and emphasis(i - 1, j, val):
                colr, bold = INK, True
            _run(p, str(val), size=size, bold=bold, color=colr,
                 font=FONT if (j == 0 or j in prose_cols) else MONO)
    return shape


# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
N = 0


def newslide(kicker=None, heading=None):
    global N
    N += 1
    s = _blank(prs)
    y = Inches(0.86)
    if kicker:
        _running_head(s, kicker)
    if heading:
        y = _heading(s, heading)
    _footer(s, N)
    return s, y


# ── 1 · title ─────────────────────────────────────────────────────────────────
s = _blank(prs)
N += 1
_rect(s, 0, 0, W, H, INK)
_rect(s, ML, Inches(2.02), Inches(1.5), Inches(0.05), PTO_C)
tf = _txbox(s, ML, Inches(2.35), CW - Inches(1.4), Inches(2.4))
p = _para(tf, first=True)
_run(p, "Three questions,", size=44, bold=True, color=PAPER)
p = _para(tf, space_before=2)
_run(p, "and what the tables now answer", size=44, bold=True, color=RGBColor(0x8F, 0xB6, 0xCB))
p = _para(tf, space_before=20)
_run(p, "Look-ahead depth · optimizer family · the oracle instrument",
     size=16, color=RGBColor(0xB9, 0xC4, 0xD4))
tf = _txbox(s, ML, H - Inches(1.5), CW, Inches(0.9))
p = _para(tf, first=True)
_run(p, "Lior Baruch   ·   Reichman University   ·   21 August 2026",
     size=13, color=RGBColor(0x9F, 0xAD, 0xC0))
p = _para(tf, space_before=6)
_run(p, "Exp3 · Llama-3.2-1B therapist · gpt-4o-mini patient + oracle · "
        "Claude Haiku 4.5 held-out judge · 96 personas, persona-paired",
     size=10.5, color=RGBColor(0x7E, 0x8C, 0xA2))

# ── 2 · the three questions ───────────────────────────────────────────────────
s, y = newslide("Agenda", "The thesis asks exactly three questions")
y += Inches(0.06)
qs = [
    ("i", "Does looking ahead help?", "K = 0 vs K = 5, within each optimizer",
     "Answered — depends on the optimizer", VERDICT),
    ("ii", "Which optimizer?", "PTO vs GRPO, at matched K and matched compute",
     "Answered — depends on K", VERDICT),
    ("iii", "Does the oracle questionnaire change the conclusion?",
     "Q1+Q2 vs WAI-SR vs CSQ-8 vs MI-SAT / MITI",
     "Held — not yet bought", CAVEAT),
]
for num, q, sub, state, col in qs:
    _rect(s, ML, y, CW, Inches(1.28), WASH)
    _rect(s, ML, y, Inches(0.055), Inches(1.28), col)
    tfn = _txbox(s, ML + Inches(0.34), y + Inches(0.2), Inches(0.9), Inches(0.9))
    pn = _para(tfn, first=True)
    _run(pn, "RQ-" + num, size=19, bold=True, color=col)
    tfq = _txbox(s, ML + Inches(1.35), y + Inches(0.16), CW - Inches(5.6), Inches(1.0))
    pq = _para(tfq, first=True)
    _run(pq, q, size=17.5, bold=True, color=INK)
    pq2 = _para(tfq, space_before=4)
    _run(pq2, sub, size=11.5, color=MUTED)
    tfs = _txbox(s, W - MR - Inches(4.05), y + Inches(0.42), Inches(3.9), Inches(0.5))
    ps = _para(tfs, first=True, align=PP_ALIGN.RIGHT)
    _run(ps, state, size=12, bold=True, color=col)
    y += Inches(1.46)
tf = _txbox(s, ML, y + Inches(0.04), CW, Inches(0.5))
p = _para(tf, first=True)
_run(p, "Everything after this slide is one of these three, plus what it cost and what to decide.",
     size=12, italic=True, color=MUTED)

# ── 3 · setup ─────────────────────────────────────────────────────────────────
s, y = newslide("Setup", "One slide of setup, then straight to the questions")
left_w = CW * 0.52
_bullets(s, y, [
    ("Therapist.", "Llama-3.2-1B, bf16, LoRA r=16. Trained by PTO (preference tree + DPO) "
                   "or GRPO (group-relative), K ∈ {0, 5} look-ahead turns."),
    ("Patient + training oracle.", "gpt-4o-mini. 96 fixed personas; every contrast below is "
                                   "paired on persona identity, n = 96."),
    ("Held-out judge.", "Claude Haiku 4.5 — different model family, never played the patient, "
                        "never was the training reward."),
    ("Training reward.", "Q1 + Q2 only (matching the ICLR paper). Evaluation reads all 8 instruments."),
], size=13, gap=11, width=left_w)

bx = ML + left_w + Inches(0.4)
bw = CW - left_w - Inches(0.4)
_rect(s, bx, y, bw, Inches(3.55), WASH)
tfb = _txbox(s, bx + Inches(0.28), y + Inches(0.24), bw - Inches(0.56), Inches(3.1))
pb = _para(tfb, first=True)
_run(pb, "THE GRID, AS OF TODAY", size=9.5, bold=True, color=MUTED)
for lead, rest in [
    ("4 arms", "PTO/GRPO × K ∈ {0, 5}"),
    ("40 model states", "11 + 11 + 11 + 7 (base + iterations)"),
    ("2 graders", "every state scored on both"),
    ("30,720 cells", "40 × 8 rubrics × 96 personas, per grader"),
]:
    pp = _para(tfb, space_before=13)
    _run(pp, lead, size=15, bold=True, color=INK)
    _run(pp, "   " + rest, size=11.5, color=BODY)
pn = _para(tfb, space_before=17)
_run(pn, "Complete except GRPO K=5, which stopped at iteration 6 (slide 16).",
     size=10.5, italic=True, color=CAVEAT)
_provenance(s, ["measurement/validity/tables/multijudge_coverage.md"])

# ── 4 · RQ-i divider ──────────────────────────────────────────────────────────
s = _blank(prs)
N += 1
_rect(s, 0, 0, W, H, INK)
tf = _txbox(s, ML, Inches(2.5), CW, Inches(2.6))
p = _para(tf, first=True)
_run(p, "RQ-i", size=15, bold=True, color=PTO_C)
p = _para(tf, space_before=12)
_run(p, "Does looking ahead help?", size=40, bold=True, color=PAPER)
p = _para(tf, space_before=16)
_run(p, "Scoring a reply on its own rewards openings that look good in isolation. Scoring it "
        "after K simulated turns rewards openings that lead somewhere.",
     size=15, color=RGBColor(0x9F, 0xAD, 0xC0))
_footer(s, N)

# ── 5 · RQ-i answer ───────────────────────────────────────────────────────────
s, y = newslide("RQ-i · does looking ahead help?", "The lever works — but it points in "
                                                   "opposite directions")
y = _band(s, y, "VERDICT",
          "Look-ahead HELPS GRPO and HURTS PTO. At iteration 6 both effects are significant, "
          "on both graders.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
rows = [
    ["4", "+0.120  (0.20)", "+0.123  (0.21)", "−0.115  (−0.25) *", "−0.233  (−0.37) **"],
    ["5", "−0.002  (−0.00)", "+0.173  (0.33) *", "−0.070  (−0.13)", "−0.311  (−0.43) **"],
    ["6", "+0.257  (0.42) ***", "+0.343  (0.51) ***", "−0.263  (−0.42) ***", "−0.533  (−0.55) ***"],
]
_table(s, y + Inches(0.05),
       ["iteration", "PTO · primary", "PTO · held-out", "GRPO · primary", "GRPO · held-out"],
       rows, col_w=[Inches(1.5), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)],
       emphasis=lambda i, j, v: i == 2)
ty = y + Inches(0.05) + Inches(0.34) * 4 + Inches(0.18)
tf = _txbox(s, ML, ty, CW, Inches(0.4))
p = _para(tf, first=True)
_run(p, "Q1+Q2, persona-paired mean difference (Cohen's dz).  Sign: + = K=0 higher.  "
        "* p_holm < .05   ** < .01   *** < .001", size=10.5, color=MUTED)
_band(s, ty + Inches(0.42), "NOT THIS",
      "Not that look-ahead is worthless. On PTO it never significantly WINS on either grader "
      "— that is a null, not a harm, everywhere except iteration 6.",
      fill=CWASH, edge=CAVEAT, label_color=CAVEAT, size=13.5)
_provenance(s, ["lookahead/reward/tables/k_table1.md"])

# ── 6 · RQ-i figure ───────────────────────────────────────────────────────────
s, y = newslide("RQ-i · does looking ahead help?", "The same lever, both optimizers, both graders")
fig = os.path.join(RES, "lookahead", "reward", "figures", "k_contrast_both_judges.png")
_figband(s, y, fig, "READ IT AS",
         "Two panels that mirror each other. Whatever look-ahead is doing, it is not a property of "
         "the lever alone — it is an interaction with the optimizer.",
         fill=WASH, edge=PTO_C, label_color=MUTED, size=12.5, bold_text=False)
_provenance(s, [_rel(fig)])

# ── 7 · RQ-i · the DiD, and what changed this week ────────────────────────────
s, y = newslide("RQ-i · does looking ahead help?",
                "New this week: the training grader can now see it too")
y = _band(s, y, "VERDICT",
          "The K × optimizer interaction is now significant on the PRIMARY oracle, not just "
          "the held-out judge.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
rows = [
    ["4", "+0.235", "0.286", ".104   (ns)", "+0.356", "0.401", ".008"],
    ["5", "+0.068", "0.095", "1.000  (ns)", "+0.484", "0.525", "<.001"],
    ["6", "+0.520", "0.605", "<.001", "+0.876", "0.754", "<.001"],
]
_table(s, y + Inches(0.05),
       ["iter", "DiD · primary", "dz", "p_holm", "DiD · held-out", "dz", "p_holm"],
       rows, col_w=[Inches(1.0), Inches(1.95), Inches(1.0), Inches(1.35),
                    Inches(1.95), Inches(1.0), Inches(1.35)],
       emphasis=lambda i, j, v: i == 2)
ty = y + Inches(0.05) + Inches(0.34) * 4 + Inches(0.16)
_bullets(s, ty, [
    ("Until iteration 5,", "the interaction was visible only to the held-out judge — which was "
                           "the single sharpest argument in the thesis for having a second grader."),
    ("At iteration 6,", "the primary sees it as well (dz 0.605). The held-out judge still sees it "
                        "1.68× more strongly (0.876 / 0.520)."),
], size=12.5, gap=7)
_bandbot(s, "RETRACTED",
      "“The grader that WAS the training reward cannot see this effect.”  True through "
      "iteration 5; false at 6. The defensible claim is now “sees it less sharply”.",
      fill=CWASH, edge=CAVEAT, label_color=CAVEAT, size=12.5)
_provenance(s, ["lookahead/reward/tables/k_did.md"])

# ── 8 · RQ-i · does the gain survive a grader who never trained on it ─────────
s, y = newslide("RQ-i · does looking ahead help?",
                "Does the gain survive a grader that was never the reward?")
y = _band(s, y, "VERDICT",
          "On GRPO, K=5 keeps 84% of its gain under the held-out judge; K=0 keeps 57%. "
          "Confidence intervals disjoint.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
rows = [
    ["GRPO", "6", "0.567", "[0.379, 0.728]", "0.844", "[0.745, 0.959]", "yes"],
    ["PTO", "10", "0.823", "[0.720, 0.947]", "0.639", "[0.551, 0.738]", "no"],
]
_table(s, y + Inches(0.1),
       ["method", "iter", "K=0 retention", "95% CI", "K=5 retention", "95% CI", "disjoint?"],
       rows, col_w=[Inches(1.5), Inches(0.85), Inches(1.9), Inches(2.1),
                    Inches(1.9), Inches(2.1), Inches(1.25)],
       emphasis=lambda i, j, v: i == 0 and j in (4, 6))
ty = y + Inches(0.1) + Inches(0.34) * 3 + Inches(0.2)
_bullets(s, ty, [
    ("Retention", "= (gain measured by the held-out judge) ÷ (gain measured by the training "
                  "oracle). 1.0 means the improvement is fully real to a grader with no stake in it."),
    ("Reward hacking, quantified.", "A low retention number is a gain the training oracle believes "
                                    "and an independent grader does not."),
], size=12.5, gap=8)
_bandbot(s, "NOT THIS",
      "Not that PTO reward-hacks and GRPO does not. PTO's K=0 arm retains 0.823 — higher than "
      "either GRPO arm. The contrast is WITHIN method, across K.",
      fill=CWASH, edge=CAVEAT, label_color=CAVEAT, size=12.5)
_provenance(s, ["lookahead/transfer/tables/k_retention_summary.md"])

# ── 9 · RQ-ii divider ─────────────────────────────────────────────────────────
s = _blank(prs)
N += 1
_rect(s, 0, 0, W, H, INK)
tf = _txbox(s, ML, Inches(2.5), CW, Inches(2.6))
p = _para(tf, first=True)
_run(p, "RQ-ii", size=15, bold=True, color=GRPO_C)
p = _para(tf, space_before=12)
_run(p, "Which optimizer?", size=40, bold=True, color=PAPER)
p = _para(tf, space_before=16)
_run(p, "PTO grows a best-of-M reranked trunk; GRPO slices an on-policy rollout. Same candidate "
        "budget (M = G = 8), same temperature, same oracle.",
     size=15, color=RGBColor(0x9F, 0xAD, 0xC0))
_footer(s, N)

# ── 10 · RQ-ii at K=0 ─────────────────────────────────────────────────────────
s, y = newslide("RQ-ii · which optimizer?", "At K = 0, PTO wins clearly")
y = _band(s, y, "VERDICT",
          "At the matched 10-iteration endpoint with no look-ahead, PTO beats GRPO on both "
          "graders, with a large effect on the held-out judge.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
rows = [
    ["Q1+Q2", "+0.507", "0.729", "<.001", "+0.609", "1.265", "<.001"],
    ["MITI (MI integrity)", "+0.352", "0.459", "<.001", "+0.253", "0.648", "<.001"],
    ["MICI (inconsistency) ↓", "−0.346", "−0.989", "<.001", "−0.225", "−0.667", "<.001"],
]
_table(s, y + Inches(0.1),
       ["metric", "Δ primary", "dz", "p_holm", "Δ held-out", "dz", "p_holm"],
       rows, col_w=[Inches(2.9), Inches(1.55), Inches(1.0), Inches(1.25),
                    Inches(1.55), Inches(1.0), Inches(1.25)],
       emphasis=lambda i, j, v: i == 0)
ty = y + Inches(0.1) + Inches(0.34) * 4 + Inches(0.2)
_bullets(s, ty, [
    ("Sign.", "Δ = PTO − GRPO. Positive favours PTO; MICI is lower-is-better, so a "
              "negative Δ also favours PTO."),
    ("Mechanism.", "GRPO K=0 peaks at iteration 8 (4.082) then regresses to 3.753 while its "
                   "MI-inconsistency climbs. PTO climbs steadily to 4.260."),
], size=12.5, gap=8)
_provenance(s, ["method/contrast/tables/method_paired_by_K.md",
                "arms/outcomes/tables/<judge>/leaderboard_scorecard.md"])

# ── 11 · RQ-ii at K=5 — the flip ──────────────────────────────────────────────
s, y = newslide("RQ-ii · which optimizer?", "At K = 5, the answer flips")
y = _band(s, y, "VERDICT",
          "Turn look-ahead on and GRPO wins. At matched iteration 6 it beats PTO on both graders; "
          "best-vs-best it ties on the primary and WINS on the held-out judge.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
rows = [
    ["matched iteration 6", "−0.332", "−0.437", "<.001", "−0.397", "−0.599", "<.001"],
    ["best vs best", "+0.078", "0.133", "1.000 (ns)", "−0.168", "−0.352", ".009"],
]
_table(s, y + Inches(0.08),
       ["comparison", "Δ primary", "dz", "p_holm", "Δ held-out", "dz", "p_holm"],
       rows, col_w=[Inches(2.9), Inches(1.55), Inches(1.0), Inches(1.3),
                    Inches(1.55), Inches(1.0), Inches(1.3)],
       emphasis=lambda i, j, v: True)
ty = y + Inches(0.08) + Inches(0.34) * 3 + Inches(0.2)
_bullets(s, ty, [
    ("Δ = PTO − GRPO,", "so every negative number on this slide favours GRPO."),
    ("Best-vs-best pairs", "PTO@10 against GRPO@6 (primary) and PTO@7 against GRPO@6 (held-out) "
                           "— each arm at its own peak."),
    ("So RQ-i and RQ-ii are not separable.", "“Which optimizer is better” has no answer "
                                             "that is not conditional on K."),
], size=12.5, gap=8)
_bandbot(s, "NOT THIS",
      "Not that GRPO overtakes PTO given enough iterations. GRPO K=5 has SIX; PTO has ten. "
      "The flip is at matched iteration, and the arm is censored.",
      fill=CWASH, edge=CAVEAT, label_color=CAVEAT, size=12.5)
_provenance(s, ["method/contrast/tables/method_paired_{by_K,best}.md"])

# ── 12 · RQ-ii trajectories ───────────────────────────────────────────────────
s, y = newslide("RQ-ii · which optimizer?", "The whole picture, held-out judge")
fig = os.path.join(RES, "arms", "outcomes", "figures", HELDOUT, "trajectories",
                   "trajectory_Q1Q2.png")
_figband(s, y, fig, "READ IT AS",
         "GRPO K=0 (orange) peaks early then collapses to 2.257. GRPO K=5 climbs monotonically to "
         "2.903 in six iterations — the highest endpoint of any arm on this grader.",
         fill=WASH, edge=GRPO_C, label_color=MUTED, size=12.5, bold_text=False)
_provenance(s, [_rel(fig)])

# ── 13 · RQ-ii on the compute axis ────────────────────────────────────────────
s, y = newslide("RQ-ii · which optimizer?", "An iteration is not a unit of spend")
y = _band(s, y, "VERDICT",
          "On the compute axis PTO dominates outright: it reaches iteration 10 for 8.1 GPU-h "
          "against GRPO's 27.9 — 3.4× cheaper, and it scores higher.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
rows = [
    ["PTO  K=0", "10", "8.12", "0.81", "4.260", "2.866"],
    ["PTO  K=5", "10", "19.68", "1.97", "4.307", "2.667"],
    ["GRPO K=0", "10", "27.91", "2.79", "3.753", "2.257"],
    ["GRPO K=5", "6 (+ partial 7)", "30.53", "5.09", "4.229", "2.903"],
]
_table(s, y + Inches(0.1),
       ["arm", "iterations", "GPU-h", "GPU-h / iter", "Q1+Q2 primary", "Q1+Q2 held-out"],
       rows, col_w=[Inches(1.9), Inches(2.0), Inches(1.4), Inches(1.8),
                    Inches(2.2), Inches(2.2)],
       emphasis=lambda i, j, v: (i == 0 and j in (2, 3)) or (i == 3 and j == 5))
ty = y + Inches(0.1) + Inches(0.34) * 5 + Inches(0.2)
_bullets(s, ty, [
    ("Why PTO is cheap.", "Its dominant phase is building preference pairs (5.7 of 8.1 h), which "
                          "does not scale with the optimizer step count. GRPO computes its reward "
                          "inside the training loop."),
    ("The pair that is no longer matched.", "The two GRPO arms were “budget-matched within 3%” "
                                            "at iteration 5. At iteration 6 the K=5 arm is 9.4% MORE "
                                            "expensive; that framing is retired."),
], size=12.5, gap=8)
_provenance(s, ["compute/cost/tables/compute_by_arm.md"])

# ── 14 · RQ-iii ───────────────────────────────────────────────────────────────
s = _blank(prs)
N += 1
_rect(s, 0, 0, W, H, INK)
tf = _txbox(s, ML, Inches(1.5), CW, Inches(1.4))
p = _para(tf, first=True)
_run(p, "RQ-iii", size=15, bold=True, color=RGBColor(0xD8, 0xA6, 0x5C))
p = _para(tf, space_before=12)
_run(p, "Does the oracle questionnaire change the conclusion?", size=34, bold=True, color=PAPER)
bx, by, bw = ML, Inches(3.15), CW
_rect(s, bx, by, bw, Inches(2.45), RGBColor(0x1B, 0x2C, 0x4B))
tfb = _txbox(s, bx + Inches(0.42), by + Inches(0.32), bw - Inches(0.84), Inches(1.85))
pb = _para(tfb, first=True)
_run(pb, "STATUS — HELD, DELIBERATELY", size=9.5, bold=True, color=RGBColor(0xD8, 0xA6, 0x5C))
pb = _para(tfb, space_before=10)
_run(pb, "All four arms train on Q1+Q2 only. Every state is EVALUATED on all 8 instruments, so we "
         "can already say the conclusions do not depend on which rubric we READ.",
     size=15, color=PAPER)
pb = _para(tfb, space_before=12)
_run(pb, "What is untested is whether they depend on which rubric we TRAIN ON. That needs new "
         "training runs, not new scoring — four more arms at roughly 8–28 GPU-h each.",
     size=15, color=RGBColor(0xB9, 0xC4, 0xD4))
pb = _para(tfb, space_before=14)
_run(pb, "Recommendation: keep it held. RQ-i and RQ-ii are the thesis; this is a fourth chapter "
         "we do not need.", size=14, bold=True, color=RGBColor(0x8F, 0xC4, 0xAA))
_footer(s, N)

# ── 15 · measurement validity ─────────────────────────────────────────────────
s, y = newslide("Can we trust any of this?", "The measurement thread, in one slide")
y = _band(s, y, "VERDICT",
          "An independent grader, in a different model family, reproduces 88.5% of all 6,240 "
          "arm × metric contrasts — and 99.5% of the large ones.",
          fill=VWASH, edge=VERDICT, label_color=VERDICT)
left_w = CW * 0.47
rows = [
    ["all contrasts", "6,240", "88.5%"],
    ["|Δ| ≥ 0.10", "4,035", "94.5%"],
    ["|Δ| ≥ 0.25", "2,539", "97.4%"],
    ["|Δ| ≥ 0.50", "1,309", "99.5%"],
]
sh = _table(s, y + Inches(0.1), ["subset", "n", "same sign"], rows,
            col_w=[Inches(2.2), Inches(1.4), Inches(1.7)])
sh.left = int(ML)
_bullets(s, y + Inches(0.1) + Inches(0.34) * 5 + Inches(0.14), [
    ("6,240", "= 8 metrics × C(40, 2) = 8 × 780."),
], size=11, gap=4, width=left_w, x=ML)

bx = ML + left_w + Inches(0.35)
bw = CW - left_w - Inches(0.35)
_bullets(s, y + Inches(0.02), [
    ("Oracle repeatability.", "ICC(2,1) 0.86–0.99 across four re-draws of the same conversations."),
    ("Dependability (1 judge).", "Strong for Q1/Q2/WAI-SR/CSQ-8/PCT (0.91–0.96) and MICI (0.85). "
                                 "MITI is the weak instrument at 0.62 — name it whenever a "
                                 "channel result rests on MITI."),
    ("Corrected today.", "This deck's predecessors quoted MITI 0.55 / MICI 0.63. Neither figure "
                         "appears in any table; the rendered values are 0.622 and 0.845."),
], size=12.5, gap=10, width=bw, x=bx)
_provenance(s, ["measurement/validity/tables/multijudge_{sign_preservation,variance_components}.md"])

# ── 16 · the infrastructure slide ─────────────────────────────────────────────
s, y = newslide("Why GRPO K=5 stops at 6", "One honest slide about the run that stopped")
y = _band(s, y, "WHAT HAPPENED",
          "Iteration 7 needs 106 optimizer steps. 40 are on disk. Across four Colab sessions, "
          "132 further steps were computed and then thrown away.",
          fill=CWASH, edge=CAVEAT, label_color=CAVEAT)
rows = [
    ["1", "19 Aug", "1 → 103", "1 → 30", "writes to Drive stopped mid-save"],
    ["2", "20 Aug", "—", "—", "OpenAI org spend limit (384 of 395 log lines)"],
    ["3", "20 Aug", "—", "—", "OpenAI org spend limit"],
    ["4", "20 Aug", "31 → 99", "31 → 40", "writes to Drive stopped mid-save"],
]
_table(s, y + Inches(0.05), ["session", "date", "steps trained", "steps saved", "outcome"],
       rows, col_w=[Inches(1.1), Inches(1.1), Inches(1.7), Inches(1.6), Inches(5.4)],
       emphasis=lambda i, j, v: j == 2, prose_cols=(4,))
ty = y + Inches(0.05) + Inches(0.34) * 5 + Inches(0.18)
_bullets(s, ty, [
    ("Training was never the problem.", "Steps after the stall ran at 162.1 s against 161.1 s "
                                        "before it. The GPU and the oracle were both healthy the "
                                        "whole time — only new files stopped reaching Drive."),
    ("It is not chronic.", "All 16 previously completed iterations have exactly as many saved "
                           "artifacts as optimizer steps. Iteration 7 is the only anomaly."),
    ("The fix is operational.", "Write checkpoints to local Colab disk and copy to Drive once per "
                                "iteration. Nothing in the experiment name, the reward or the data "
                                "changes. Credits are already topped up."),
], size=12, gap=7)

# ── 17 · where the evidence is thin ───────────────────────────────────────────
s, y = newslide("Threats", "Where the evidence is thin — said out loud")
_bullets(s, y, [
    ("Every contested endpoint is a single 96-conversation draw.", "The only noise floor we have "
     "is at the base: 4 independent draws of the identical base policy give 54 same-policy "
     "contrasts, 0 of which reach even uncorrected p < .05 (max |dz| 0.128). That is reassuring "
     "for the base and says nothing about a trained checkpoint."),
    ("GRPO K=5 is censored at 6.", "Every statement about it is “within six iterations”. "
     "Whether its lead persists, or whether it regresses the way GRPO K=0 did after iteration 8, "
     "is simply unobserved."),
    ("MITI dependability is 0.62 off one judge,", "and MITI carries the MI-integrity channel "
     "results. There is no channel-level ICC at all."),
    ("All 96 personas are used for both training and evaluation", "at every iteration, so every "
     "number is in-sample with respect to the patient distribution."),
    ("Two statistics quoted in earlier decks have no owning table", "and are withdrawn until they "
     "are rendered: MI acts per 1,000 therapist characters, and the per-conversation over-praise "
     "share."),
], size=12.5, gap=11)
_provenance(s, ["results/LIMITATIONS.md"])

# ── 18 · decisions ────────────────────────────────────────────────────────────
s, y = newslide("Decisions", "Three things I would like decided today")
y += Inches(0.02)
decisions = [
    ("1", "Finish GRPO K=5?", GRPO_C,
     "Resume iteration 7 and run to 10: ~16–18 GPU-h, ~$85–130. The arm lands near "
     "48–50 GPU-h against its K=0 sibling's 28.",
     "My recommendation: yes — it is the only arm without an endpoint, and it is currently "
     "the best-scoring arm on the held-out judge."),
    ("2", "Which axis does the thesis report?", PTO_C,
     "Matched iteration and matched compute answer different questions and disagree at K=5. "
     "PTO wins both at K=0.",
     "My recommendation: report both, and name the axis in every sentence that quotes a number."),
    ("3", "Buy the replicate draw?", VERDICT,
     "A second independent 96-conversation draw of ~5 contested endpoints: ~$10 and about one "
     "A100-hour, or four free hours locally. No code change needed.",
     "My recommendation: yes — it either retires the endpoint-fragility objection thesis-wide "
     "or tells us the headline is fragile."),
]
for num, q, col, body, rec in decisions:
    _rect(s, ML, y, CW, Inches(1.62), WASH)
    _rect(s, ML, y, Inches(0.055), Inches(1.62), col)
    tfn = _txbox(s, ML + Inches(0.32), y + Inches(0.16), Inches(0.6), Inches(0.6))
    pn = _para(tfn, first=True)
    _run(pn, num, size=22, bold=True, color=col)
    tfq = _txbox(s, ML + Inches(1.0), y + Inches(0.14), CW - Inches(1.5), Inches(1.4))
    pq = _para(tfq, first=True)
    _run(pq, q, size=15.5, bold=True, color=INK)
    pb = _para(tfq, space_before=5)
    _run(pb, body, size=11.5, color=BODY)
    pr = _para(tfq, space_before=5)
    _run(pr, rec, size=11.5, bold=True, color=col)
    y += Inches(1.78)

# ── 19 · in one line ──────────────────────────────────────────────────────────
s = _blank(prs)
N += 1
_rect(s, 0, 0, W, H, INK)
tf = _txbox(s, ML, Inches(1.9), CW, Inches(3.6))
p = _para(tf, first=True)
_run(p, "IN ONE LINE", size=11, bold=True, color=PTO_C)
p = _para(tf, space_before=18)
_run(p, "Look-ahead is not a knob with a sign.", size=32, bold=True, color=PAPER)
p = _para(tf, space_before=8)
_run(p, "It helps group-relative optimization and hurts preference-tree optimization,",
     size=24, color=RGBColor(0x8F, 0xB6, 0xCB))
p = _para(tf, space_before=4)
_run(p, "and which optimizer you should prefer depends on whether you turned it on.",
     size=24, color=RGBColor(0x8F, 0xB6, 0xCB))
p = _para(tf, space_before=26)
_run(p, "PTO remains 3.4× cheaper and wins outright at K = 0. GRPO with look-ahead is the "
        "best arm the held-out judge has seen — with six iterations to PTO's ten.",
     size=14, color=RGBColor(0x9F, 0xAD, 0xC0))
_footer(s, N)

# ── 20 · appendix: the method schematics ──────────────────────────────────────
# One per slide: the framework diagrams are 2.70:1, so side-by-side renders them unreadably small.
for name, png, col, note in [
    ("PTO — preference tree + DPO", "pto_framework.png", PTO_C,
     "Branch M candidates per therapist turn, look ahead K, score, keep (chosen, rejected) where "
     "the margin clears τ, then a DPO update."),
    ("GRPO — group-relative rollout", "grpo_framework.png", GRPO_C,
     "Slice the rollout after every patient turn, sample G completions per prompt, score all of "
     "them, and use the group-relative advantage. The oracle sits INSIDE the update."),
]:
    s, y = newslide("Appendix", name)
    _figband(s, y, os.path.join(RES, "schematics", png), "IN ONE LINE", note,
             fill=WASH, edge=col, label_color=MUTED, size=12.5, bold_text=False)
    _provenance(s, ["schematics/" + png])

# ── 21 · appendix: compute trajectory ─────────────────────────────────────────
s, y = newslide("Appendix", "GPU-hours per iteration, reconstructed from artifact mtimes")
fig = os.path.join(RES, "compute", "cost", "figures", "compute_trajectory.png")
_pic(s, fig, ML, y + Inches(0.02), CW, H - y - Inches(1.05))
_provenance(s, [_rel(fig)])

# ── write ─────────────────────────────────────────────────────────────────────
os.makedirs(OUTDIR, exist_ok=True)
prs.save(OUT)
print("wrote %s  (%d slides)" % (os.path.relpath(OUT, REPO), len(prs.slides._sldIdLst)))
