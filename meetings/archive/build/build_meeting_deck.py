"""Build the MEETING deck for the 2026-08-03 supervision meeting.

Written for a ROOM, not for one reader: several people, not all of whom have followed the thread,
so every slide has to stand on its own. No "new since the email", no "as discussed" — if something
is worth a slide it is stated in full, and if it is only interesting relative to a previous
conversation it does not belong here.

Same lean visual language as `build_results_snapshot.py`, but written to be *talked through*
rather than read cold:

- background + setup, enough that someone seeing the project today can follow the results;
- the results block, condensed;
- a measurement-validity block: the oracle's own repeatability, the second-judge sweep, sign
  preservation, gain retention, and the places where the two graders do NOT agree;
- the open look-ahead comparison, at its first matched-budget point;
- a mechanism block reading the training signal directly: the drift is a compounding loop rather
  than a hard pull, and the PTO-vs-GRPO gap survives swapping the weighting rule — i.e. it is
  about exploration, not the loss. That result is a direct input to Decision 1;
- and it closes on the open decisions: framing, budget, and the asks.

Numbers are owned by `eda/results/L0/SUMMARY.md` and `eda/docs/LIMITATIONS.md` — this script only
restates them. If a number here disagrees with those files, those files are right.

Exports .pptx, then convert to .pdf via export_pdf.ps1 (PowerPoint COM).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))       # meetings/build/
REPO = os.path.dirname(os.path.dirname(HERE))           # repo root (meetings/ lives here)
ROOT = os.path.join(REPO, "Exp3_PTO_GRPO")             # the experiment the artifacts come from
L0F  = os.path.join(ROOT, "eda", "results", "L0", "figures")
L0T  = os.path.join(ROOT, "eda", "results", "L0", "tables")
L5F  = os.path.join(ROOT, "eda", "results", "L5", "figures")
METH = os.path.join(ROOT, "figures")   # hand-authored method schematics (build_method_figures.py)
OUT  = os.path.join(REPO, "meetings", "2026-08-03", "meeting_2026-08-03.pptx")

# Since 2026-07-28 every grader has its own leaf: results/<view>/figures/<family>/<judge>/<name>.
# The decks show the PRIMARY grader's figures, so the judge segment is injected here rather than
# spelled out at ~40 call sites. Point JUDGE at another grader to rebuild the same deck off it.
JUDGE = "gpt-4o-mini"

# Families that are ABOUT the judges rather than produced BY one export with no <JUDGE> level
# (eda_analysis.exports.JUDGE_INVARIANT_GROUPS). Keep in sync, or a deck asks for a path that the
# EDA never writes -- which fails loudly at add_picture, but only when someone builds a deck.
JUDGE_INVARIANT_FAMILIES = {"8_measurement"}

def _jp(base, p):
    """``<base>/<family>[/<sub>]/<JUDGE>/<name>`` — judge goes ahead of the FILENAME, not the family.

    A judge-invariant family has no <JUDGE> segment at all: its artifacts contain every grader.
    """
    *parts, name = p.split("/")
    if parts and parts[0] in JUDGE_INVARIANT_FAMILIES:
        return os.path.join(base, *parts, name)
    return os.path.join(base, *parts, JUDGE, name)

def f0(p): return _jp(L0F, p)
def t0(p): return _jp(L0T, p)
def f5(p): return _jp(L5F, p)
def fm(p): return os.path.join(METH, p)   # schematics: no view, no judge, no family

NAVY  = RGBColor(0x1F,0x3A,0x5F); PTO = RGBColor(0x00,0x72,0xB2)
GRPO  = RGBColor(0xE6,0x9F,0x00); GREY= RGBColor(0x5A,0x5A,0x5A)
LIGHT = RGBColor(0xEF,0xF2,0xF6); GREEN=RGBColor(0x00,0x8A,0x63)
RED   = RGBColor(0xC0,0x4A,0x1A); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x22,0x28,0x30); ROWALT=RGBColor(0xF4,0xF7,0xFA)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_N = 0

def slide(fnum=True):
    global _N
    s = prs.slides.add_slide(BLANK)
    if fnum:
        _N += 1
        tb = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"Exp3 · {_N}"; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xB0,0xB8,0xC0); r.font.name = FONT
    return s

def box(s,l,tp,w,h):
    tb = s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tb.text_frame.word_wrap = True; return tb

def rect(s,l,tp,w,h,color,line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp

def set_runs(para, segs, size, align=PP_ALIGN.LEFT):
    para.alignment = align
    for txt,bold,color in segs:
        r = para.add_run(); r.text = txt; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

def title_bar(s, text, kicker=None):
    rect(s,0,0,13.333,1.15,NAVY); rect(s,0,1.15,13.333,0.06,PTO)
    tb = box(s,0.55,0.18,12.2,0.9); tf = tb.text_frame
    p = tf.paragraphs[0]
    if kicker:
        set_runs(p,[(kicker,True,RGBColor(0x8F,0xC7,0xEC))],13)
        set_runs(tf.add_paragraph(),[(text,True,WHITE)],26)
    else:
        set_runs(p,[(text,True,WHITE)],27)

def bullets(s, items, left=0.6, top=1.5, width=12.1, size=16, gap=8):
    tb = box(s,left,top,width,5.6); tf = tb.text_frame
    first = True
    for segs,level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level = level; p.space_before = Pt(gap); p.space_after = Pt(2)
        mk = "•  " if level==0 else "–  "
        set_runs(p,[("    "*level+mk, False, PTO if level==0 else GREY)]+segs, size)

def fit(img,mw,mh):
    with Image.open(img) as im: iw,ih = im.size
    ar = iw/ih; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    return w,h

def figure(s,img,mw,mh,left,top,anchor="center"):
    """Place img fitted inside the (mw × mh) band at (left, top). Returns the image's bottom edge."""
    w,h = fit(img,mw,mh)
    dy = 0.0 if anchor=="top" else (mh-h)/2
    s.shapes.add_picture(img,Inches(left+(mw-w)/2),Inches(top+dy),Inches(w),Inches(h))
    return top+dy+h

def caption(s,text,left,top,width,color=GREY,size=11.5,align=PP_ALIGN.CENTER):
    p = box(s,left,top,width,0.5).text_frame.paragraphs[0]
    set_runs(p,[(text,False,color)],size,align)

def side_notes(s, heading, rows, left, top, width, size=12, gap=10):
    """Heading + bulleted rows in a narrow right-hand column next to a figure."""
    p = box(s,left,top,width,0.45).text_frame.paragraphs[0]
    set_runs(p,[(heading,True,NAVY)],14.5)
    tb = box(s,left,top+0.45,width,5.2); tf = tb.text_frame
    first = True
    for segs in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(gap)
        set_runs(p,[("•  ",False,PTO)]+segs,size)

def cards(s, items, top, left=0.6, width=12.1, height=1.2, gap=0.12):
    """Option cards: (title, body, accent colour, tag)."""
    y = top
    for name,desc,col,tag in items:
        rect(s,left,y,width,height,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,left,y,0.14,height,col)
        tb = box(s,left+0.35,y+0.10,width-0.55,height); tf = tb.text_frame
        seg = [(name,True,col)]
        if tag: seg.append(("    ("+tag+")",True,GREEN))
        set_runs(tf.paragraphs[0],seg,14.5)
        set_runs(tf.add_paragraph(),[(desc,False,DARK)],12)
        y += height + gap
    return y

# ---- native pptx table from a markdown table file ----
def _num(x):
    try: float(x); return True
    except: return False

def _fmt(col,v):
    if not _num(v): return v
    fv = float(v)
    if col in ("p","p_holm","wilcoxon_p") and fv == 0: return "<.001"
    if col in ("iteration","target_iter"): return f"{fv:.0f}"
    s = f"{fv:.3f}".rstrip("0").rstrip(".")
    return s if s else "0"

def md_table(s, md_path, left, top, width, height, drop=(), keep=None,
             fontsize=9.0, rename=None, order=None):
    raw = [l for l in open(md_path,encoding="utf-8").read().splitlines() if l.strip().startswith("|")]
    rows = [[c.strip() for c in l.strip().strip("|").split("|")] for l in raw]
    header, body = rows[0], rows[2:]
    if keep: body = [r for r in body if keep(dict(zip(header,r)))]
    if order: body = sorted(body, key=lambda r: order(dict(zip(header,r))))
    idx = [i for i,h in enumerate(header) if h not in drop]
    disp = [(rename or {}).get(header[i],header[i]) for i in idx]
    data = [[_fmt(header[i], r[i]) for i in idx] for r in body]
    if not data:
        raise ValueError(
            f"md_table: no rows left after keep() on {os.path.basename(md_path)} — the filter "
            f"matched nothing. Header={header}. This usually means an EDA re-render changed a "
            f"column's number formatting (e.g. '10.000' -> '10'); fix the keep() predicate.")
    nr, nc = len(data)+1, len(idx)
    gt = s.shapes.add_table(nr,nc,Inches(left),Inches(top),Inches(width),Inches(height)).table
    lens = [max(len(disp[c]), *(len(data[r][c]) for r in range(len(data)))) for c in range(nc)]
    tot = sum(lens)
    for c in range(nc):
        gt.columns[c].width = Inches(width*lens[c]/tot)
    for c,htext in enumerate(disp):
        cell = gt.cell(0,c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_top=Pt(1);cell.margin_bottom=Pt(1);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        r = pr.add_run(); r.text = htext; r.font.size=Pt(fontsize); r.font.bold=True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for ri,row in enumerate(data,1):
        for c,val in enumerate(row):
            cell = gt.cell(ri,c); cell.fill.solid()
            cell.fill.fore_color.rgb = ROWALT if ri%2 else WHITE
            cell.margin_top=Pt(0);cell.margin_bottom=Pt(0);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
            pr = cell.text_frame.paragraphs[0]
            col = DARK
            if disp[c] in ("dz","delta","mean_delta") and val.startswith("-"): col = RED
            rr = pr.add_run(); rr.text = val; rr.font.size=Pt(fontsize)
            rr.font.color.rgb = col; rr.font.name = FONT
    return gt

def grid_table(s, rows, left, top, width, colw, fontsize=12.5, rowh=0.42):
    """rows = list of (cells:list[str], color_or_None, is_header:bool)"""
    y = top
    for ri,(cells, color, is_hdr) in enumerate(rows):
        x = left
        for i,c in enumerate(cells):
            w = width*colw[i]
            rect(s,x,y,w,rowh, NAVY if is_hdr else (LIGHT if ri%2==0 else WHITE))
            p = box(s,x+0.1,y+0.04,w-0.15,rowh).text_frame.paragraphs[0]
            col = WHITE if is_hdr else (color if (color and i==len(cells)-1) else DARK)
            set_runs(p,[(c,is_hdr or i==0,col)],fontsize)
            x += w
        y += rowh + 0.04
    return y

# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.6,13.333,0.08,PTO)
tb = box(s,0.9,1.6,11.5,2.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Looking Ahead in Goal-Oriented Dialogue",True,WHITE)],34)
set_runs(tf.add_paragraph(),[("Preference-Tree (PTO) vs Group-Relative (GRPO) optimization "
                             "of a small Motivational-Interviewing therapist",False,
                             RGBColor(0xBD,0xD6,0xEA))],18)
tb2 = box(s,0.9,4.85,11.5,2.0); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("Results, measurement validity, and the open decisions",True,WHITE)],20)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("The main comparison is complete and has been re-scored end-to-end by a second, "
             "independent judge.",False,RGBColor(0x9F,0xB4,0xC8))],14)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_runs(p,[("Lior Baruch · Reichman University · 3 August 2026",
             False,RGBColor(0x9F,0xB4,0xC8))],13)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Llama-3.2-1B therapist (bf16)  ·  gpt-4o-mini simulated patient + oracle  ·  "
             "Claude Haiku 4.5 held-out judge  ·  96 personas, persona-paired",
             False,RGBColor(0x7E,0x93,0xA8))],12)

# =====================================================================
# 2 · AGENDA
# =====================================================================
s = slide(); title_bar(s,"Where things stand, and what I'd like to decide today","AGENDA")

p = box(s,0.6,1.45,6.0,0.4).text_frame.paragraphs[0]
set_runs(p,[("In hand",True,GREEN)],16.5)
bullets(s,[
  ([("The main comparison is finished. ",True,NAVY),("PTO and GRPO, 10 training iterations each "
    "under matched settings, fully scored on the whole battery over 96 fixed patient personas.",
    False,DARK)],0),
  ([("The measurement instrument is measured, not assumed. ",True,NAVY),
    ("The oracle's own repeatability, plus a complete re-scoring of every conversation by a judge "
     "from a different model family that never played the patient and took no part in training.",
     False,DARK)],0),
  ([("The training signal itself is measurable — for both methods. ",True,NAVY),
    ("What each update actually rewards, on one scale, which turns the reward-hacking story from "
     "an inference about outcomes into a direct measurement.",False,DARK)],0),
  ([("Everything is reproducible from one command",False,DARK),(" — figures, tables and summaries "
    "regenerate under either grader, with seeded confidence intervals.",False,DARK)],0),
], left=0.6, top=1.95, width=6.0, size=13.5, gap=11)

rect(s,6.85,1.45,0.02,4.4,RGBColor(0xD5,0xDD,0xE5))
p = box(s,7.15,1.45,5.5,0.4).text_frame.paragraphs[0]
set_runs(p,[("Open — the agenda",True,RED)],16.5)
bullets(s,[
  ([("How to tell the story. ",True,NAVY),("Method paper, look-ahead paper, or MI paper — the "
    "choice sets both the venue and the shape of the thesis chapter.",False,DARK)],0),
  ([("Whether to finish the look-ahead (K=5) arms",True,NAVY),(", which are paused on API budget "
    "and are the one comparison still open.",False,DARK)],0),
  ([("What else is worth running",True,NAVY),(" — in particular human MI-coder validation, which "
    "is the one remaining validity gap that money cannot close.",False,DARK)],0),
  ([("Scope. ",True,NAVY),("What goes in a paper, what stays in the thesis, and in what order they "
    "get written.",False,DARK)],0),
], left=7.15, top=1.95, width=5.5, size=14.5, gap=18)

rect(s,0.6,6.1,12.1,0.85,LIGHT); rect(s,0.6,6.1,0.06,0.85,PTO)
set_runs(box(s,0.85,6.22,11.7,0.7).text_frame.paragraphs[0],
  [("In one line: ",True,NAVY),("the PTO-vs-GRPO result is finished and now survives a grader that "
    "took no part in training. What is left to decide is how we frame it, and whether the "
    "look-ahead comparison is worth buying.",False,DARK)],13.5)

# =====================================================================
# 3 · REMINDER — THE PUBLISHED PAPER (Exp1)
# =====================================================================
s = slide(); title_bar(s,"The starting point — the ICLR 2025 paper","BACKGROUND · 1")

rect(s,0.6,1.45,12.1,1.12,LIGHT)
tb = box(s,0.85,1.55,11.6,1.0); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Preference Tree Optimization: Enhancing Goal-Oriented Dialogue with "
                           "Look-Ahead Simulations",True,NAVY)],17)
p = tf.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Baruch, Butman, Bar, Friedman · ICLR 2025",False,GREY)],13)

bullets(s,[
  ([("The method. ",True,NAVY),("At each therapist turn, branch several candidate replies; for each "
    "candidate simulate ",False,DARK),("K further turns",True,DARK),(" of the conversation; let the "
    "oracle score the resulting trajectory; keep the best and worst as a preference pair; update "
    "with DPO. Repeat, regenerating the tree from the improved model each round.",False,DARK)],0),
  ([("Why look ahead. ",True,NAVY),("Scoring a reply on its own rewards openings that look good in "
    "isolation. Scoring the reply plus K simulated turns rewards openings that lead somewhere.",
    False,DARK)],0),
  ([("Setup. ",True,NAVY),("Llama-2-7B therapist, GPT-3.5 as both simulated patient and oracle, "
    "96 patient personas, K ∈ {0, 5}, 7 iterations, reward = mean(Q1, Q2).",False,DARK)],0),
  ([("Headline finding. ",True,NAVY),("Every PTO model beat the untrained baseline on session "
    "satisfaction and working alliance, and ",False,DARK),("K = 5 gave higher and more stable "
    "scores than K = 0",True,DARK),(".",False,DARK)],0),
], top=2.85, size=15, gap=15)

# =====================================================================
# 4 · WHAT CHANGED SINCE
# =====================================================================
s = slide(); title_bar(s,"What changed since the paper","BACKGROUND · 2")

grid_table(s,[
  (["","Exp1 — ICLR paper","Exp2","Exp3 — current"],None,True),
  (["Therapist","Llama-2-7B","Llama-3.2-1B, 4-bit","Llama-3.2-1B, bf16"],None,False),
  (["Patient + oracle","GPT-3.5","gpt-4o-mini","gpt-4o-mini"],None,False),
  (["Patient personas","cooperative","less cooperative","less cooperative"],None,False),
  (["Oracle output","regex-parsed","JSON schema","JSON schema"],None,False),
  (["Evaluation","Q1, Q2","6 questionnaires","8 metrics × 2 independent judges"],None,False),
  (["Methods","PTO","PTO (4 oracles) + weak GRPO","PTO vs GRPO, matched"],None,False),
  (["Iterations","7","7 per arm","10 per arm (K=0)"],None,False),
],0.6,1.45,12.1,[0.16,0.24,0.28,0.32],fontsize=12,rowh=0.36)

p = box(s,0.6,5.15,12.1,0.4).text_frame.paragraphs[0]
set_runs(p,[("Each experiment is a fresh re-implementation — no data carries over between them.",
             True,NAVY)],13.5)

rect(s,0.6,5.62,12.1,1.25,RGBColor(0xFC,0xF3,0xEC))
rect(s,0.6,5.62,0.06,1.25,RED)
tb = box(s,0.85,5.72,11.7,1.1); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("One caveat worth stating up front — ",True,RED),
                           ("Exp2 and Exp3 absolute scores are not on the same axis.",True,DARK)],13)
p = tf.add_paragraph(); p.space_before = Pt(3)
set_runs(p,[("Same therapist model, but Exp2 generated its conversations in 4-bit and Exp3 in bf16. "
             "4-bit produces far more degenerate, looping therapist turns (≈9.5% of turns vs "
             "≈0.3%), which the oracle floors — so Exp2's base model scores 2.38 and Exp3's "
             "scores 3.00 for the same model. The non-degenerate Exp2 subset scores ≈2.93. "
             "Compare within an experiment, not across.",False,DARK)],12)

# =====================================================================
# 5 · WHAT WAS RUN
# =====================================================================
s = slide(); title_bar(s,"What was run in this experiment","SETUP")

bullets(s,[
  ([("Task. ",True,NAVY),("A 1B therapist model converses with a simulated patient (96 fixed "
    "personas). A larger oracle model then grades the full conversation on validated MI "
    "questionnaires.",False,DARK)],0),
  ([("Training signal. ",True,NAVY),("Oracle score on Q1+Q2 only. All other questionnaires are "
    "held out for evaluation.",False,DARK)],0),
  ([("Methods compared. ",True,NAVY),("PTO",True,PTO),(" (branch each therapist turn → oracle → "
    "preference pairs → DPO update) vs ",False,DARK),("GRPO",True,GRPO),
    (" (group-relative update over the same generations). Matched hyper-parameters, matched "
     "look-ahead K, matched min-conversation-length filter.",False,DARK)],0),
  ([("Loop. ",True,NAVY),("Each iteration regenerates all conversations from the current policy; "
    "those same conversations are the evaluation set.",False,DARK)],0),
  ([("Evaluation. ",True,NAVY),("Full-conversation scores on eight oracle instruments — ",
    False,DARK),("Q1, Q2, WAI-SR, CSQ-8, MI-SAT, MITI",True,DARK),(" (Q1 and Q2 are scored "
    "separately and reported throughout as the ",False,DARK),("Q1+Q2",True,DARK),(" composite), "
    "plus ",False,DARK),("PCT",True,DARK),(" (patient change-talk) and ",False,DARK),
    ("MICI",True,DARK),(" (MI-inconsistent therapist behaviour, lower is better) — and three "
    "ratios derived from the MITI behaviour counts, ",False,DARK),("R:Q / %CR / %MICO",True,DARK),
    (". Every conversation is graded twice: once by the training oracle, once by a held-out "
     "judge.",False,DARK)],0),
], top=1.45, size=15.5, gap=13)

rect(s,0.6,5.35,12.1,0.03,RGBColor(0xD5,0xDD,0xE5))
grid_table(s,[
  (["Arm","Iterations trained","Iterations scored","Status"],None,True),
  (["PTO   K=0","1–10","base + 1–10","complete, both judges"],GREEN,False),
  (["GRPO  K=0","1–10","base + 1–10","complete, both judges"],GREEN,False),
  (["PTO   K=5","1–5","base + 1–5","paused after iteration 5"],RED,False),
  (["GRPO  K=5","1","base + 1","paused"],RED,False),
],0.6,5.55,12.1,[0.16,0.19,0.19,0.46],fontsize=12,rowh=0.33)

# =====================================================================
# 6 · METHOD — the PTO loop  (redraw of the ICLR paper's Figure 1)
# =====================================================================
s = slide(); title_bar(s,"How PTO works — one training iteration","METHOD · 1")
figure(s,fm("pto_framework.png"),12.4,5.4,0.45,1.4)
caption(s,"Figure 1 of the ICLR 2025 paper, redrawn with this experiment's models.",
        0.45,6.95,12.4,size=11.5)

# =====================================================================
# 7 · METHOD — the GRPO loop, same picture
# =====================================================================
s = slide(); title_bar(s,"How GRPO works — the same loop, one method over","METHOD · 2")
figure(s,fm("grpo_framework.png"),12.4,5.4,0.45,1.4)
caption(s,"The same diagram for GRPO — no counterpart exists in the paper, which is PTO-only.",
        0.45,6.95,12.4,size=11.5)

# =====================================================================
# 8 · METHOD — one PTO branch point  (redraw of the ICLR paper's Figure 2)
# =====================================================================
s = slide(); title_bar(s,"PTO — how a single training pair is made","METHOD · 3")
figure(s,fm("pto_preference_tree.png"),5.1,5.85,0.35,1.28,anchor="top")
side_notes(s,"Figure 2 of the paper, in Exp3's greedy mode",[
  [("At each therapist turn the policy samples ",False,DARK),("M = 8 candidate replies",True,DARK),
   (".",False,DARK)],
  [("Each candidate is rolled forward ",False,DARK),("K further turns",True,PTO),
   (" — patient, therapist, patient — ",False,DARK),("before",True,DARK),(" the oracle sees it. "
   "That rollout is the entire look-ahead lever: K changes what the oracle scores, and nothing "
   "else in the method.",False,DARK)],
  [("Of the 8 scored candidates exactly two survive — best and worst — and only if they are far "
    "enough apart. ",False,DARK),("A tie at a branch point produces no training data at all.",
    True,DARK)],
  [("The winner is then ",False,DARK),("appended to the trunk",True,PTO),(", so the tree is one "
   "conversation growing under its own decisions rather than eight independent ones. This is what "
   "makes the search sequential.",False,DARK)],
  [("The pair that leaves this picture — (trunk, chosen, rejected) — is the only thing DPO ever "
    "sees. The scores themselves are not part of the loss.",False,DARK)],
],5.75,1.35,7.15,size=12,gap=11)

# =====================================================================
# 9 · METHOD — one GRPO group
# =====================================================================
s = slide(); title_bar(s,"GRPO — how a single update is made","METHOD · 4")
figure(s,fm("grpo_group_rollout.png"),5.1,5.85,0.35,1.28,anchor="top")
side_notes(s,"Deliberately the same rows as the previous slide",[
  [("Everything above the oracle is identical",True,NAVY),(" — same policy, same G = 8 samples, "
    "same K-turn look-ahead, same full-trajectory scoring. The methods diverge only in what "
    "happens to the scores.",False,DARK)],
  [("Nothing is discarded. ",True,GRPO),("All G completions carry a gradient, weighted by the "
    "group-relative advantage — how far above or below the group mean each one scored.",
    False,DARK)],
  [("There is no threshold, ",True,GRPO),("so every prompt produces an update. PTO can spend a "
    "whole branch point and emit nothing.",False,DARK)],
  [("There is no trunk. ",True,GRPO),("The prompt list is sliced from the rollout at the start of "
    "the iteration and then fixed — a completion winning here does not change which prompt comes "
    "next.",False,DARK)],
  [("That last difference is the substantive one: ",True,NAVY),("PTO searches sequentially, "
    "conditioning each choice on the previous one; GRPO improves the policy against a fixed set "
    "of situations.",False,DARK)],
],5.75,1.35,7.15,size=12,gap=11)

# =====================================================================
# 10 · MAIN FIGURE — trajectories
# =====================================================================
s = slide(); title_bar(s,"All metrics across 10 iterations (K = 0)","RESULTS · 1")
figure(s,f0("0_headline/trajectories_all_metrics.png"),12.3,5.15,0.5,1.35)
caption(s,"Mean over the 96 personas at each iteration; shaded band = 95% CI. "
          "Iteration 0 = untrained base model. MICI is inverted in meaning (lower = fewer "
          "MI-inconsistent therapist behaviours).",0.5,6.65,12.3)

# =====================================================================
# 7 · ENDPOINT NUMBERS
# =====================================================================
s = slide(); title_bar(s,"Endpoint numbers (K = 0)","RESULTS · 2")
caption(s,"Two endpoints reported for each arm: the matched final iteration (10) and each arm's "
          "own best iteration on its training rubric (Q1+Q2).",
        0.5,1.32,12.3,size=12.5,align=PP_ALIGN.LEFT)
md_table(s,t0("0_headline/leaderboard_scorecard.md"),0.5,1.85,12.3,1.9,
         fontsize=10.5,
         rename={"target":"endpoint","iteration":"iter",
                 "WAI-SR (Working Alliance)":"WAI-SR","CSQ-8 (Client Satisfaction)":"CSQ-8",
                 "MI-SAT (MI Satisfaction)":"MI-SAT","MITI (MI Integrity)":"MITI",
                 "PCT (Patient Change-Talk)":"PCT","MICI (MI-Inconsistency) ↓":"MICI ↓",
                 "Reflection:Question (MITI)":"R:Q","% Complex Reflections (MITI)":"%CR",
                 "% MI-Consistent (MITI)":"%MICO"})

rect(s,0.5,4.1,12.3,0.03,RGBColor(0xD5,0xDD,0xE5))
p = box(s,0.5,4.25,12.3,0.4).text_frame.paragraphs[0]
set_runs(p,[("Paired contrasts on the 96 shared personas (PTO − GRPO), Holm-corrected",True,NAVY)],14)
grid_table(s,[
  (["Contrast","Q1+Q2 Δ","dz","p (Holm)"],None,True),
  (["PTO iter-10  vs  GRPO iter-10  (matched final)","+0.51","0.73","< .001"],None,False),
  (["PTO iter-10  vs  GRPO iter-8  (each at its own best)","+0.18","0.30",".010"],None,False),
],0.5,4.75,8.6,[0.55,0.16,0.13,0.16],fontsize=12,rowh=0.36)

p = box(s,0.5,6.05,12.3,1.0).text_frame.paragraphs[0]
set_runs(p,[("Trajectory shape: ",True,NAVY),
            ("PTO Q1+Q2 rises 3.00 → 4.26 with its maximum at the last iteration (OLS slope "
             "0.120/iter). GRPO rises 3.07 → 4.08 by iteration 8, then falls to 3.75 by "
             "iteration 10 (slope 0.072/iter). Every arm × rubric effect vs base is significant "
             "at Holm p < .001.",False,DARK)],12.5)

# =====================================================================
# 8 · EFFECT VS BASE
# =====================================================================
s = slide(); title_bar(s,"Effect vs the untrained base model (K = 0)","RESULTS · 3")
figure(s,f0("0_headline/effect_vs_base_forest_final.png"),6.05,4.9,0.35,1.45)
figure(s,f0("0_headline/effect_vs_base_forest_best.png"),6.05,4.9,6.75,1.45)
caption(s,"at the matched final iteration (10)",0.35,6.45,6.05,color=NAVY,size=12.5)
caption(s,"at each arm's own best iteration (PTO 10, GRPO 8)",6.75,6.45,6.05,color=NAVY,size=12.5)
caption(s,"Paired Cohen's dz vs base over the 96 personas, with 95% CI. Positive = improvement, "
          "except MICI where positive = more MI-inconsistent behaviour.",0.35,6.85,12.45,size=11.5)

# =====================================================================
# 9 · WHAT MOVED ALONGSIDE THE SCORES
# =====================================================================
s = slide(); title_bar(s,"What else moved as the scores rose (K = 0)","RESULTS · 4")
bot = figure(s,f0("0_headline/reward_hack_panel.png"),7.85,5.15,0.35,1.5,anchor="top")
caption(s,"Global rubric score (left axis) vs MI-inconsistent behaviour and patient change-talk "
          "(right axis), per iteration",0.35,bot+0.1,7.85,size=11.5)

side_notes(s,"Measured alongside the gains",[
  [("MICI",True,DARK),(" (MI-inconsistent behaviour) rises 0.21 at base → ",False,DARK),
   ("0.49",True,PTO),(" PTO and ",False,DARK),("0.84",True,GRPO),(" GRPO at iter 10 "
   "(GRPO 0.54 at its iter-8 peak).",False,DARK)],
  [("GRPO collapses to 0.15 questions per turn by iter 10, against PTO's 0.55; "
    "reflection-to-question ratio 1.44 vs 0.75.",False,DARK)],
  [("Both arms drift toward affirmation-heavy language, more so in GRPO's late iterations.",
    False,DARK)],
  [("The reward's own composition is part of it: the Q2 items that move most reward therapist "
    "self-disclosure, which MI does not prescribe.",False,DARK)],
  [("Whether this is really reward-hacking, or just a harsher rubric, is what the held-out judge "
    "settles — next section.",False,NAVY)],
], 8.35, 1.5, 4.6, size=12.5)

# =====================================================================
# 10 · MICI DETAIL
# =====================================================================
s = slide(); title_bar(s,"MI-inconsistency, decomposed (K = 0)","RESULTS · 5")
figure(s,f0("0_headline/mici_detail_grid.png"),12.3,5.15,0.5,1.35)
caption(s,"Per-conversation rate of each MI-inconsistent therapist behaviour, by iteration. "
          "Lower is better throughout.",0.5,6.65,12.3)

# =====================================================================
# 11 · VALIDITY 1 — the instrument is measured
# =====================================================================
s = slide(); title_bar(s,"The instrument itself is now measured, not assumed",
                       "MEASUREMENT VALIDITY · 1")
figure(s,f0("8_measurement/oracle_repeatability_icc.png"),11.6,2.55,0.85,1.32,anchor="top")
caption(s,"ICC(2,1) per model and metric across four independent scorings of the same "
          "conversations. Dotted lines: Koo & Li \"good\" (0.75) and \"excellent\" (0.90).",
        0.85,3.92,11.6,size=11)

grid_table(s,[
  (["Metric","Training oracle (gpt-4o-mini)","Held-out judge (Haiku 4.5)"],None,True),
  (["Q1","ICC 0.982 – 0.994","ICC 0.951 – 0.978"],None,False),
  (["Q2","ICC 0.955 – 0.992","ICC 0.938 – 0.963"],None,False),
  (["MICI","ICC 0.864 – 0.943","ICC 0.525 – 0.929"],RED,False),
],0.6,4.42,6.55,[0.17,0.43,0.40],fontsize=11.5,rowh=0.32)
caption(s,"Repeated scorings of 4 anchor model states × 96 conversations. Mean |Δ| between "
          "repeats 0.04–0.09 — the project's informal \"oracle noise ≈ 0.10\" was a conservative "
          "upper bound, and it shrinks by ~√96 at the arm-mean level everything here reports.",
        0.6,5.85,6.55,size=10.5,align=PP_ALIGN.LEFT)

side_notes(s,"And a second grader over the whole grid",[
  [("Claude Haiku 4.5",True,NAVY),(" — different model family, never played the patient, never "
    "touched training — re-scored ",False,DARK),("22,272 of 22,272 cells",True,DARK),
   (" (29 model states × 8 rubrics × 96 conversations).",False,DARK)],
  [("Deliberately 1 repeat, not 3: judge noise adds ≈0.01 to a 96-conversation arm mean against "
    "≈0.09 from persona sampling, so breadth beats depth at equal cost.",False,DARK)],
  [("Cost ",False,DARK),("$42",True,DARK),(" for the sweep plus ",False,DARK),("$9",True,DARK),
   (" for the repeatability repeats, using Anthropic's batch API.",False,DARK)],
],7.5,4.35,5.3,size=11.5,gap=9)

# =====================================================================
# 12 · VALIDITY 2 — the rankings survive
# =====================================================================
s = slide(); title_bar(s,"Does the result survive a judge that never played the patient? Yes",
                       "MEASUREMENT VALIDITY · 2")
figure(s,f0("8_measurement/judge_contrast_preservation.png"),11.9,4.05,0.7,1.35,anchor="top")
caption(s,"Endpoint contrasts computed separately under each grader. Same-sign bars mean the "
          "claim does not depend on the grader also having simulated the patient.",
        0.7,5.42,11.9,size=11)

grid_table(s,[
  (["All arm × metric contrasts","|Δ| ≥ 0.10","|Δ| ≥ 0.25","|Δ| ≥ 0.50","Anchor contrasts"],None,True),
  (["88.3%  (1,632 / 1,848)","94.1%","97.0%","98.9%","18 / 18"],GREEN,False),
],0.6,5.88,12.1,[0.28,0.15,0.15,0.15,0.27],fontsize=12,rowh=0.34)
caption(s,"Share of contrasts keeping their sign under the held-out judge. The two graders "
          "disagree only about differences too small to claim — and the held-out judge actually "
          "widens the headline PTO − GRPO Q1 gap, to +0.77 against the primary's +0.53.",
        0.6,6.62,12.1,size=11.5,align=PP_ALIGN.LEFT)

# =====================================================================
# 13 · VALIDITY 3 — gain retention
# =====================================================================
s = slide(); title_bar(s,"The sharpest evidence for reward-hacking — gain retention",
                       "MEASUREMENT VALIDITY · 3")
figure(s,f0("8_measurement/multijudge_retention_trajectory.png"),7.2,4.55,0.4,1.5,anchor="top")
caption(s,"Share of each arm's measured gain that the held-out judge also sees, by iteration.",
        0.4,6.15,7.2,size=11)

side_notes(s,"Read it as a train / test ratio",[
  [("The quantity is ",False,DARK),("Δ under the held-out judge ÷ Δ under the judge that was the "
    "training reward",True,DARK),(". Near 1.0 = a real gain; falling = a policy fitting the grader "
   "it was trained against.",False,DARK)],
  [("At iteration 10, Q1 retention is ",False,DARK),("PTO 0.80 [0.68, 0.93]",True,PTO),
   (" against ",False,DARK),("GRPO 0.28 [0.06, 0.43]",True,GRPO),(" — non-overlapping intervals.",
   False,DARK)],
  [("In plain terms: under a grader that never took part in training, GRPO's net 10-iteration Q1 "
    "gain is about ",False,DARK),("0.19 points, not the 0.68",True,DARK),(" the training oracle "
   "credits it with.",False,DARK)],
  [("It is an onset, not an endpoint accident. The arms are indistinguishable for three "
    "iterations, then PTO holds 0.80–0.98 for the whole run while GRPO decays to 0.28.",False,DARK)],
  [("Not a scale artifact: every Q2 retention interval overlaps (0.80–0.85). Only Q1 separates.",
    False,DARK)],
],8.0,1.5,4.85,size=12,gap=11)

# =====================================================================
# 14 · VALIDITY 4 — where they disagree
# =====================================================================
s = slide(); title_bar(s,"Where the two graders do not agree — stated as limitations",
                       "MEASUREMENT VALIDITY · 4")
figure(s,f0("8_measurement/multijudge_variance_decomposition.png"),8.0,2.90,2.65,1.28,anchor="top")
caption(s,"Share of arm-mean variance by source. A large grader-level slice is harmless — it "
          "cancels in every contrast. The arm × judge slice is the only one that can invalidate a "
          "claim, and it is 1.2–6.9% on every metric.",2.65,4.24,8.0,size=11)

bullets(s,[
  ([("MITI is the exception, and it is a thesis limitation. ",True,RED),
    ("Only 3.6% of MITI's arm-mean variance is between-arm signal — 94.5% is grader level. "
     "Dependability from one judge is 0.65, and it keeps its sign on only 77.5% of contrasts, the "
     "worst of the eight. ",False,DARK),
    ("MITI arm differences are reported as provisional.",True,DARK)],0),
  ([("MICI agrees weakly across graders (r 0.20–0.55), ",True,RED),("and the held-out judge's own "
    "MICI repeatability falls exactly where the sycophancy claim lives (0.53 on GRPO at iter 10). "
    "So that claim is made at the ",False,DARK),("contrast level, not as a precise rate",
    True,DARK),(" — gain retention is the load-bearing evidence.",False,DARK)],0),
  ([("Never average the two graders. ",True,NAVY),("The primary was the training reward and the "
    "second is held out — train-vs-test, not two raters. The level offset is 1.2–1.7 points and "
    "model-dependent, so averaging would silently shrink every effect.",False,DARK)],0),
  ([("Still uncovered: no human MI/MITI-coder validation. ",True,RED),("A judge can be perfectly "
    "repeatable and consistently wrong, and two LLM judges agreeing does not fix that — the "
    "largest remaining validity gap, and it costs time, not budget.",False,DARK)],0),
], top=4.86, size=12.5, gap=9)

# =====================================================================
# 15 · K=5 — WHAT EXISTS
# =====================================================================
s = slide(); title_bar(s,"Look-ahead (K = 5): exactly what exists today","THE OPEN COMPARISON")

grid_table(s,[
  (["Arm","Trained","Scored","Missing"],None,True),
  (["PTO  K=5","iters 1–5","base + iters 1–5  (complete)","iters 6–10"],GREEN,False),
  (["GRPO K=5","iter 1","base + iter 1","iters 2–10"],RED,False),
],0.5,1.4,12.3,[0.13,0.15,0.28,0.44],fontsize=12,rowh=0.34)

bot5 = figure(s,f5("7_stats/k_trajectory_Q1Q2.png"),7.5,3.85,0.35,2.6,anchor="top")
side_notes(s,"The first matched-budget K comparison",[
  [("The K=5 arm's last trained adapter has now been generated and scored on both graders, so "
    "K=0 and K=5 can be compared at ",False,DARK),
   ("equal iteration count",True,DARK),(".",False,DARK)],
  [("K=5 trails K=0 by 0.08–0.16 on Q1+Q2 through iterations 1–4 — small (dz ≤ 0.20) and ",False,DARK),
   ("never significant",True,DARK),(" — then ties at iteration 5 (4.016 vs 4.014).",False,DARK)],
  [("Under the held-out judge there is ",False,DARK),("no tie",True,RED),
   (": K=0 leads at every iteration, at iteration 5 by +0.173 (Holm p = .017).",False,DARK)],
  [("So: ",False,DARK),("\"K=5 never wins\" holds under both graders",True,DARK),
   ("; \"they converge\" is the primary oracle's picture only.",False,DARK)],
  [("GRPO K=5 still has one trained iteration — this is within PTO, not a K × method result.",
    False,GREY)],
],8.2,2.5,4.75,size=11.5,gap=8)

rect(s,8.2,6.25,4.75,0.85,RGBColor(0xFB,0xEE,0xE6)); rect(s,8.2,6.25,0.06,0.85,RED)
tb2 = box(s,8.42,6.32,4.45,0.78); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("The awkward part: ",True,RED),
                            ("K=5 costs materially more per iteration and buys nothing measurable "
                             "at matched budget. The ICLR paper found the opposite on a 7B model "
                             "over 7 iterations.",False,DARK)],11)
caption(s,"Q1+Q2 across iterations, both K arms of both methods. K=5 lines stop where the arms were "
          "paused — everything right of that is K=0 only.",
        0.35,bot5+0.03,7.5,size=10.5)

# =====================================================================
# 16 · TRAINING SIGNAL 1 — the reward-hack is a loop, not a pull
# =====================================================================
s = slide(); title_bar(s,"Inside the training signal — the drift is a loop, not a pull",
                       "MECHANISM · 1")

caption(s,"Both methods weight the candidates of a group and step along the weighted sum — DPO puts "
          "±1 on chosen/rejected, GRPO uses the standardized advantage. Rescaled to a common size "
          "they become one probe, so for the first time the GRPO side is measurable too.",
        0.6,1.3,12.1,size=12,align=PP_ALIGN.LEFT)

bot6 = figure(s,f0("6_preference/generation_vs_selection.png"),7.6,4.0,0.35,1.85,anchor="top")
side_notes(s,"Two rows, two different questions",[
  [("Top row = what the policy ",False,DARK),("generates",True,DARK),
   (" (mean over every candidate). Bottom row = what the update ",False,DARK),
   ("selects for",True,DARK),(" inside a group.",False,DARK)],
  [("Affirmation: selection pressure runs 0.01 → 0.10, while the generated rate runs ",False,DARK),
   ("0.02 → 0.54",True,RED),(" (GRPO) and 0.04 → 0.57 (PTO).",False,DARK)],
  [("Over-praise reaches ",False,DARK),("74%",True,RED),(" of GRPO's candidates; questions collapse "
    "from 0.71 to 0.06 per turn.",False,DARK)],
  [("So the reward-hack is not one hard pull. It is a small, persistent, same-signed pressure "
    "compounding through an on-policy loop — each iteration branches from an already-more-"
    "effusive policy.",False,DARK)],
  [("By the last iterations the update is choosing between ",False,DARK),
   ("two effusive completions",True,DARK),(", which is why the selection contrast understates "
    "how far the pool has moved.",False,DARK)],
],8.25,1.85,4.7,size=11.5,gap=8)

caption(s,"Training-side measurement — independent of the eval scores. GRPO's dip at iteration 9 "
          "is the same iteration the outcome curves dip.",
        0.35,bot6+0.03,7.6,size=10.5)

# The same finding in words: what the update was actually choosing between, early vs late.
rect(s,0.35,5.45,12.6,1.55,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,0.35,5.45,0.1,1.55,PTO)
p = box(s,0.6,5.5,12.2,0.35).text_frame.paragraphs[0]
set_runs(p,[("The same finding in words — the two completions PTO's update was choosing between",
             True,NAVY)],12.5)
tb = box(s,0.6,5.85,12.2,1.1); tf = tb.text_frame
set_runs(tf.paragraphs[0],
  [("iteration 1    ",True,GREY),("kept: ",True,GREEN),
   ("“That’s really helpful to hear. How about we start by asking you about your "
    "motivations…”      ",False,DARK),
   ("dropped: ",True,RED),("“Great! I’m happy you’ve agreed to work on your motivation…”",
                            False,DARK)],11)
set_runs(tf.add_paragraph(),
  [("iteration 10  ",True,GREY),("kept: ",True,GREEN),
   ("“You’re brave, open, and strong. It takes tremendous courage to choose to explore "
    "and transform…”",False,DARK)],11)
set_runs(tf.add_paragraph(),
  [("                     ",True,GREY),("dropped: ",True,RED),
   ("“You’re courageous and courageous, and your resilience is inspiring. Trusting your "
    "ability to find answers…”",False,DARK)],11)

# =====================================================================
# 17 · TRAINING SIGNAL 2 — loss or exploration?
# =====================================================================
s = slide(); title_bar(s,"Is the PTO–GRPO gap the loss, or the exploration?","MECHANISM · 2")

caption(s,"The two methods differ in two ways at once: they see different candidate pools, and they "
          "weight them differently. Every group logs all its candidates' scores, so the group can be "
          "held fixed and only the weighting rule swapped — which separates the two.",
        0.6,1.3,12.1,size=12.5,align=PP_ALIGN.LEFT)

grid_table(s,[
  (["What is held fixed","Cosine between the two update directions",""],None,True),
  (["As trained — rule and data both differ","0.27  (0.32 corrected for estimation noise)",
    "the headline gap"],GREY,False),
  (["Same groups, PTO's rule swapped for GRPO's","0.91","the rule barely matters"],GREEN,False),
  (["Same groups, GRPO's rule swapped for PTO's","0.99","the rule barely matters"],GREEN,False),
  (["Same rule, each method's own groups","0.36 / 0.27  (0.40 / 0.32 corrected)",
    "the gap survives"],RED,False),
],0.6,2.3,12.1,[0.34,0.40,0.26],fontsize=12.5,rowh=0.42)

rect(s,0.6,4.75,12.1,1.15,LIGHT); rect(s,0.6,4.75,0.06,1.15,PTO)
tb = box(s,0.85,4.85,11.7,1.05); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Read: ",True,NAVY),
  ("swapping DPO for group-relative weighting on the same eight completions barely changes what the "
   "update points at. Holding the rule fixed and swapping the candidates leaves the two methods as "
   "far apart as they ever were.",False,DARK)],13.5)
set_runs(tf.add_paragraph(),[("At matched look-ahead and a shared oracle, the two losses extract "
  "nearly the same direction from the same candidates — ",False,DARK),
  ("the difference between PTO and GRPO is which candidates they generate, not how they weight them.",
   True,NAVY)],13.5)

rect(s,0.6,6.1,12.1,0.9,RGBColor(0xEC,0xF5,0xF1)); rect(s,0.6,6.1,0.06,0.9,GREEN)
set_runs(box(s,0.85,6.2,11.7,0.8).text_frame.paragraphs[0],
  [("Why it matters for the framing: ",True,GREEN),
   ("this makes the result a statement about exploration — preference-tree search versus "
    "group sampling — rather than about DPO versus group-relative PPO. It is a cleaner and more "
    "general claim, and it is the one the data actually supports.",False,DARK)],13)

# =====================================================================
# 18 · DECISION 1 — the framing
# =====================================================================
s = slide(); title_bar(s,"Decision 1 — which story do we tell, and where","DECISION")
caption(s,"The three framings are not exclusive, but one of them has to lead: it decides the venue, "
          "what still needs running, and how the thesis chapter is shaped.",
        0.6,1.32,12.1,size=12.5,align=PP_ALIGN.LEFT)

cards(s,[
  ("A · Search story — preference-tree exploration vs group sampling under an expensive LLM judge",
   "The stability and generalization gap: GRPO is competitive up to its peak, then overshoots into "
   "a grader-specific optimum; PTO sustains gains. Gain retention under a held-out judge is the "
   "evidence. The training-signal analysis sharpens it — the two losses want the same thing from "
   "the same candidates, so the gap is about what each method explores.   In hand.",
   GREEN,"complete today"),
  ("B · Look-ahead story — the ICLR lever, extended to both method families",
   "The most direct continuation of the published work. K=0 and K=5 now meet "
   "at a matched iteration and look-ahead buys nothing — a real but underpowered negative (one "
   "method, 5 iterations). Making it conclusive needs at least one K=5 arm finished.   Needs "
   "budget — see Decision 2.",
   NAVY,"needs the K=5 arms"),
  ("C · MI story — what it takes to train a small model toward genuine MI quality",
   "The most interesting story to a clinical audience, and the one the added metrics were built "
   "for: the questionnaires go up while MI-inconsistent behaviour also goes up. To carry it, the "
   "rubrics need a human MI-coder anchor.   Needs coder time, not API budget.",
   PTO,"needs human validation"),
],2.0)

rect(s,0.6,6.15,12.1,0.85,LIGHT)
set_runs(box(s,0.8,6.27,11.7,0.7).text_frame.paragraphs[0],
  [("My reading: ",True,NAVY),("lead with A because it is finished and defensible, fold B in as "
    "the look-ahead section if the budget is approved, and keep C as the framing of the thesis "
    "chapter rather than of the paper. Happy to be argued out of it.",False,DARK)],13)

# =====================================================================
# 17 · DECISION 2 — what to run next
# =====================================================================
s = slide(); title_bar(s,"Decision 2 — what, if anything, to run next","DECISION · BUDGET")

rect(s,0.6,1.4,12.1,0.95,RGBColor(0xFB,0xEE,0xE6)); rect(s,0.6,1.4,0.06,0.95,RED)
tb = box(s,0.85,1.48,11.6,0.9); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Spent so far: ",True,RED),
  ("≈ $300 on OpenAI (training rollouts, look-ahead, oracle scoring) and ≈ $51 on Anthropic (the "
   "full second-judge sweep plus repeatability).",False,DARK)],13)
set_runs(tf.add_paragraph(),[("Cost scales with the number of candidates scored — prompts × G, or "
  "branch-points × M — times iterations. Prompt caching is already maxed, so the only lever is "
  "call count, not price.",False,DARK)],12)

cards(s,[
  ("1 · The nearly-free look-ahead point",
   "Already done: a generate-only pass with the PTO K=5 iteration-5 adapter, scored on "
   "both graders for $1.33. It bought the first matched-budget K comparison — and the answer so "
   "far is that look-ahead does not pay at equal iteration count.",
   GREEN,"DONE — $1.33"),
  ("2 · Resume one K=5 arm, cost-capped",
   "Halve the candidate count (M/G 8 → 4) and cap at 5–6 iterations — the curves plateau by "
   "iteration 4 anyway. Turns the current one-method, five-iteration negative into a conclusive "
   "one. Keep K and the oracle model fixed, since those are the variables under test.",
   NAVY,"needed for framing B"),
  ("3 · Human MI-coder validation on a sample",
   "A trained coder scores a sample of conversations against the same rubrics. Costs time, not "
   "budget, and is the only thing that closes the \"repeatable but possibly wrong\" gap.",
   PTO,"needed for framing C"),
  ("4 · A different starting model",
   "Train from an instruction-tuned base rather than the raw base. Changes the starting point "
   "substantially and may change how much headroom either method has. Largest cost of the four.",
   GREY,""),
],2.55,height=1.02,gap=0.1)

# =====================================================================
# 18 · NEXT STEPS + WHAT'S ON DISK
# =====================================================================
s = slide(); title_bar(s,"What I'd like to leave with","ASKS")
bullets(s,[
  ([("A decision on the lead framing ",True,NAVY),("(A / B / C) — everything else follows from it.",
    False,DARK)],0),
  ([("A yes or no on resuming one K=5 arm, cost-capped",True,NAVY),(", and if yes, roughly what "
    "budget I should plan against. The matched iteration-5 point is already in — it says look-ahead "
    "does not pay, and the question is whether that is worth confirming properly.",False,DARK)],0),
  ([("A view on human coder validation",True,NAVY),(" — whether it is worth arranging, and who "
    "could do the coding.",False,DARK)],0),
  ([("Agreement on scope",True,NAVY),(" — what goes in a paper, what stays in the thesis, and "
    "whether they are written in parallel.",False,DARK)],0),
], top=1.5, size=16, gap=13)

rect(s,0.6,4.0,12.1,0.03,RGBColor(0xD5,0xDD,0xE5))
p = box(s,0.6,4.15,12.1,0.4).text_frame.paragraphs[0]
set_runs(p,[("Already on disk, in case a question turns out to be cheap to answer",True,NAVY)],14.5)
bullets(s,[
  ([("2,784 scored conversations",True,DARK),(" across 29 model checkpoints, every one graded on "
    "all eight instruments by ",False,DARK),("both",True,DARK),(" judges.",False,DARK)],0),
  ([("Per-item decomposition of every instrument",True,DARK),(", persona-level heterogeneity "
    "splits, and per-candidate generation records from every training iteration — all computed, "
    "no further API calls needed.",False,DARK)],0),
  ([("Figures, tables and written summaries regenerate from one command",True,DARK),(", under "
    "either grader, with seeded confidence intervals.",False,DARK)],0),
], top=4.55, size=13.5, gap=10)

prs.save(OUT)
print("wrote", OUT, f"({_N+1} slides)")
