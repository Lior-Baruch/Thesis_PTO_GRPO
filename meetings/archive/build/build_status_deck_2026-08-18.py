"""Build the GENERAL STATUS deck for the 2026-08-18 supervision meeting.

Copied from build_status_deck.py (the 2026-08-16 status deck) per the README convention —
that script stays frozen as the snapshot of what was presented on its date; this one carries
the current numbers. Written for a ROOM: no "since the email"; every slide states its own claim.

What is new since the 2026-08-16 deck, and why this is not that deck re-dated:

- **GRPO at K=5 exists.** Five iterations trained, scored 0-5 on both graders — and the two GRPO
  arms turn out to be budget-matched within 3% (27.08 vs 27.91 GPU-h). RQ-i is now a genuine
  K x method comparison, and the answer is method-dependent: on GRPO, K=5 leads on the reward at
  matched iterations AND at matched budget; on PTO it never does.
- **The COMPUTE axis.** GPU-hours per (arm, iteration) reconstructed from artifact mtimes.
  PTO reaches iteration 10 for 3.4x less compute than GRPO and scores higher — but at matched
  spend it is worse on MI-inconsistency, because the hack tracks optimization depth.
- **Retention by K (new 2026-08-18).** The L5 gain-retention table was a silent 0-byte file until
  a table-first audit caught it; once rendered, it shows GRPO K=5 retaining its FULL Q1 gain
  under the held-out judge (1.08 vs K=0's 0.73, disjoint) while PTO K=5 retains the same or less.
- **Both drafts were revised 2026-08-18** — an audit pass fixed a factual §5 claim in the CLPsych
  draft, wrote the substitution draft's related work + full channel appendix, and added the
  compute axis to the CLPsych discussion. Both compile at/under the 8-page body limit.
- The budget decision changed: GRPO-at-K=5 is DONE and off the list; the live purchase is the
  replicate draw (~$11.4), with a third grader as the runner-up.

Numbers are owned by the artifacts they are read from — `eda/results/{L0,L5}/` tables, STATUS.md,
and each paper's `NUMBERS.md`. This script only restates them; every quoted cell was re-verified
against its table on 2026-08-18. If a number here disagrees with those, they are right.

Exports .pptx, then convert to .pdf via export_pdf.ps1 (PowerPoint COM).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))       # meetings/build/
REPO = os.path.dirname(os.path.dirname(HERE))           # repo root (meetings/ lives here)
ROOT = os.path.join(REPO, "Exp3_PTO_GRPO")             # the experiment the artifacts come from
L0F  = os.path.join(ROOT, "eda", "results", "L0", "figures")
L0T  = os.path.join(ROOT, "eda", "results", "L0", "tables")
L5F  = os.path.join(ROOT, "eda", "results", "L5", "figures")
L5T  = os.path.join(ROOT, "eda", "results", "L5", "tables")
METH = os.path.join(ROOT, "figures")   # hand-authored method schematics (build_method_figures.py)
OUT  = os.path.join(REPO, "meetings", "2026-08-18", "status_2026-08-18.pptx")

JUDGE = "gpt-4o-mini"
JUDGE_INVARIANT_FAMILIES = {"8_measurement"}

def _jp(base, p):
    *parts, name = p.split("/")
    if parts and parts[0] in JUDGE_INVARIANT_FAMILIES:
        return os.path.join(base, *parts, name)
    return os.path.join(base, *parts, JUDGE, name)

def f0(p): return _jp(L0F, p)
def t0(p): return _jp(L0T, p)
def f5(p): return _jp(L5F, p)
def t5(p): return _jp(L5T, p)
def fm(p): return os.path.join(METH, p)

NAVY  = RGBColor(0x1F,0x3A,0x5F); PTO = RGBColor(0x00,0x72,0xB2)
GRPO  = RGBColor(0xE6,0x9F,0x00); GREY= RGBColor(0x5A,0x5A,0x5A)
LIGHT = RGBColor(0xEF,0xF2,0xF6); GREEN=RGBColor(0x00,0x8A,0x63)
RED   = RGBColor(0xC0,0x4A,0x1A); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x22,0x28,0x30); ROWALT=RGBColor(0xF4,0xF7,0xFA)
AMBER = RGBColor(0xB0,0x77,0x00)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_N = 0

def slide(fnum=True):
    global _N
    s = prs.slides.add_slide(BLANK)
    if fnum:
        _N += 1
        tb = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"Exp3 · {_N}"; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xB0,0xB8,0xC0); r.font.name = FONT
    return s

def box(s,l,tp,w,h):
    tb = s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tb.text_frame.word_wrap = True; return tb

def rect(s,l,tp,w,h,color,line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp

def set_runs(para, segs, size, align=PP_ALIGN.LEFT):
    para.alignment = align
    for txt,bold,color in segs:
        r = para.add_run(); r.text = txt; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

def title_bar(s, text, kicker=None):
    rect(s,0,0,13.333,1.15,NAVY); rect(s,0,1.15,13.333,0.06,PTO)
    tb = box(s,0.55,0.18,12.2,0.9); tf = tb.text_frame
    p = tf.paragraphs[0]
    if kicker:
        set_runs(p,[(kicker,True,RGBColor(0x8F,0xC7,0xEC))],13)
        set_runs(tf.add_paragraph(),[(text,True,WHITE)],26)
    else:
        set_runs(p,[(text,True,WHITE)],27)

def bullets(s, items, left=0.6, top=1.5, width=12.1, size=16, gap=8):
    tb = box(s,left,top,width,5.6); tf = tb.text_frame
    first = True
    for segs,level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level = level; p.space_before = Pt(gap); p.space_after = Pt(2)
        mk = "•  " if level==0 else "–  "
        set_runs(p,[("    "*level+mk, False, PTO if level==0 else GREY)]+segs, size)

def fit(img,mw,mh):
    with Image.open(img) as im: iw,ih = im.size
    ar = iw/ih; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    return w,h

def figure(s,img,mw,mh,left,top,anchor="center"):
    w,h = fit(img,mw,mh)
    dy = 0.0 if anchor=="top" else (mh-h)/2
    s.shapes.add_picture(img,Inches(left+(mw-w)/2),Inches(top+dy),Inches(w),Inches(h))
    return top+dy+h

def caption(s,text,left,top,width,color=GREY,size=11.5,align=PP_ALIGN.CENTER):
    p = box(s,left,top,width,0.5).text_frame.paragraphs[0]
    set_runs(p,[(text,False,color)],size,align)

def side_notes(s, heading, rows, left, top, width, size=12, gap=10):
    p = box(s,left,top,width,0.45).text_frame.paragraphs[0]
    set_runs(p,[(heading,True,NAVY)],14.5)
    tb = box(s,left,top+0.45,width,5.2); tf = tb.text_frame
    first = True
    for segs in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(gap)
        set_runs(p,[("•  ",False,PTO)]+segs,size)

def cards(s, items, top, left=0.6, width=12.1, height=1.2, gap=0.12):
    y = top
    for name,desc,col,tag in items:
        rect(s,left,y,width,height,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,left,y,0.14,height,col)
        tb = box(s,left+0.35,y+0.10,width-0.55,height); tf = tb.text_frame
        seg = [(name,True,col)]
        if tag: seg.append(("    ("+tag+")",True,GREEN))
        set_runs(tf.paragraphs[0],seg,14.5)
        set_runs(tf.add_paragraph(),[(desc,False,DARK)],12)
        y += height + gap
    return y

def grid_table(s, rows, left, top, width, colw, fontsize=12.5, rowh=0.42):
    y = top
    for ri,(cells, color, is_hdr) in enumerate(rows):
        x = left
        for i,c in enumerate(cells):
            w = width*colw[i]
            rect(s,x,y,w,rowh, NAVY if is_hdr else (LIGHT if ri%2==0 else WHITE))
            p = box(s,x+0.1,y+0.04,w-0.15,rowh).text_frame.paragraphs[0]
            col = WHITE if is_hdr else (color if (color and i==len(cells)-1) else DARK)
            set_runs(p,[(c,is_hdr or i==0,col)],fontsize)
            x += w
        y += rowh + 0.04
    return y

def callout(s, left, top, width, height, segs, accent=PTO, fill=LIGHT, size=13):
    rect(s,left,top,width,height,fill); rect(s,left,top,0.06,height,accent)
    set_runs(box(s,left+0.25,top+0.10,width-0.45,height-0.1).text_frame.paragraphs[0], segs, size)

# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.5,13.333,0.08,PTO)
tb = box(s,0.9,1.5,11.5,2.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Looking Ahead in Goal-Oriented Dialogue",True,WHITE)],34)
set_runs(tf.add_paragraph(),[("Preference-Tree (PTO) vs Group-Relative (GRPO) optimization "
                             "of a small Motivational-Interviewing therapist",False,
                             RGBColor(0xBD,0xD6,0xEA))],18)
tb2 = box(s,0.9,4.75,11.5,2.2); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("Project status — all four arms complete, both graders, "
                             "the compute axis, and a K × method answer",True,WHITE)],20)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Every arm is trained and fully scored. Look-ahead's value turns out to be "
             "method-dependent — and visible mostly to the grader that was NOT the reward.",
             False,RGBColor(0x9F,0xB4,0xC8))],14)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_runs(p,[("Lior Baruch · Reichman University · 18 August 2026",
             False,RGBColor(0x9F,0xB4,0xC8))],13)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Llama-3.2-1B therapist (bf16)  ·  gpt-4o-mini simulated patient + oracle  ·  "
             "Claude Haiku 4.5 held-out judge  ·  96 personas, persona-paired",
             False,RGBColor(0x7E,0x93,0xA8))],12)

# =====================================================================
# 2 · AGENDA
# =====================================================================
s = slide(); title_bar(s,"Where the project stands, and what I'd like to decide","AGENDA")

p = box(s,0.6,1.45,6.0,0.4).text_frame.paragraphs[0]
set_runs(p,[("Settled",True,GREEN)],16.5)
bullets(s,[
  ([("All four arms are finished and fully scored. ",True,NAVY),("PTO and GRPO at K ∈ {0,5}, "
    "39 model states, every conversation scored on eight instruments by two independent graders "
    "— 29,952 cells per grader (39 × 8 × 96).",False,DARK)],0),
  ([("The comparison now has a COST axis. ",True,NAVY),("GPU-hours per iteration, reconstructed "
    "from run artifacts. PTO reaches iteration 10 for 3.4× less compute than GRPO — and the two "
    "GRPO arms turn out to be budget-matched within 3%.",False,DARK)],0),
  ([("The look-ahead question is answered, per method. ",True,NAVY),("On GRPO, K=5 wins — reward, "
    "MI-consistency at matched budget, and how much of the gain a held-out judge credits. On PTO "
    "it closes the flattery channel and buys nothing else.",False,DARK)],0),
  ([("Both drafts were revised today",True,NAVY),(" after a table-first audit; both compile at or "
    "under the 8-page body limit with complete related-work sections.",False,DARK)],0),
], left=0.6, top=1.95, width=6.0, size=13.5, gap=11)

rect(s,6.85,1.45,0.02,4.4,RGBColor(0xD5,0xDD,0xE5))
p = box(s,7.15,1.45,5.5,0.4).text_frame.paragraphs[0]
set_runs(p,[("Open — the agenda",True,RED)],16.5)
bullets(s,[
  ([("Replication. ",True,NAVY),("Every contested endpoint is a single 96-conversation draw. A "
    "second draw from 5 adapters costs ~$11 and either retires the objection or changes the "
    "headline. Go?",False,DARK)],0),
  ([("A third grader. ",True,NAVY),("The one purchase that would break the substitution paper's "
    "central tie (the aggregate the two judges disagree on).",False,DARK)],0),
  ([("Scope and order. ",True,NAVY),("Which draft is finished first; what the thesis chapter is "
    "built around.",False,DARK)],0),
  ([("Human MI coders · authorship · venue",True,NAVY),(" — unchanged from last time, still open.",
    False,DARK)],0),
], left=7.15, top=1.95, width=5.5, size=14.5, gap=16)

callout(s,0.6,6.1,12.1,0.9,
  [("In one line: ",True,NAVY),("the therapist learns to flatter; trajectory-level scoring stops "
    "the flattery and the therapist advises instead; which method you optimise with decides "
    "whether look-ahead pays — and only the grader outside the training loop can see most of it.",
    False,DARK)],size=13.5)

# =====================================================================
# 3 · BACKGROUND — the published paper
# =====================================================================
s = slide(); title_bar(s,"The starting point — the ICLR 2025 paper","BACKGROUND")

rect(s,0.6,1.45,12.1,1.12,LIGHT)
tb = box(s,0.85,1.55,11.6,1.0); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Preference Tree Optimization: Enhancing Goal-Oriented Dialogue with "
                           "Look-Ahead Simulations",True,NAVY)],17)
p = tf.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Baruch, Butman, Bar, Friedman · ICLR 2025 (SSI-FM workshop poster)",False,GREY)],13)

bullets(s,[
  ([("The method. ",True,NAVY),("At each therapist turn, branch several candidate replies; for "
    "each, simulate ",False,DARK),("K further turns",True,DARK),("; let the oracle score the "
    "resulting trajectory; keep best and worst as a preference pair; update with DPO. Repeat, "
    "regenerating from the improved model.",False,DARK)],0),
  ([("Why look ahead. ",True,NAVY),("Scoring a reply on its own rewards openings that look good in "
    "isolation. Scoring the reply plus K simulated turns rewards openings that lead somewhere.",
    False,DARK)],0),
  ([("Headline finding. ",True,NAVY),("Every PTO model beat the untrained baseline, and ",
    False,DARK),("K = 5 scored higher and more stably than K = 0",True,DARK),
   (" — on a 7B therapist over 7 iterations.",False,DARK)],0),
  ([("Why that matters today. ",True,AMBER),("On a 1B model over 10 iterations, that result "
    "reproduces on one optimiser and not the other — and where it reproduces, most of the effect "
    "is visible only to a held-out grader. Slides 12–16.",False,DARK)],0),
], top=2.85, size=15, gap=15)

# =====================================================================
# 4 · WHAT WAS RUN — the four arms
# =====================================================================
s = slide(); title_bar(s,"What was run, what it cost, and what is scored","SETUP")

grid_table(s,[
  (["Arm","Trained","Scored (both graders)","GPU-h","Status"],None,True),
  (["PTO  K=0","iterations 1–10","base + 1–10","8.1","complete"],GREEN,False),
  (["GRPO K=0","iterations 1–10","base + 1–10","27.9","complete"],GREEN,False),
  (["PTO  K=5","iterations 1–10","base + 1–10","19.7","complete"],GREEN,False),
  (["GRPO K=5","iterations 1–5","base + 1–5","27.1","complete — budget-matched to K=0"],GREEN,False),
],0.5,1.45,12.3,[0.12,0.16,0.20,0.10,0.42],fontsize=12.5,rowh=0.40)

bullets(s,[
  ([("39 scored model states",True,NAVY),(" — 29,952 cells per grader (39 × 8 × 96), complete on "
    "both, no partial cells.",False,DARK)],0),
  ([("“GRPO K=5 stopped at iteration 5” is a statement about iteration count, not "
    "spend",True,AMBER),(" — 27.08 vs 27.91 GPU-h is a 3% difference. A K=5 optimizer step costs "
    "~1.9× a K=0 step (measured per iteration), so matched-iteration tables hand K=5 double the "
    "compute per cell. Every K contrast in this deck is read on BOTH axes.",False,DARK)],0),
  ([("Matched by construction ",True,NAVY),("across every arm: sampling temperature 1.2, 8 "
    "candidates per branch point, minimum context length 12, the same τ filter, the same DPO β, "
    "the same oracle. Verified against each run's stored metadata, not assumed.",False,DARK)],0),
  ([("Training reward = two rubrics ",True,NAVY),("(session satisfaction + working alliance). "
    "The other six instruments were never optimised against, which is what makes them evidence. "
    "All comparisons are persona-paired on the recovered persona identity.",False,DARK)],0),
], top=3.75, size=13.5, gap=12)

# =====================================================================
# 5 · METHOD — the two loops
# =====================================================================
s = slide(); title_bar(s,"The two methods — same loop, one step different","METHOD")
figure(s,fm("pto_framework.png"),6.1,4.5,0.35,1.45,anchor="top")
figure(s,fm("grpo_framework.png"),6.1,4.5,6.85,1.45,anchor="top")
caption(s,"PTO — branch, score, keep best vs worst, DPO",0.35,6.15,6.1,color=PTO,size=12.5)
caption(s,"GRPO — sample a group, standardize rewards, policy gradient",6.85,6.15,6.1,
        color=GRPO,size=12.5)
callout(s,0.6,6.6,12.1,0.72,
  [("Both regenerate their own training data from the current policy each iteration. ",True,NAVY),
   ("Look-ahead (K) changes only what CONTEXT the oracle grades — never the loss, the "
    "hyperparameters, or the sampling.",False,DARK)],size=12.5)

# =====================================================================
# 6 · MAIN RESULT — trajectories
# =====================================================================
s = slide(); title_bar(s,"All eight instruments across 10 iterations (K = 0)","RESULTS · 1")
bot = figure(s,f0("0_headline/trajectories_all_metrics.png"),12.3,5.2,0.5,1.35,anchor="top")
caption(s,"Mean ± 95% CI over the 96 personas, primary oracle. MICI is lower-is-better. "
          "PTO climbs steadily to iteration 10; GRPO peaks at iteration 8 and regresses.",
        0.5,bot+0.05,12.3,size=11.5)

# =====================================================================
# 7 · ENDPOINT NUMBERS
# =====================================================================
s = slide(); title_bar(s,"Endpoint numbers (K = 0)","RESULTS · 2")
caption(s,"Primary oracle, final and best endpoint. Read off the generated scorecard "
          "(results/L0/tables/0_headline/gpt-4o-mini/) — every cell traces to the score lake.",
        0.6,1.3,12.1,size=12,align=PP_ALIGN.LEFT)
grid_table(s,[
  (["","Q1+Q2","MITI","MICI ↓","R:Q","% MI-consistent"],None,True),
  (["PTO  K=0 @ iter 10","4.26","4.27","0.49","0.75","0.70"],None,False),
  (["GRPO K=0 @ iter 10","3.75","3.92","0.84","1.44","0.83"],None,False),
  (["GRPO K=0 @ iter 8 (its peak)","4.08","4.23","0.54","1.04","0.71"],None,False),
],0.6,1.85,12.1,[0.30,0.14,0.14,0.14,0.14,0.14],fontsize=13,rowh=0.44)

bullets(s,[
  ([("PTO beats GRPO at the matched endpoint",True,NAVY),(" — Q1+Q2 4.26 vs 3.75, persona-paired "
    "+0.51, dz 0.73.",False,DARK)],0),
  ([("Credit GRPO at its peak instead",True,NAVY),(" and PTO still leads (+0.18, dz 0.30, "
    "p = .010) — but the MITI and MICI gaps stop being significant. ",False,DARK),
   ("We claim the weaker version.",True,GREEN)],0),
  ([("GRPO's higher R:Q is not a win. ",True,AMBER),("Its reflection-to-question ratio crosses the "
    "MITI competency bar because it stopped asking questions — the regex question rate falls "
    "0.83 → 0.15 per turn. A ratio is gameable through its denominator.",False,DARK)],0),
  ([("“Best” is itself grader-dependent. ",True,RED),("The held-out judge orders the "
    "arms the same way (PTO@10 2.87 vs GRPO@10 2.26) — but it puts GRPO's best at iteration ",
    False,DARK),("3",True,RED),(", not 8. Peak selection performed on the very signal being "
    "optimised is not a neutral operation. Never average the two graders' raw scores — this is "
    "train-vs-test, and only contrasts combine.",False,DARK)],0),
], top=3.9, size=13.5, gap=13)

# =====================================================================
# 8 · WHAT ELSE MOVED — the reward hack
# =====================================================================
s = slide(); title_bar(s,"What the model actually learned (K = 0)","RESULTS · 3")
bot = figure(s,f0("0_headline/mici_detail_grid.png"),7.9,4.6,0.35,1.4,anchor="top")
side_notes(s,"One channel carries the whole thing",[
  [("Every ",False,DARK),("coercive",True,DARK),(" MI violation FALLS with training — confronting, "
    "warning, directing and judging all go to ~0.",False,DARK)],
  [("Over-praise rises from ",False,DARK),("0.02 to 0.70 per turn",True,RED),(" in GRPO and 0.01 "
    "to 0.30 in PTO, and accounts for more than 100% of the increase in total MI-inconsistency.",
    False,DARK)],
  [("The rubrics measure how the session ",False,DARK),("felt to the client",True,DARK),
   (". Flattery is the cheapest way to move that.",False,DARK)],
  [("So the model did not get broadly worse at MI. It swapped coercion for flattery.",False,DARK)],
],8.55,1.35,4.4,size=12,gap=11)
caption(s,"MI-inconsistent behaviour, decomposed, per therapist turn. Higher is worse.",
        0.35,bot+0.05,7.9,size=11)

# =====================================================================
# 9 · MEASUREMENT VALIDITY
# =====================================================================
s = slide(); title_bar(s,"How much of this survives a grader that took no part in training",
                       "MEASUREMENT VALIDITY")
bot = figure(s,f0("8_measurement/multijudge_gain_retention.png"),7.7,4.5,0.35,1.4,anchor="top")
side_notes(s,"Gain retention = a train/test split on the reward",[
  [("The primary oracle ",False,DARK),("was the training signal",True,DARK),("; Claude Haiku 4.5 "
    "is held out, from another family, and never played the patient. Agreement between them is a "
    "generalisation check, not inter-rater reliability.",False,DARK)],
  [("Oracle repeatability ICC(2,1) ",False,DARK),("0.86–0.99",True,GREEN),
   ("; sign preservation across all arm × metric contrasts rises to 98.9% at |Δ| ≥ 0.50.",
    False,DARK)],
  [("Q1 retention: PTO@10 ",False,DARK),("0.80",True,GREEN),(" vs GRPO@10 ",False,DARK),
   ("0.28",True,RED),(", non-overlapping — GRPO's late gains are increasingly grader-specific. "
    "The retention curve separates the arms from iteration 4, four iterations before the reward "
    "curve turns over.",False,DARK)],
  [("⚠ Standing caveats: MITI dependability is ",False,DARK),("0.65",True,RED),
   (" off one judge, and MICI agreement is weak at the conversation level. Both are limitations "
    "in the write-up, not footnotes.",False,DARK)],
],8.35,1.35,4.6,size=11.5,gap=9)
caption(s,"Fraction of each arm's gain over base that survives the judge swap. ~1.0 = a real "
          "behaviour change; ~0 = a gain that existed only in the optimised grader.",
        0.35,bot+0.05,7.7,size=11)

# =====================================================================
# 10 · NEW — THE COMPUTE AXIS
# =====================================================================
s = slide(); title_bar(s,"The compute axis — an iteration is not a unit of spend","NEW · COMPUTE")
bot = figure(s,f5("7_stats/compute_trajectory.png"),7.4,4.6,0.35,1.4,anchor="top")
side_notes(s,"Reconstructed from run artifacts, not from logs",[
  [("PTO reaches iteration 10 for ",False,DARK),("8.1 GPU-h vs GRPO's 27.9",True,GREEN),
   (" — 3.4× cheaper, and it scores higher. Its preference-tree build runs once per iteration; "
    "GRPO recomputes its reward inside the training loop on every step.",False,DARK)],
  [("But at MATCHED spend (~8 GPU-h, GRPO only at iteration 3), PTO@10 wins the reward "
    "(+0.27, dz 0.53) while being ",False,DARK),("worse on MI-inconsistency",True,RED),
   (" (+0.26, dz 0.90; both replicate held-out): the hack tracks optimization depth, so more "
    "reward per GPU-hour buys more reward-hacking per GPU-hour.",False,DARK)],
  [("For look-ahead, the lever's sign is a ",False,DARK),("function of budget",True,DARK),
   (": K=5 (GRPO) is clearly worse at ~13 GPU-h (dz −0.74), null at 7.8 and 23–27 on the primary "
    "oracle — and ahead at the top budgets only under the held-out judge (dz +0.31–0.33).",
    False,DARK)],
  [("Quote the curve, never one row",True,AMBER),(" — and name the axis (iteration vs budget) on "
    "every K claim; the MICI contrast reverses sign between them.",False,DARK)],
],8.0,1.35,4.95,size=11.5,gap=9)
caption(s,"The primary metric against cumulative GPU-hours instead of iteration; iteration "
          "numbers annotated on the markers. Unequal marker spacing IS the finding.",
        0.35,bot+0.03,7.4,size=11)

# =====================================================================
# 11 · SECTION BREAK — look-ahead
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,3.95,13.333,0.08,PTO)
tb = box(s,1.1,2.5,11.2,1.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Now a K × method comparison",True,RGBColor(0x8F,0xC7,0xEC))],16)
set_runs(tf.add_paragraph(),[("Does look-ahead fix the reward hack?",True,WHITE)],36)
tb2 = box(s,1.1,4.3,11.2,1.6); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("It closes the channel it targets, on both methods. What else it "
                             "buys depends on the optimiser — and on who is grading.",
                             False,RGBColor(0xBD,0xD6,0xEA))],19)
p = tf2.add_paragraph(); p.space_before = Pt(12)
set_runs(p,[("PTO: eleven matched iterations (0–10). GRPO: six (0–5), budget-matched to its K=0 "
             "arm within 3%. Both graders, persona-paired.",False,RGBColor(0x7E,0x93,0xA8))],13)

# =====================================================================
# 12 · LOOK-AHEAD 1 — the channel closes
# =====================================================================
s = slide(); title_bar(s,"Look-ahead closes the channel it targets, and keeps it closed",
                       "LOOK-AHEAD · 1")
bot = figure(s,f5("7_stats/k_overpraise_trajectory.png"),7.4,4.5,0.35,1.4,anchor="top")
side_notes(s,"Exactly what it was supposed to do",[
  [("The two PTO arms are ",False,DARK),("indistinguishable for six iterations",True,DARK),
   (" and then separate for good. Over-praise per session at iteration 10: ",False,DARK),
   ("3.042 (K=0) vs 0.625 (K=5)",True,DARK),(".",False,DARK)],
  [("Persona-paired, n = 96: ",False,DARK),("dz 0.887, p < 10⁻⁴",True,GREEN),
   (". The held-out judge reports a larger gap in the same direction (dz 0.999).",False,DARK)],
  [("Prevention, not delay. ",True,GREEN),("K=5's over-praise stays LOW to the end of the run "
    "— 0.115 → 0.625 across eleven points, while K=0 goes 0.167 → 3.042.",False,DARK)],
  [("Not a denominator artifact — these are ",False,DARK),("raw per-session counts",True,DARK),
   (". K=5 takes MORE therapist turns (14.4 vs 10.2), which would push a rate the other way. "
    "Not data starvation — K=5 trained on more preference groups (6,416 vs 4,935).",False,DARK)],
  [("Replicates on GRPO at its matched iterations: over-praise as a share of MI-inconsistent "
    "acts falls 0.178 → 0.086 (primary) and 0.182 → 0.063 (held-out) at iteration 5.",
    False,DARK)],
],8.0,1.35,4.95,size=11.5,gap=8)
caption(s,"Over-praise per therapist turn, all four arms. K=0 solid, K=5 dashed. Shading marks "
          "the measured divergence onset. Higher is worse.",0.35,bot+0.03,7.4,size=11)

# =====================================================================
# 13 · LOOK-AHEAD 2 — the substitution
# =====================================================================
s = slide(); title_bar(s,"…and the substitution replicates — the reduction does not",
                       "LOOK-AHEAD · 2")

grid_table(s,[
  (["per session, iteration 10 (PTO)","Δ (K=0 − K=5)","dz","p (Holm)"],None,True),
  (["Over-praise","+2.417","+0.887","< 10⁻⁴"],RED,False),
  (["Advice without permission","−0.313","−0.239",".057  n.s."],GREEN,False),
  (["Directing the client","−0.375","−0.389",".0016"],GREEN,False),
  (["ALL MI-inconsistent acts","+1.615","+0.446",".0003"],AMBER,False),
],0.5,1.4,6.6,[0.44,0.22,0.15,0.19],fontsize=12,rowh=0.40)

side_notes(s,"The therapist changed how it fails; whether it fails less depends on who grades",[
  [("Table above is the ",False,DARK),("primary oracle",True,DARK),(", which WAS the training "
    "reward: totals 4.958 (K=0) vs 3.344 (K=5) acts per session.",False,DARK)],
  [("The held-out judge reproduces every component and ",False,DARK),("larger",True,DARK),
   (" — over-praise dz 0.999, advice −0.709, directing −0.468 — but its total is ",False,DARK),
   ("null: +0.531, dz 0.099, p = .167",True,RED),(".",False,DARK)],
  [("So the claim that travels across graders is ",False,DARK),
   ("substitution, not reduction",True,GREEN),(". The significant aggregate is a "
    "primary-oracle-only result and is labelled as one wherever it appears.",False,DARK)],
  [("A second, independently prompted coder agrees on the trade: ",True,NAVY),
   ("in the MITI behaviour counts, ",False,DARK),("persuade",True,DARK),(" rises under K=5 by "
    "0.67 acts/session (dz 0.52 primary) and 2.63 (dz 0.87 held-out), while affirmations fall on "
    "both measures.",False,DARK)],
  [("Counts, not rates. ",True,NAVY),("The per-TURN rate is large under both judges only because "
    "K=5 takes 14.4 therapist turns to K=0's 10.2. Per-session counts are the primary reading.",
    False,DARK)],
  [("And the reward the run optimised still cannot separate the arms at the endpoint (dz −0.096, "
    "p_holm .695 primary; dz 0.308, p_holm .130 held-out — both null).",False,DARK)],
],7.35,1.35,5.55,size=11,gap=6)

bot = figure(s,f5("7_stats/k_mici_composition_grid.png"),6.6,3.05,0.5,3.80,anchor="top")
caption(s,"Per-session counts — no denominator involved.",0.5,bot+0.02,6.6,size=10.5)

# =====================================================================
# 14 · LOOK-AHEAD 3 — where it acts
# =====================================================================
s = slide(); title_bar(s,"Where the intervention acts — and why the pressure relocates",
                       "LOOK-AHEAD · 3")
figure(s,f5("6_preference/k_mechanism_overpraise.png"),4.5,5.4,0.4,1.35,anchor="top")
side_notes(s,"A behaviour has to clear three gates",[
  [("Every training group logs all its candidates and their oracle scores, so we can measure what "
    "the reward ",False,DARK),("selected for",True,DARK),(" — not just what the policy produced.",
    False,DARK)],
  [("Under K=0 the selection weight on over-praise is ~0 for ",False,DARK),("six",True,DARK),
   (" iterations, then jumps — 0.063, 0.034, ",False,DARK),("0.083",True,RED),
   (" at policy iterations 6–8 — and the candidate pool and the evaluation follow with a lag.",
    False,DARK)],
  [("Under K=5 it ",False,DARK),("never clears +0.025",True,GREEN),(" at any iteration; the pool "
    "moves 0.004 → 0.065 and the evaluated rate 0.008 → 0.043. Look-ahead acted on the ",
    False,DARK),("reward",True,DARK),(", not on the policy.",False,DARK)],
  [("Which is also why the substitution is unsurprising. The objective did not change; only the "
    "price of one route went up. A trajectory-extended oracle still rewards a turn that makes the "
    "next five look productive — and ",False,DARK),("giving advice does that too",True,AMBER),
   (". Nothing penalises advice; the intervention penalises praise that does not compound.",
    False,DARK)],
  [("Scale worth holding together: the largest per-iteration selection pressure anywhere in this "
    "chain is ",False,DARK),("0.083",True,DARK),(", while the generated over-praise rate moves ",
    False,DARK),("0.002 → 0.318",True,RED),(". A bias too small to see in one update compounds "
    "through the on-policy loop.",False,DARK)],
],5.35,1.35,7.55,size=11.5,gap=8)

# =====================================================================
# 15 · NEW — K × METHOD
# =====================================================================
s = slide(); title_bar(s,"The K × method answer — look-ahead pays on GRPO, not on PTO",
                       "NEW · K × METHOD")
bot = figure(s,f5("7_stats/k_trajectory_Q1Q2.png"),7.4,4.5,0.35,1.4,anchor="top")
side_notes(s,"Same lever, opposite verdicts",[
  [("On PTO, K=5 never significantly leads on the reward",True,DARK),
   (" — over eleven matched iterations, under either grader.",False,DARK)],
  [("On GRPO, K=5 leads",True,GREEN),(" — at iteration 4 on both graders (Δ 0.115 primary, "
    "0.233 held-out, both Holm-sig.) and at iteration 5 under the held-out judge (Δ 0.311, "
    "dz 0.43).",False,DARK)],
  [("At matched BUDGET the GRPO gain is larger and MICI reverses: ",True,DARK),
   ("K=5@5 vs K=0@10 (27.1 vs 27.9 GPU-h): Q1+Q2 +0.289 (dz 0.36, primary) / +0.540 (dz 0.84, "
    "held-out); MI-inconsistency per turn ",False,DARK),("−0.497 (dz −1.34) / −0.403 (dz −1.23)",
    True,GREEN),(" — far cleaner MI at equal spend.",False,DARK)],
  [("Look-ahead FLIPS which method wins at iteration 5",True,RED),(" — K=0: PTO leads (+0.265, "
    "p .014); K=5: GRPO leads (−0.219, p .005); difference-in-differences dz 0.525, p .0001. ",
    False,DARK),("On the primary oracle the same interaction is null",True,DARK),
   (" — the grader that WAS the reward cannot see it.",False,DARK)],
],8.0,1.35,4.95,size=11.5,gap=8)
caption(s,"Both K arms of both methods, Q1+Q2, mean ± 95% CI. K=5 lines stop at their last "
          "scored iteration (PTO 10, GRPO 5); GRPO's two arms cost the same total GPU-hours.",
        0.35,bot+0.03,7.4,size=11)

# =====================================================================
# 16 · NEW — RETENTION BY K
# =====================================================================
s = slide(); title_bar(s,"New this week — look-ahead decides how REAL the gains are, per method",
                       "NEW · RETENTION × K")

grid_table(s,[
  (["Q1 gain retention (held-out ÷ primary)","K = 5","K = 0","read"],None,True),
  (["GRPO @ iteration 5","1.08  [0.94, 1.27]","0.73  [0.57, 0.92]","disjoint — K=5 fully credited"],GREEN,False),
  (["PTO @ iteration 10  (Q1)","0.72  [0.61, 0.84]","0.80  [0.68, 0.93]","overlapping"],GREY,False),
  (["PTO @ iteration 10  (Q2)","0.56  [0.47, 0.66]","0.85  [0.74, 0.98]","disjoint — K=5 WORSE"],RED,False),
],0.5,1.45,12.3,[0.30,0.24,0.24,0.22],fontsize=12.5,rowh=0.42)

bullets(s,[
  ([("Under K=5, GRPO's gains are fully credited by the judge that never trained it",True,GREEN),
    (" — where its K=0 gains were already leaking by iteration 5 and collapsed to 0.28 by 10. "
    "The retention-space counterpart of the interaction on the previous slide.",False,DARK)],0),
  ([("For PTO, closing the flattery channel bought no transferability",True,NAVY),(" — the same "
    "on Q1, significantly less on Q2. Consistent with substitution: the pressure found another "
    "outlet with the same grader-specific character.",False,DARK)],0),
  ([("How this number almost didn't exist: ",True,AMBER),("the K=5 retention table had rendered "
    "as a silent 0-byte file (a hardcoded reference model the K=5 view excludes). A table-first "
    "audit caught it on 2026-08-18; the export layer now refuses to write a silently-empty "
    "artifact. Next slide.",False,DARK)],0),
], top=3.6, size=13.5, gap=13)

callout(s,0.6,6.35,12.1,0.85,
  [("Caveat: ",True,NAVY),("cross-view comparison against different draws of the identical base "
    "policy; the measured base noise floor (54 same-policy contrasts, max |dz| 0.15) bounds the "
    "draw effect. The Q2 gap is far beyond it.",False,DARK)],size=12.5)

# =====================================================================
# 17 · THE AUDIT SLIDE
# =====================================================================
s = slide(); title_bar(s,"What a table-first audit caught this week — and why it stays on a slide",
                       "HOW WE KEEP CATCHING IT")

cards(s,[
  ("A silent empty artifact",
   "The L5 gain-retention table was 0 bytes for weeks — the producing notebook hardcoded a "
   "reference model the view excludes, and an empty frame serialized to nothing. Fixed, guarded "
   "(empty tables now write an explicit marker), and the un-broken table produced slide 16.",
   RED,"prose never noticed; the file 'existed'"),
  ("A factual claim in the CLPsych draft",
   "§5 said the same three reward items lead the endpoint gains in both arms. The table says "
   "self-disclosure tops only GRPO; warmth tops PTO. Corrected, and added to the paper's ledger "
   "of claims that are easy to get subtly wrong.",
   AMBER,"a summary hardened into a finding"),
  ("A placeholder citation and a miscited paper",
   "One bib entry was literally PLACEHOLDER; another (a conversation-planning paper) was cited "
   "for therapist-behaviour analysis. Both verified against sources and re-homed.",
   GREY,"caught before any reviewer could"),
],1.45,left=0.6,width=12.1,height=1.24,gap=0.16)

bullets(s,[
  ([("Same shape as the two corrections presented on 2026-08-16: ",True,NAVY),
    ("prose about tables drifts; the tables themselves have never been wrong. The discipline — "
    "read the tables cold, then diff the narrative — is now written into the repo's own rules, "
    "and it caught the auditor too (a sign-convention misread, fixed before anything shipped).",
    False,DARK)],0),
], top=5.5, size=13.5, gap=10)

# =====================================================================
# 18 · TWO DRAFTS
# =====================================================================
s = slide(); title_bar(s,"Two drafts, deliberately disjoint — both revised today","WRITE-UP")

cards(s,[
  ("Affirmation Without Inquiry  —  K = 0 only",
   "What the reward hack IS. Revised 2026-08-18: the §5 item-level claim corrected, the "
   "compute axis added to Discussion (3.4× cheaper, hackier at matched budget), the retention "
   "figure now a legible single panel, citations verified. Body ends exactly at p.8; 0 overfull "
   "boxes, 0 undefined references.",
   PTO,"submission-ready pending decisions"),
  ("The Hack Moves  —  PTO only, K ∈ {0,5}",
   "Whether a principled fix removes the hack. Revised 2026-08-18: related work written (was a "
   "scaffold), the full channel table in Appendix A — surfacing the MITI-persuade corroboration "
   "— the retention result added to §4, reproducibility filled from run metadata. Body at "
   "68% of p.7; ~1.3 columns of slack.",
   GREEN,"complete draft"),
],1.45,left=0.6,width=12.1,height=1.55,gap=0.18)

bullets(s,[
  ([("They share no claims. ",True,NAVY),("The split is the K axis: the first varies the optimiser "
    "at K=0; the second varies K within PTO. Each carries a claim→artifact ledger; every number "
    "traces to a generated table, which is what caught this week's error.",False,DARK)],0),
  ([("Neither is a K × method paper — by scope, not by data. ",True,AMBER),("GRPO K=5 is now "
    "trained 1–5 and scored on both graders, so the slide-15 contrast exists and both drafts say "
    "in Limitations that they state their claims outside that window. Where the K × method result "
    "goes — thesis chapter or third paper — is part of Decision 1.",False,DARK)],0),
], top=4.85, size=13.5, gap=13)

# =====================================================================
# 19 · DECISION 1 — scope
# =====================================================================
s = slide(); title_bar(s,"Decision 1 — scope and order","DECISION")
caption(s,"Two drafts exist and the thesis needs a spine. These are not mutually exclusive; the "
          "question is what gets finished first and where the K × method result lives.",
        0.6,1.3,12.1,size=13,align=PP_ALIGN.LEFT)
cards(s,[
  ("Submit the K=0 paper first; look-ahead + K × method become the thesis' core",
   "The reward-hacking result is the more complete and self-contained. The look-ahead story — "
   "substitution within PTO, the method-dependent payoff, the compute axis — is bigger than 8 "
   "pages and is exactly a thesis chapter.",
   PTO,"my read: lowest risk, best fit"),
  ("Submit the substitution paper first",
   "The more novel claim, with a methodological contribution beyond MI (nominate the aggregate "
   "before the intervention; composition breaks grader-invariance). One optimiser, one run per "
   "arm; its own TODO names a third grader as the tie-breaker.",
   AMBER,"higher ceiling, higher risk"),
  ("Fold the K × method result into a third short paper",
   "GRPO K=5's full-credit retention + the interaction only the held-out grader sees is a clean "
   "standalone story — but it competes with the thesis chapter for the same material.",
   GREY,None),
],1.95,left=0.6,width=12.1,height=1.42,gap=0.16)
callout(s,0.6,6.5,12.1,0.8,
  [("Whatever the order: ",True,NAVY),("both drafts are frozen against their ledgers, so finishing "
    "either is now a decisions problem (authors, venue, human coders), not an analysis problem.",
    False,DARK)],size=13)

# =====================================================================
# 20 · DECISION 2 — budget
# =====================================================================
s = slide(); title_bar(s,"Decision 2 — what, if anything, is worth buying","DECISION · BUDGET")
caption(s,"Total API spend to date ≈ $317 and it is the binding constraint. GRPO at K=5 — last "
          "meeting's big question — has been bought and settled. What's left:",
        0.6,1.3,12.1,size=13,align=PP_ALIGN.LEFT)
cards(s,[
  ("A second independent draw from 5 contested adapters   (~$11.4 + ~1 A100-h)",
   "GRPO_LA0 @3/@8/@10 and GRPO_LA5 @4/@5. Every contested endpoint is a single 96-conversation "
   "draw; the only measured noise floor is at the base (54 same-policy contrasts, all null). "
   "Either it retires the endpoint-fragility objection thesis-wide, or GRPO's iteration-9/10 "
   "collapse fails to reproduce and the headline changes. Decisive both ways.",
   GREEN,"recommended"),
  ("A third grader on the substitution paper's aggregate",
   "The two judges agree on every component and disagree on the total. A third family breaks the "
   "tie either way and converts the paper's central hedge into a finding. Priced off the "
   "receipt-calibrated basis before any go.",
   AMBER,"the runner-up"),
  ("Extending GRPO K=5 to iteration 10   (~$118 + 23–34 A100-h)",
   "NOT recommended: it would push that arm to ~50 GPU-h against its K=0 twin's 28, making the "
   "arms LESS comparable on the axis that currently makes them a clean pair.",
   RED,"declined, with the reason on the slide"),
],1.95,left=0.6,width=12.1,height=1.46,gap=0.16)

# =====================================================================
# 21 · ASKS
# =====================================================================
s = slide(); title_bar(s,"What I'd like to leave with","ASKS")
bullets(s,[
  ([("A go / no-go on the replicate draw. ",True,NAVY),("~$11.4; the single highest-value purchase "
    "left, and both papers' framing depends on the endpoint it tests.",False,DARK)],0),
  ([("A decision on order. ",True,NAVY),("Which draft is finished first, and does the K × method "
    "result live in the thesis or in a third paper?",False,DARK)],0),
  ([("A view on the third grader",True,NAVY),(" for the substitution paper's aggregate tie.",
    False,DARK)],0),
  ([("A view on human MI coders. ",True,NAVY),("Two graders from different families agree, but "
    "agreement is not validity, and no MI-trained human has read a single transcript. Even 40–50 "
    "coded sessions would change what we can claim — do we have access to a coder?",False,DARK)],0),
  ([("Co-author list and order",True,NAVY),(" for both drafts, plus venue and target date.",
    False,DARK)],0),
], top=1.6, size=15, gap=17)

callout(s,0.6,6.35,12.1,0.85,
  [("Everything in this deck regenerates from one command. ",True,NAVY),
   ("Figures, tables and summaries rebuild under either grader with seeded intervals; the "
    "analysis package carries a 22-check self-test that runs after every change (22/22 passing "
    "as of this deck).",False,DARK)],size=13)

# =====================================================================
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
