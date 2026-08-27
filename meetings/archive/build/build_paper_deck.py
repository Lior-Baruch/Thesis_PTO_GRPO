"""Build the PAPER deck for the 2026-08-16 supervision meeting.

About the CLPsych draft *Affirmation Without Inquiry*
(papers/2026_clpsych_mi_reward_hacking/) rather than about run status: what the paper argues,
what evidence carries each move, what it explicitly does NOT claim, and the decisions needed to
submit it (co-authors, venue, timing).

It deliberately carries a 'what the paper does not claim' slide. Three claims were weakened
during drafting after checking the artifacts -- the retention intervals overlap at best-vs-best,
the PTO/GRPO gap is a state-distribution effect rather than 'exploration', and the held-out
grid for this paper is 16,896 cells not 22,272 -- and a deck that hides that invites a supervisor
to repeat the stronger version in public.

Numbers are owned by papers/2026_clpsych_mi_reward_hacking/NUMBERS.md, which maps each one to the
eda/results/ artifact it came from. If a number here disagrees with that ledger, the ledger wins.

Exports .pptx, then convert to .pdf via export_pdf.ps1 (PowerPoint COM).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))       # meetings/build/
REPO = os.path.dirname(os.path.dirname(HERE))           # repo root (meetings/ lives here)
ROOT = os.path.join(REPO, "Exp3_PTO_GRPO")             # the experiment the artifacts come from
L0F  = os.path.join(ROOT, "eda", "results", "L0", "figures")
L0T  = os.path.join(ROOT, "eda", "results", "L0", "tables")
L5F  = os.path.join(ROOT, "eda", "results", "L5", "figures")
METH = os.path.join(ROOT, "figures")   # hand-authored method schematics (build_method_figures.py)
OUT  = os.path.join(REPO, "meetings", "2026-08-16", "paper_2026-08-16.pptx")

# Since 2026-07-28 every grader has its own leaf: results/<view>/figures/<family>/<judge>/<name>.
# The decks show the PRIMARY grader's figures, so the judge segment is injected here rather than
# spelled out at ~40 call sites. Point JUDGE at another grader to rebuild the same deck off it.
JUDGE = "gpt-4o-mini"

# Families that are ABOUT the judges rather than produced BY one export with no <JUDGE> level
# (eda_analysis.exports.JUDGE_INVARIANT_GROUPS). Keep in sync, or a deck asks for a path that the
# EDA never writes -- which fails loudly at add_picture, but only when someone builds a deck.
JUDGE_INVARIANT_FAMILIES = {"8_measurement"}

def _jp(base, p):
    """``<base>/<family>[/<sub>]/<JUDGE>/<name>`` — judge goes ahead of the FILENAME, not the family.

    A judge-invariant family has no <JUDGE> segment at all: its artifacts contain every grader.
    """
    *parts, name = p.split("/")
    if parts and parts[0] in JUDGE_INVARIANT_FAMILIES:
        return os.path.join(base, *parts, name)
    return os.path.join(base, *parts, JUDGE, name)

def f0(p): return _jp(L0F, p)
def t0(p): return _jp(L0T, p)
def f5(p): return _jp(L5F, p)
def fm(p): return os.path.join(METH, p)   # schematics: no view, no judge, no family

NAVY  = RGBColor(0x1F,0x3A,0x5F); PTO = RGBColor(0x00,0x72,0xB2)
GRPO  = RGBColor(0xE6,0x9F,0x00); GREY= RGBColor(0x5A,0x5A,0x5A)
LIGHT = RGBColor(0xEF,0xF2,0xF6); GREEN=RGBColor(0x00,0x8A,0x63)
RED   = RGBColor(0xC0,0x4A,0x1A); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x22,0x28,0x30); ROWALT=RGBColor(0xF4,0xF7,0xFA)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_N = 0

def slide(fnum=True):
    global _N
    s = prs.slides.add_slide(BLANK)
    if fnum:
        _N += 1
        tb = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"Exp3 · {_N}"; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xB0,0xB8,0xC0); r.font.name = FONT
    return s

def box(s,l,tp,w,h):
    tb = s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tb.text_frame.word_wrap = True; return tb

def rect(s,l,tp,w,h,color,line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp

def set_runs(para, segs, size, align=PP_ALIGN.LEFT):
    para.alignment = align
    for txt,bold,color in segs:
        r = para.add_run(); r.text = txt; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

def title_bar(s, text, kicker=None):
    rect(s,0,0,13.333,1.15,NAVY); rect(s,0,1.15,13.333,0.06,PTO)
    tb = box(s,0.55,0.18,12.2,0.9); tf = tb.text_frame
    p = tf.paragraphs[0]
    if kicker:
        set_runs(p,[(kicker,True,RGBColor(0x8F,0xC7,0xEC))],13)
        set_runs(tf.add_paragraph(),[(text,True,WHITE)],26)
    else:
        set_runs(p,[(text,True,WHITE)],27)

def bullets(s, items, left=0.6, top=1.5, width=12.1, size=16, gap=8):
    tb = box(s,left,top,width,5.6); tf = tb.text_frame
    first = True
    for segs,level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level = level; p.space_before = Pt(gap); p.space_after = Pt(2)
        mk = "•  " if level==0 else "–  "
        set_runs(p,[("    "*level+mk, False, PTO if level==0 else GREY)]+segs, size)

def fit(img,mw,mh):
    with Image.open(img) as im: iw,ih = im.size
    ar = iw/ih; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    return w,h

def figure(s,img,mw,mh,left,top,anchor="center"):
    """Place img fitted inside the (mw × mh) band at (left, top). Returns the image's bottom edge."""
    w,h = fit(img,mw,mh)
    dy = 0.0 if anchor=="top" else (mh-h)/2
    s.shapes.add_picture(img,Inches(left+(mw-w)/2),Inches(top+dy),Inches(w),Inches(h))
    return top+dy+h

def caption(s,text,left,top,width,color=GREY,size=11.5,align=PP_ALIGN.CENTER):
    p = box(s,left,top,width,0.5).text_frame.paragraphs[0]
    set_runs(p,[(text,False,color)],size,align)

def side_notes(s, heading, rows, left, top, width, size=12, gap=10):
    """Heading + bulleted rows in a narrow right-hand column next to a figure."""
    p = box(s,left,top,width,0.45).text_frame.paragraphs[0]
    set_runs(p,[(heading,True,NAVY)],14.5)
    tb = box(s,left,top+0.45,width,5.2); tf = tb.text_frame
    first = True
    for segs in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = Pt(gap)
        set_runs(p,[("•  ",False,PTO)]+segs,size)

def cards(s, items, top, left=0.6, width=12.1, height=1.2, gap=0.12):
    """Option cards: (title, body, accent colour, tag)."""
    y = top
    for name,desc,col,tag in items:
        rect(s,left,y,width,height,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,left,y,0.14,height,col)
        tb = box(s,left+0.35,y+0.10,width-0.55,height); tf = tb.text_frame
        seg = [(name,True,col)]
        if tag: seg.append(("    ("+tag+")",True,GREEN))
        set_runs(tf.paragraphs[0],seg,14.5)
        set_runs(tf.add_paragraph(),[(desc,False,DARK)],12)
        y += height + gap
    return y

# ---- native pptx table from a markdown table file ----
def _num(x):
    try: float(x); return True
    except: return False

def _fmt(col,v):
    if not _num(v): return v
    fv = float(v)
    if col in ("p","p_holm","wilcoxon_p") and fv == 0: return "<.001"
    if col in ("iteration","target_iter"): return f"{fv:.0f}"
    s = f"{fv:.3f}".rstrip("0").rstrip(".")
    return s if s else "0"

def md_table(s, md_path, left, top, width, height, drop=(), keep=None,
             fontsize=9.0, rename=None, order=None):
    raw = [l for l in open(md_path,encoding="utf-8").read().splitlines() if l.strip().startswith("|")]
    rows = [[c.strip() for c in l.strip().strip("|").split("|")] for l in raw]
    header, body = rows[0], rows[2:]
    if keep: body = [r for r in body if keep(dict(zip(header,r)))]
    if order: body = sorted(body, key=lambda r: order(dict(zip(header,r))))
    idx = [i for i,h in enumerate(header) if h not in drop]
    disp = [(rename or {}).get(header[i],header[i]) for i in idx]
    data = [[_fmt(header[i], r[i]) for i in idx] for r in body]
    if not data:
        raise ValueError(
            f"md_table: no rows left after keep() on {os.path.basename(md_path)} — the filter "
            f"matched nothing. Header={header}. This usually means an EDA re-render changed a "
            f"column's number formatting (e.g. '10.000' -> '10'); fix the keep() predicate.")
    nr, nc = len(data)+1, len(idx)
    gt = s.shapes.add_table(nr,nc,Inches(left),Inches(top),Inches(width),Inches(height)).table
    lens = [max(len(disp[c]), *(len(data[r][c]) for r in range(len(data)))) for c in range(nc)]
    tot = sum(lens)
    for c in range(nc):
        gt.columns[c].width = Inches(width*lens[c]/tot)
    for c,htext in enumerate(disp):
        cell = gt.cell(0,c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_top=Pt(1);cell.margin_bottom=Pt(1);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        r = pr.add_run(); r.text = htext; r.font.size=Pt(fontsize); r.font.bold=True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for ri,row in enumerate(data,1):
        for c,val in enumerate(row):
            cell = gt.cell(ri,c); cell.fill.solid()
            cell.fill.fore_color.rgb = ROWALT if ri%2 else WHITE
            cell.margin_top=Pt(0);cell.margin_bottom=Pt(0);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
            pr = cell.text_frame.paragraphs[0]
            col = DARK
            if disp[c] in ("dz","delta","mean_delta") and val.startswith("-"): col = RED
            rr = pr.add_run(); rr.text = val; rr.font.size=Pt(fontsize)
            rr.font.color.rgb = col; rr.font.name = FONT
    return gt

def grid_table(s, rows, left, top, width, colw, fontsize=12.5, rowh=0.42):
    """rows = list of (cells:list[str], color_or_None, is_header:bool)"""
    y = top
    for ri,(cells, color, is_hdr) in enumerate(rows):
        x = left
        for i,c in enumerate(cells):
            w = width*colw[i]
            rect(s,x,y,w,rowh, NAVY if is_hdr else (LIGHT if ri%2==0 else WHITE))
            p = box(s,x+0.1,y+0.04,w-0.15,rowh).text_frame.paragraphs[0]
            col = WHITE if is_hdr else (color if (color and i==len(cells)-1) else DARK)
            set_runs(p,[(c,is_hdr or i==0,col)],fontsize)
            x += w
        y += rowh + 0.04
    return y

# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.5,13.333,0.08,PTO)
tb = box(s,0.9,1.35,11.5,2.8); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Affirmation Without Inquiry",True,WHITE)],40)
p = tf.add_paragraph(); p.space_before = Pt(6)
set_runs(p,[("Reward hacking when an LLM judge trains a Motivational Interviewing therapist",
             False,RGBColor(0xBD,0xD6,0xEA))],19)
p = tf.add_paragraph(); p.space_before = Pt(14)
set_runs(p,[("A complete 8-page draft, ready to read",True,RGBColor(0x8F,0xC7,0xEC))],15)

tb2 = box(s,0.9,4.75,11.5,2.2); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("What the paper argues, what carries it, and what it refuses to claim",
                             True,WHITE)],19)
p = tf2.add_paragraph(); p.space_before = Pt(6)
set_runs(p,[("Target: CLPsych / clinical-NLP workshop, ACL style. Scope: Exp3 at matched "
             "look-ahead K=0 — the optimizer is the only variable.",
             False,RGBColor(0x9F,0xB4,0xC8))],13.5)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_runs(p,[("Lior Baruch · Reichman University · 16 August 2026",
             False,RGBColor(0x9F,0xB4,0xC8))],13)

# =====================================================================
# 2 · THE ARGUMENT IN ONE LINE
# =====================================================================
s = slide(); title_bar(s,"The argument, in one sentence","THE PAPER")

rect(s,0.6,1.45,12.1,1.5,LIGHT); rect(s,0.6,1.45,0.07,1.5,PTO)
tb = box(s,0.95,1.58,11.5,1.4); tf = tb.text_frame
set_runs(tf.paragraphs[0],
  [("An LLM-judge reward takes a 1B therapist from below basic competence into the "
    "fair-to-good MITI band — but what it teaches is ",False,DARK),
   ("affirmation without inquiry",True,NAVY),
   (", and a held-out judge from another model family credits ",False,DARK),
   ("0.80",True,PTO),(" of one optimizer's Q1 gain and only ",False,DARK),
   ("0.28",True,GRPO),(" of the other's.",False,DARK)],17)

p = box(s,0.6,3.2,12.1,0.4).text_frame.paragraphs[0]
set_runs(p,[("One claim, three moves — not five parallel findings",True,NAVY)],16)

cards(s,[
 ("1 · The loop works, on its own terms",
  "Large effects on all five global rubrics; MITI relational ratings cross the manual's \"good\" "
  "threshold in both arms. Read only through the grader that supplied the reward, this is a "
  "strong positive result.", GREEN, "§4"),
 ("2 · What it actually learned: affirm, stop asking",
  "Turns 2.3–3.4× longer, affirmations 5–6× up, MI-inconsistency 2.3× (PTO) / 4.0× (GRPO), "
  "questions collapse. The clinical failure mode has a name, and MI says it is the wrong thing "
  "to learn.", GRPO, "§5"),
 ("3 · How much of the gain is real",
  "A grader from a different family that never played the patient and never produced a reward "
  "re-scores all 16,896 cells. Gain retention = held-out improvement ÷ trained-against "
  "improvement.", PTO, "§6"),
],3.65,height=1.12,gap=0.12)

rect(s,0.6,7.08,12.1,0.03,RGBColor(0xD5,0xDD,0xE5))
set_runs(box(s,0.6,6.95,12.1,0.45).text_frame.paragraphs[0],
  [("PTO vs GRPO is not a section — it is the contrast variable running through all three moves.",
    True,NAVY)],13.5)

# =====================================================================
# 3 · MOVE 1 — THE LOOP WORKS
# =====================================================================
s = slide(); title_bar(s,"Move 1 — on its own terms, the loop works well","EVIDENCE · §4")

b = figure(s,f0("1_outcomes/trajectories/trajectory_Q1Q2.png"),7.1,4.5,0.5,1.5)
caption(s,"Training reward (Q1+Q2) by iteration, mean over 96 personas, primary grader. "
          "Bands are bootstrap 95% CIs.",0.5,b+0.08,7.1)

side_notes(s,"What the numbers say",[
  [("PTO: Q1+Q2 3.00 → 4.26",True,PTO),(" (dz 1.43, large). Every global rubric a large effect, "
    "Holm p≈0.",False,DARK)],
  [("GRPO: peaks 4.08 at iteration 8",True,GRPO),(", then regresses to 3.75 by iteration 10.",
    False,DARK)],
  [("Climb rate",True,NAVY),(" — OLS 0.120/iter (PTO) vs 0.072 (GRPO).",False,DARK)],
  [("Absolute anchor",True,NAVY),(": MITI relational crosses the manual's \"good\" threshold in "
    "both arms (PTO 4.61, GRPO 4.20) — but ",False,DARK),
   ("neither reaches \"good\" on the technique ratios",True,RED),(".",False,DARK)],
],7.9,1.55,4.9,size=12.5,gap=13)

rect(s,7.9,6.05,4.85,1.1,LIGHT); rect(s,7.9,6.05,0.06,1.1,GREEN)
set_runs(box(s,8.1,6.15,4.6,1.0).text_frame.paragraphs[0],
  [("Both arms also fix something real: ",True,NAVY),
   ("the base model degenerates into phrase loops in ~49% of conversations; both drive that to "
    "zero. Part of the gain is genuine repair.",False,DARK)],12)

# =====================================================================
# 4 · FINAL vs BEST — THE FAIRNESS POINT
# =====================================================================
s = slide(); title_bar(s,"Comparing only final iterations is unfair to GRPO","EVIDENCE · §4")

set_runs(box(s,0.6,1.35,12.1,0.5).text_frame.paragraphs[0],
  [("GRPO regresses after iteration 8, so the head-to-head is reported ",False,DARK),
   ("both ways",True,NAVY),(". The weaker reading is the one the paper defends.",False,DARK)],15)

grid_table(s,[
  (["Paired PTO − GRPO","final  (10 v 10)","best  (10 v 8)","verdict"],None,True),
  (["Q1+Q2  (the reward)","+0.51    dz 0.73","+0.18    dz 0.30","PTO leads either way"],
   GREEN,False),
  (["MITI","+0.35    dz 0.46","+0.04    n.s. (p .57)","gap vanishes"],RED,False),
  (["MICI ↓  (sycophancy)","−0.35    dz −0.99","−0.04    n.s. (p .52)","gap vanishes"],RED,False),
  (["Q1 / MI-SAT / CSQ-8 / WAI-SR","all favour PTO","all survive Holm","holds"],GREEN,False),
],0.6,1.95,12.1,[0.28,0.24,0.24,0.24],fontsize=13,rowh=0.5)

rect(s,0.6,4.85,12.1,1.05,LIGHT); rect(s,0.6,4.85,0.07,1.05,RED)
set_runs(box(s,0.9,4.95,11.6,1.0).text_frame.paragraphs[0],
  [("So the sycophancy separation is a property of GRPO's ",False,DARK),
   ("post-peak run-off",True,NAVY),(", not of its best state. At GRPO's own peak the MI-consistency "
    "gap is not established. That is materially weaker than the endpoint numbers alone suggest — "
    "and it is what the paper claims.",False,DARK)],14)

rect(s,0.6,6.15,12.1,1.0,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,0.6,6.15,0.07,1.0,PTO)
set_runs(box(s,0.9,6.25,11.6,0.95).text_frame.paragraphs[0],
  [("And \"best\" is itself grader-dependent. ",True,NAVY),
   ("The primary oracle puts GRPO's peak at iteration 8; the held-out judge puts it at iteration 3. "
    "Peak selection performed on the very signal being optimised is not a neutral operation — which "
    "is part of the paper's argument rather than an aside.",False,DARK)],13.5)

# =====================================================================
# 5 · MOVE 2 — WHAT THE THERAPIST LEARNED
# =====================================================================
s = slide(); title_bar(s,"Move 2 — the therapist becomes warmer and stops asking",
                       "EVIDENCE · §5")

grid_table(s,[
  (["Behaviour (per therapist turn)","PTO base → 10","GRPO base → 10"],None,True),
  (["Turn length (characters)","301 → 686","266 → 896"],None,False),
  (["Affirmations  (MITI B6)","0.025 → 0.142","0.029 → 0.154"],None,False),
  (["Over-praise  (MICI) ↓","0.013 → 0.299","0.019 → 0.698"],RED,False),
  (["MI-inconsistency  (MICI) ↓","0.21 → 0.49","0.21 → 0.84"],RED,False),
  (["Questions, regex \"?\"","0.93 → 0.55","0.83 → 0.15"],RED,False),
  (["Confrontation ↓","0.007 → 0.000","0.008 → 0.000"],GREEN,False),
],0.6,1.5,7.5,[0.44,0.28,0.28],fontsize=12.5,rowh=0.46)

caption(s,"96 conversations per cell, primary grader.",0.6,5.15,7.5,align=PP_ALIGN.LEFT)

side_notes(s,"The clinical reading",[
  [("MI's core move is the open question",True,NAVY),(" — eliciting the client's own reasons for "
    "change. GRPO nearly stops asking.",False,DARK)],
  [("The rise in MI-inconsistency is almost entirely over-praise",True,NAVY),
   (". Confrontation and judging go to zero; unsolicited advice is flat. The model is not becoming "
    "adversarial — it is becoming ingratiating.",False,DARK)],
  [("Which is exactly what a patient-satisfaction reward would be expected to select for.",
    True,RED)],
],8.4,1.5,4.35,size=12.5,gap=13)

rect(s,0.6,5.6,12.1,1.5,LIGHT); rect(s,0.6,5.6,0.07,1.5,GRPO)
set_runs(box(s,0.9,5.72,11.6,1.4).text_frame.paragraphs[0],
  [("The reward's own composition pulls this way. ",True,NAVY),
   ("The three Q2 items gaining most in both arms are \"revealed what he was thinking\" (+1.07), "
    "\"put himself in my shoes\" (+1.01) and \"took charge\" (+0.99) — two rewarding therapist "
    "self-disclosure and one therapist direction. Neither is prescribed by MI; the latter is "
    "actively discouraged. A general-psychotherapy alliance scale, repurposed as an MI reward, "
    "contains items that pay for non-MI behaviour. No choice of optimizer fixes that.",False,DARK)],
  13)

# =====================================================================
# 6 · THE DIAGNOSTIC — REGEX vs THE ORACLE'S OWN COUNT
# =====================================================================
s = slide(); title_bar(s,"The grader keeps crediting inquiry after it has stopped",
                       "EVIDENCE · §5 — the paper's novel diagnostic")

b = figure(s,f0("3_validity/question_rate_crosscheck.png"),7.7,4.35,0.45,1.45)
caption(s,"Solid: deterministic count of \"?\" per therapist turn. Dashed: the oracle's own MITI "
          "B3 question code, same conversations.",0.45,b+0.06,7.7)

side_notes(s,"Why this matters",[
  [("Two ways of counting the same thing",True,NAVY),(". One is an LLM judgement; the other is a "
    "regular expression and costs nothing.",False,DARK)],
  [("They agree at initialisation",True,DARK),(" — the untrained model emits question marks the "
    "coder declines to treat as MI questions.",False,DARK)],
  [("In GRPO they cross between iterations 4 and 5",True,GRPO),(" and end inverted: 0.15 regex "
    "against 0.32 coded. PTO never inverts.",False,DARK)],
  [("So part of the apparent competence at the GRPO endpoint is a property of the ",False,DARK),
   ("measuring instrument",True,RED),(", not of the transcript.",False,DARK)],
],8.35,1.5,4.4,size=12.5,gap=12)

rect(s,8.35,5.85,4.4,1.25,LIGHT); rect(s,8.35,5.85,0.06,1.25,PTO)
set_runs(box(s,8.55,5.95,4.15,1.2).text_frame.paragraphs[0],
  [("The crossover point is not arbitrary: ",True,NAVY),
   ("iterations 4–5 is also where the held-out grader starts withdrawing credit from GRPO. Two "
    "views of one event.",False,DARK)],12)

# =====================================================================
# 7 · MOVE 3 — THE HELD-OUT JUDGE
# =====================================================================
s = slide(); title_bar(s,"Move 3 — a grader that took no part in training","EVIDENCE · §6")

rect(s,0.6,1.4,12.1,1.15,LIGHT); rect(s,0.6,1.4,0.07,1.15,PTO)
set_runs(box(s,0.9,1.5,11.6,1.1).text_frame.paragraphs[0],
  [("Claude Haiku 4.5 re-scored the complete grid — ",False,DARK),
   ("16,896 cells",True,NAVY),(", 8 instruments × 22 model states × 96 conversations, no partial "
    "cells. Decoupled three ways: a different model family, it never played the patient, and it "
    "never produced a reward. Cost ",False,DARK),("$42",True,NAVY),
   (" batched — about 13% of the project's API spend, for a complete second reading of every "
    "result.",False,DARK)],14)

p = box(s,0.6,2.75,6.0,0.4).text_frame.paragraphs[0]
set_runs(p,[("The rankings survive",True,GREEN)],16)
bullets(s,[
  ([("88.3% of all 1,848 arm × instrument contrasts keep their sign",True,DARK),
    (" — rising to 94.1% at |Δ|≥0.10 and 98.9% at |Δ|≥0.50. The judges disagree only about "
     "differences too small to claim.",False,DARK)],0),
  ([("Only 1.2–6.9% of arm-mean variance",True,DARK),(" sits in the arm × judge interaction — the "
    "one component that could invalidate a ranking.",False,DARK)],0),
  ([("Haiku is systematically harsher",True,DARK),(" (1.2–1.7 points), which cancels in every "
    "contrast.",False,DARK)],0),
],left=0.6,top=3.2,width=6.0,size=13,gap=11)

rect(s,6.85,2.75,0.02,3.6,RGBColor(0xD5,0xDD,0xE5))
p = box(s,7.15,2.75,5.55,0.4).text_frame.paragraphs[0]
set_runs(p,[("One exception, stated plainly",True,RED)],16)
bullets(s,[
  ([("MITI does not meet that bar.",True,RED),(" Only 3.6% of its arm-mean variance is between-arm "
    "signal (94.5% is grader level), dependability 0.65, and it preserves sign on just 77.5% of "
    "contrasts — the worst of the eight.",False,DARK)],0),
  ([("So every MITI-based statement is flagged provisional",True,DARK),(", including the "
    "competency-threshold placement on the previous slide.",False,DARK)],0),
  ([("Q1, Q2, PCT and MICI are unaffected.",True,DARK)],0),
],left=7.15,top=3.2,width=5.55,size=13,gap=11)

rect(s,0.6,6.5,12.1,0.75,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,0.6,6.5,0.07,0.75,RED)
set_runs(box(s,0.9,6.58,11.6,0.7).text_frame.paragraphs[0],
  [("Never average the two graders. ",True,RED),
   ("The primary oracle WAS the training reward — this is train-versus-test, not two raters. "
    "Only contrasts are compared.",False,DARK)],13)

# =====================================================================
# 8 · GAIN RETENTION
# =====================================================================
s = slide(); title_bar(s,"Gain retention — how much of the improvement a neutral grader credits",
                       "EVIDENCE · §6 — the paper's contribution")

rect(s,0.6,1.38,12.1,0.72,LIGHT)
set_runs(box(s,0.9,1.46,11.6,0.7).text_frame.paragraphs[0],
  [("retention  =  Δ(held-out judge)  ÷  Δ(trained-against oracle)",True,NAVY),
   ("        both measured against the same base policy, so the graders' level offset cancels. "
    "1.0 = the full gain is credited; 0 = none of it is.",False,DARK)],14)

grid_table(s,[
  (["Q1  (session satisfaction)","Δ trained","Δ held-out","retention  [95% CI]"],None,True),
  (["PTO @ 10","1.22","0.97","0.80   [0.68, 0.93]"],GREEN,False),
  (["GRPO @ 10","0.68","0.19","0.28   [0.06, 0.43]"],RED,False),
  (["PTO @ 8","1.18","1.05","0.89   [0.76, 1.04]"],None,False),
  (["GRPO @ 8  (its peak)","1.01","0.65","0.64   [0.52, 0.78]"],None,False),
],0.6,2.3,7.55,[0.34,0.2,0.2,0.26],fontsize=12.5,rowh=0.46)

caption(s,"Q2 is the control: all four intervals overlap in a 0.80–0.86 band.",
        0.6,4.85,7.55,align=PP_ALIGN.LEFT)

side_notes(s,"Read it in the instrument's units",[
  [("Under a grader that never graded during training, GRPO's net ten-iteration Q1 gain is ",
    False,DARK),("≈0.19 points",True,RED),(", not the ≈0.68 its own reward reports.",False,DARK)],
  [("Q1 and Q2 behave differently",True,NAVY),(", which is what shows retention is not simply "
    "\"the second grader is stingier\" — if it were, they would move together.",False,DARK)],
  [("Q1 asks ",False,DARK),("was this helpful, did I learn, will I act on it",True,NAVY),
   (" — precisely where warmth substitutes for substance.",False,DARK)],
],8.35,2.25,4.4,size=12.5,gap=13)

rect(s,0.6,5.35,12.1,1.75,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,0.6,5.35,0.07,1.75,RED)
tb = box(s,0.9,5.45,11.6,1.7); tf = tb.text_frame
set_runs(tf.paragraphs[0],
  [("Where this claim stops — and the paper says so. ",True,RED),
   ("The intervals are non-overlapping at iterations 9–10 ONLY. At GRPO's own iteration-8 peak "
    "PTO still leads (0.89 vs 0.64) but the intervals overlap, and so do PTO@10 against GRPO@8.",
    False,DARK)],13.5)
p = tf.add_paragraph(); p.space_before = Pt(7)
set_runs(p,[("What is robust is the ordering, not any single disjoint pair: ",True,NAVY),
            ("PTO's Q1 retention exceeds GRPO's at every iteration from 4 onward, with "
             "monotonically diverging trends.",False,DARK)],13.5)

# =====================================================================
# 9 · THE ONSET CURVE
# =====================================================================
s = slide(); title_bar(s,"It is an onset curve — visible four iterations before the reward curve "
                        "turns","EVIDENCE · §6")

b = figure(s,f0("8_measurement/multijudge_retention_trajectory.png"),8.0,4.6,0.4,1.4)
caption(s,"Retention by iteration, every instrument where the ratio is estimable. Blue PTO, "
          "orange GRPO. 1.0 = full gain credited.",0.4,b+0.05,8.0)

side_notes(s,"The shape is the finding",[
  [("Indistinguishable for three iterations",True,DARK),(" (0.84–0.97 in both arms).",False,DARK)],
  [("Then they separate.",True,NAVY),(" PTO holds 0.80–0.98 for the rest of the run; GRPO decays "
    "in trend to 0.28.",False,DARK)],
  [("Retention separates the arms from iteration 4",True,GREEN),(", while the reward curve does "
    "not distinguish them until GRPO turns over at iteration 8.",False,DARK)],
  [("A monitor computed on data you already have would have flagged this four iterations earlier.",
    True,PTO)],
  [("Honest caveat",True,RED),(": the separation is confined to Q1 and MITI. WAI-SR, CSQ-8 and "
    "MI-SAT track each other in both arms, and sit above 1.0 — Haiku credits more than the "
    "primary there.",False,DARK)],
],8.65,1.45,4.15,size=12,gap=10)

# =====================================================================
# 10 · WHAT THE PAPER DOES NOT CLAIM
# =====================================================================
s = slide(); title_bar(s,"What the paper deliberately does not claim","HONESTY")

set_runs(box(s,0.6,1.35,12.1,0.45).text_frame.paragraphs[0],
  [("Two of these were weaker versions of claims I had written down before checking the "
    "artifacts. Both were corrected in the draft.",False,DARK)],14)

cards(s,[
 ("Not \"the retention result is significant at best-vs-best\"",
  "Non-overlapping intervals hold at iterations 9–10 only. At peak-vs-peak they overlap. The "
  "defensible statement is the consistent ordering from iteration 4 onward. — corrected during drafting",
  RED,""),
 ("Not \"DPO beats GRPO\"",
  "Candidate sampling is matched by construction (same temperature 1.2, M = G = 8, same MCL). "
  "Swapping the weighting rule on the same groups barely moves the update direction (0.908 / 0.988). "
  "The gap is the STATE distribution, not the loss. — corrected during drafting",
  RED,""),
 ("Not \"these models do good MI\"",
  "Every number is produced by an LLM — patient, reward and both graders. There is no human "
  "MI-coder validation. The defended claim is comparative: this arm's gains transfer to an "
  "unrelated grader and that arm's do not.",
  GRPO,"the biggest gap"),
 ("Not a length-controlled result",
  "Turns grow 2.3–3.4× while sessions shorten. We cannot partition the gain into content and "
  "verbosity. It bears on every effect size — though not on the retention contrast, where both "
  "arms lengthen and only one loses retention.",
  GREY,""),
 ("Not a K × method result",
  "The look-ahead comparison is excluded entirely. GRPO at K=5 has one scored iteration, so it "
  "cannot be stated as an interaction — a separate piece once that arm is scored.",
  GREY,"scope"),
],1.95,height=1.0,gap=0.1)

# =====================================================================
# 11 · THE MECHANISM — STATE DISTRIBUTION, NOT THE LOSS
# =====================================================================
s = slide(); title_bar(s,"Why PTO holds up: it trains on different states, not a different loss",
                       "MECHANISM · Appendix B")

p = box(s,0.6,1.4,12.1,0.4).text_frame.paragraphs[0]
set_runs(p,[("Re-weighting each method's own candidate groups under the ",False,DARK),
            ("other",True,NAVY),(" method's rule separates the two explanations.",False,DARK)],14.5)

grid_table(s,[
  (["Comparison","cosine","corrected","what it means"],None,True),
  (["As trained  (rule AND data differ)","0.267","0.317","far apart"],None,False),
  (["Same data, weighting rule swapped","0.908 / 0.988","—","rule barely matters"],GREEN,False),
  (["Same rule, data differ","0.356 / 0.266","0.397 / 0.324","still far apart"],RED,False),
],0.6,1.9,7.6,[0.4,0.2,0.2,0.2],fontsize=12,rowh=0.46)

side_notes(s,"So what actually differs?",[
  [("GRPO",True,GRPO),(" slices an ",False,DARK),("unmodified on-policy rollout",True,NAVY),
   (" — every prefix is a state the current policy really produces.",False,DARK)],
  [("PTO",True,PTO),(" grows a trunk by appending the ",False,DARK),
   ("oracle-argmax of 8",True,NAVY),(" at each therapist turn, and that selection compounds.",
    False,DARK)],
  [("PTO therefore trains on states from a best-of-M reranked policy — closer to expert iteration "
    "than to exploration.",True,NAVY)],
],8.4,1.85,4.35,size=12,gap=11)

rect(s,0.6,4.05,7.6,1.15,LIGHT); rect(s,0.6,4.05,0.06,1.15,PTO)
set_runs(box(s,0.85,4.15,7.2,1.1).text_frame.paragraphs[0],
  [("Which suggests a mechanism the outcome scores cannot show: ",True,NAVY),
   ("GRPO's training prompts drift along with its own policy. As it becomes more effusive, the "
    "states it is next trained on are themselves more effusive.",False,DARK)],12.5)

rect(s,0.6,5.4,12.1,1.65,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,0.6,5.4,0.07,1.65,GRPO)
tb = box(s,0.9,5.5,11.6,1.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("And the hack is a compounding loop, not a hard pull. ",True,NAVY),
  ("Per-iteration selection pressure toward affirmation stays small (≈0.01 → 0.10) while what the "
   "policy generates moves an order of magnitude further (0.02 → 0.54 GRPO, 0.04 → 0.57 PTO).",
   False,DARK)],13)
p = tf.add_paragraph(); p.space_before = Pt(6)
set_runs(p,[("A weak, persistently same-signed preference applied each iteration to an "
             "already-more-effusive policy is enough. By the end the update is choosing between "
             "two effusive candidates — which is why single-iteration diagnostics miss it, and why "
             "retention (which compares policies, not candidates) sees it first.",False,DARK)],13)

# =====================================================================
# 12 · WHERE THE DRAFT IS
# =====================================================================
s = slide(); title_bar(s,"Where the draft is","STATUS")

p = box(s,0.6,1.4,6.0,0.4).text_frame.paragraphs[0]
set_runs(p,[("Complete and building",True,GREEN)],16.5)
bullets(s,[
  ([("Body fits inside 8 pages",True,DARK),(" — the ACL workshop limit — with 0 overfull boxes and "
    "0 undefined references. 13 pages including Limitations, Ethics, References and three "
    "appendices, none of which count.",False,DARK)],0),
  ([("Eleven sections plus three appendices",True,DARK),(": measurement quality, the training-signal "
    "probe, and reproducibility.",False,DARK)],0),
  ([("Five tables, three figures",True,DARK),(", all copied unmodified from the tracked EDA tree by "
    "a sync script that fails loudly if a source figure moves.",False,DARK)],0),
  ([("A claims ledger (NUMBERS.md)",True,DARK),(" maps every number in the draft to the artifact "
    "path it came from, and records the traps — so a re-render identifies every sentence that has "
    "to change.",False,DARK)],0),
],left=0.6,top=1.9,width=6.0,size=13,gap=11)

rect(s,6.85,1.4,0.02,4.6,RGBColor(0xD5,0xDD,0xE5))
p = box(s,7.15,1.4,5.55,0.4).text_frame.paragraphs[0]
set_runs(p,[("Open before submission",True,RED)],16.5)
bullets(s,[
  ([("Co-author list and order",True,NAVY),(" — currently a placeholder. This is the first thing "
    "I need from you.",False,DARK)],0),
  ([("Venue and date",True,NAVY),(" — CLPsych is the natural home; the draft is written to its "
    "8-page ACL format.",False,DARK)],0),
  ([("Two bibliography entries",True,DARK),(" inherited from an earlier draft still need checking "
    "against source.",False,DARK)],0),
  ([("Artifact release",True,DARK),(" — transcripts and scores are synthetic and releasable; "
    "adapter release needs your sign-off.",False,DARK)],0),
  ([("One figure is cramped",True,DARK),(" at column width and would read better as a Q1-only "
    "variant.",False,DARK)],0),
],left=7.15,top=1.9,width=5.55,size=13,gap=10)

rect(s,0.6,6.3,12.1,0.85,LIGHT); rect(s,0.6,6.3,0.06,0.85,PTO)
set_runs(box(s,0.85,6.4,11.7,0.8).text_frame.paragraphs[0],
  [("The body is at the page limit with no slack",True,NAVY),
   (" — any addition now needs a matching cut, and the README records where the existing slack "
    "already went.",False,DARK)],13.5)

# =====================================================================
# 13 · WHY THIS FRAMING
# =====================================================================
s = slide(); title_bar(s,"Why this framing, and what it leaves for later","SCOPE")

p = box(s,0.6,1.4,12.1,0.42).text_frame.paragraphs[0]
set_runs(p,[("Exp3 supports more than one paper. This draft takes the clinical question as the "
             "spine and uses the held-out judge as its proof.",False,DARK)],14.5)

cards(s,[
 ("In this paper — the clinical failure mode, proved by a neutral grader",
  "What an LLM-judge reward actually teaches a therapist model, and how much of the reported gain "
  "is real. Domain-first, so it belongs at CLPsych rather than a generic eval venue.",
  GREEN,"drafted"),
 ("Held back — the look-ahead (K) comparison",
  "K=5 never leads across 8 matched iterations under either grader, but GRPO at K=5 has only one "
  "scored iteration. Within-PTO today; a K × method result once that arm is scored.",
  GRPO,"blocked on budget"),
 ("Held back — \"exploration, not the loss\" as a standalone",
  "The update-direction probe is the most original result and the least battle-tested. It sits in "
  "an appendix here, answering the obvious reviewer question without carrying a paper.",
  PTO,"appendix for now"),
 ("The thesis chapter",
  "Wider than any single paper: all three research questions, all three experiment generations, no "
  "page limit. This draft becomes one chapter of it.",
  GREY,"later"),
],1.95,height=1.15,gap=0.13)

rect(s,0.6,7.0,12.1,0.03,RGBColor(0xD5,0xDD,0xE5))
set_runs(box(s,0.6,6.88,12.1,0.45).text_frame.paragraphs[0],
  [("Two earlier drafts were written and discarded — they carried five parallel claims each and "
    "ran over the limit. This one carries one claim with three moves.",False,GREY)],12.5)

# =====================================================================
# 14 · ASKS
# =====================================================================
s = slide(); title_bar(s,"What I'd like to leave with","ASKS")
bullets(s,[
  ([("Read the draft — it is complete, not an outline.",True,NAVY),(" Eight pages; the three moves "
    "are §4, §5 and §6, and §6 is the one I would most like challenged.",False,DARK)],0),
  ([("Co-author list and order",True,NAVY),(", so the title block can be filled in.",False,DARK)],0),
  ([("Confirm CLPsych as the target",True,NAVY),(" and the submission date I should be working "
    "back from.",False,DARK)],0),
  ([("A view on human MI-coder validation.",True,NAVY),(" It is the one validity gap that money "
    "cannot close, it is named as the first limitation, and a reviewer will ask. Costs time rather "
    "than budget — is it worth arranging, and who could code?",False,DARK)],0),
  ([("A yes or no on scoring GRPO at K=5",True,NAVY),(" beyond iteration 1 — not for this paper, "
    "which excludes K entirely, but it decides whether the look-ahead result is ever more than a "
    "within-PTO finding.",False,DARK)],0),
], top=1.5, size=15, gap=13)

rect(s,0.6,5.05,12.1,0.03,RGBColor(0xD5,0xDD,0xE5))
p = box(s,0.6,5.2,12.1,0.4).text_frame.paragraphs[0]
set_runs(p,[("Already on disk, if a question turns out to be cheap to answer",True,NAVY)],14.5)
bullets(s,[
  ([("2,112 scored conversations",True,DARK),(" across 22 model states, every one graded on all "
    "eight instruments by ",False,DARK),("both",True,DARK),(" judges — 16,896 cells per grader.",
    False,DARK)],0),
  ([("Per-item decomposition of every instrument",True,DARK),(", persona-level heterogeneity "
    "splits, and per-candidate training records — all computed, no further API calls.",False,DARK)],0),
  ([("Every figure, table and summary regenerates from one command",True,DARK),(", under either "
    "grader, with seeded confidence intervals.",False,DARK)],0),
], top=5.6, size=13, gap=9)

prs.save(OUT)
print("wrote", OUT, f"({_N+1} slides)")
