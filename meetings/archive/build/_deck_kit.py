"""_deck_kit.py — the house visual system, factored out of build_rq_deck_2026-08-21.py.

Every deck in ``meetings/build/`` before 2026-08-23 carried its own copy of these primitives, so a
new builder began with ~250 lines of transcription. This module is that block, verbatim in
behaviour, so a builder can be about SLIDES.

⚠ **This is a shared library, not a deck.** Past decks are records of what was presented and are
never edited; this file is imported only by builders written on or after 2026-08-23. If a future
deck needs a different look, add a primitive here — do not change an existing one's behaviour, or
a rebuild of an older deck that imports it would silently render differently from what was shown.

Geometry is 16:9 at 13.333 x 7.5 in. All public helpers take EMU/Inches, never points-as-floats.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── palette ───────────────────────────────────────────────────────────────────
INK = RGBColor(0x12, 0x20, 0x3A)    # deep navy — headings, dividers, table heads
BODY = RGBColor(0x2B, 0x33, 0x45)   # body text
MUTED = RGBColor(0x69, 0x72, 0x82)  # captions, provenance, kickers
RULE = RGBColor(0xD8, 0xDE, 0xE7)   # hairlines, bullet dashes
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF4, 0xF6, 0xF9)   # neutral band fill, alternating table rows
VERDICT = RGBColor(0x1E, 0x5B, 0x45)  # green — the answer band
VWASH = RGBColor(0xE8, 0xF1, 0xED)
CAVEAT = RGBColor(0x8A, 0x5A, 0x12)   # amber — the "not this" band
CWASH = RGBColor(0xFB, 0xF3, 0xE3)
PTO_C = RGBColor(0x2F, 0x6F, 0x8F)    # teal — PTO
GRPO_C = RGBColor(0xC2, 0x70, 0x3D)   # burnt orange — GRPO

# on-navy text tints (dividers, title, closing slide)
SKY = RGBColor(0x8F, 0xB6, 0xCB)
MIST = RGBColor(0xB9, 0xC4, 0xD4)
SLATE = RGBColor(0x9F, 0xAD, 0xC0)
DUSK = RGBColor(0x7E, 0x8C, 0xA2)
PANEL = RGBColor(0x1B, 0x2C, 0x4B)   # a raised panel on a navy slide
GOLD = RGBColor(0xD8, 0xA6, 0x5C)
MINT = RGBColor(0x8F, 0xC4, 0xAA)

FONT = "Segoe UI"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
BOTTOM = Inches(0.66)   # clearance kept under a bottom-anchored band for the source line
ML, MR = Inches(0.85), Inches(0.85)
CW = W - ML - MR


# ── primitives ────────────────────────────────────────────────────────────────
def txbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def run(p, text, *, size, bold=False, color=BODY, font=FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def para(tf, *, space_before=0, space_after=0, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.alignment = align
    return p


def rect(slide, x, y, w, h, fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def pic(slide, path, x, y, max_w, max_h):
    """Insert a PNG scaled to fit (max_w, max_h), centred in that box.

    Fails loudly if the artifact moved — a deck is never written half-updated.
    """
    if not os.path.exists(path):
        raise SystemExit(
            "MISSING ARTIFACT: %s\n"
            "The EDA tree moved or the family was not rendered. Re-run\n"
            "  python tools/render_results.py\n"
            "from Exp3_PTO_GRPO/eda/ and rebuild." % path
        )
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(path, int(x + (max_w - w) / 2), int(y + (max_h - h) / 2),
                                    width=w, height=h)


# ── slide furniture ───────────────────────────────────────────────────────────
def running_head(slide, kicker):
    tf = txbox(slide, ML, Inches(0.34), CW, Inches(0.26))
    p = para(tf, first=True)
    run(p, kicker.upper(), size=10.5, bold=True, color=MUTED)
    rect(slide, ML, Inches(0.66), CW, Emu(9525), RULE)


def heading(slide, text, y=Inches(0.86), size=27):
    tf = txbox(slide, ML, y, CW, Inches(0.62))
    p = para(tf, first=True)
    run(p, text, size=size, bold=True, color=INK)
    return y + Inches(0.66)


def band_h(text, size):
    """Height a band needs for `text` at `size`, from an empirically calibrated wrap width.

    Measured off the rendered deck: a 97-char run at 15.5 pt sits on one line inside CW, a
    152-char run wraps to two. 1600/size reproduces both, and the 12.5 pt caveat bands.
    """
    cpl = max(40, int(1600 / size))
    lines = max(1, -(-len(text) // cpl))
    return max(Inches(0.80), Inches(0.32) + Inches(0.26) * lines)


def band(slide, y, label, text, *, fill, edge, label_color, height=None,
         size=15.5, bold_text=True):
    """A labelled band — the deck's repeating device (VERDICT / NOT THIS / WHAT HAPPENED).

    `height=None` sizes the band to its text, so a long caveat can never spill over the
    provenance line.
    """
    height = band_h(text, size) if height is None else height
    rect(slide, ML, y, CW, height, fill)
    rect(slide, ML, y, Inches(0.055), height, edge)
    tf = txbox(slide, ML + Inches(0.28), y + Inches(0.11), CW - Inches(0.5),
               height - Inches(0.22))
    p = para(tf, first=True)
    run(p, label, size=9.5, bold=True, color=label_color)
    p2 = para(tf, space_before=3)
    run(p2, text, size=size, bold=bold_text, color=INK)
    return y + height + Inches(0.2)


def bandbot(slide, label, text, **kw):
    """Bottom-anchored band: its BOTTOM edge sits just above the source line, whatever its height."""
    h = band_h(text, kw.get("size", 15.5))
    return band(slide, H - BOTTOM - h, label, text, height=h, **kw)


def figband(slide, y, png, label, text, **kw):
    """Figure filling the space above a bottom-anchored caption band (band sized first)."""
    h = band_h(text, kw.get("size", 12.5))
    band_y = H - BOTTOM - h
    pic(slide, png, ML, y + Inches(0.02), CW, band_y - y - Inches(0.14))
    band(slide, band_y, label, text, height=h, **kw)


def provenance(slide, paths, y=None):
    y = y if y is not None else H - Inches(0.52)
    tf = txbox(slide, ML, y, CW, Inches(0.3))
    p = para(tf, first=True)
    run(p, "source:  ", size=8.5, color=MUTED)
    run(p, "   ·   ".join(paths), size=8.5, color=MUTED, font=MONO)


def footer(slide, n):
    tf = txbox(slide, W - MR - Inches(1.2), H - Inches(0.5), Inches(1.2), Inches(0.3))
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, str(n), size=9.5, color=MUTED)


def bullets(slide, y, items, *, size=13.5, gap=9, width=None, x=None, bottom=Inches(0.7)):
    width = width or CW
    x = ML if x is None else x
    tf = txbox(slide, x, y, width, H - y - bottom)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
        else:
            lead, rest = None, it
        p = para(tf, first=(i == 0), space_after=gap)
        run(p, "—   ", size=size, color=RULE)
        if lead:
            run(p, lead, size=size, bold=True, color=INK)
            run(p, "  ", size=size)
        run(p, rest, size=size, color=BODY)
    return tf


def table(slide, y, headers, rows, *, col_w, size=11.5, head_size=10, row_h=Inches(0.34),
          emphasis=None, prose_cols=(), left=None):
    """Native pptx table with the deck's flat styling (no banding, hairline rules).

    Returns the GraphicFrame. `left=None` centres it in the content column.
    """
    total = sum(col_w)
    x = ML + (CW - total) / 2 if left is None else left
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), int(x), int(y),
                                   int(total), int(row_h * (len(rows) + 1)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = int(cw)
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = ""
        c.fill.solid()
        c.fill.fore_color.rgb = INK
        c.margin_left = c.margin_right = Inches(0.09)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
        run(p, htxt, size=head_size, bold=True, color=PAPER)
    for i, row in enumerate(rows, start=1):
        tbl.rows[i].height = int(row_h)
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = PAPER if i % 2 else WASH
            c.margin_left = c.margin_right = Inches(0.09)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0 or j in prose_cols) else PP_ALIGN.CENTER
            colr = BODY
            bold = False
            if emphasis and emphasis(i - 1, j, val):
                colr, bold = INK, True
            run(p, str(val), size=size, bold=bold, color=colr,
                font=FONT if (j == 0 or j in prose_cols) else MONO)
    return shape


def table_bottom(y, n_rows, row_h=Inches(0.34), pad=Inches(0.18)):
    """Y coordinate just under a `table` of `n_rows` data rows drawn at `y`."""
    return y + row_h * (n_rows + 1) + pad


def factstrip(slide, fields, *, y=None, height=Inches(1.02), size=9.5, label_size=8):
    """Bottom strip of 2-5 labelled DESCRIPTIVE fields — the figure-catalogue device.

    `fields` is a list of (LABEL, text). Columns are equal width with hairline separators.
    Intended for stating what a chart plots — axes, series, n, how it was computed — WITHOUT
    saying what it means. Returns the strip's top y.
    """
    y = H - Inches(0.60) - height if y is None else y
    rect(slide, ML, y, CW, height, WASH)
    n = max(1, len(fields))
    colw = CW / n
    for i, (label, text) in enumerate(fields):
        x = ML + colw * i
        if i:
            rect(slide, x, y + Inches(0.10), Emu(9525), height - Inches(0.20), RULE)
        tf = txbox(slide, x + Inches(0.20), y + Inches(0.11), colw - Inches(0.34),
                   height - Inches(0.20))
        p = para(tf, first=True)
        run(p, label.upper(), size=label_size, bold=True, color=MUTED)
        p2 = para(tf, space_before=3)
        run(p2, text, size=size, color=BODY)
    return y


def sidepanel(slide, fields, x, y, w, h, *, size=9.5, label_size=8):
    """The fact strip's vertical twin — labelled descriptive fields stacked in a right column."""
    rect(slide, x, y, w, h, WASH)
    tf = txbox(slide, x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), h - Inches(0.34))
    for i, (label, text) in enumerate(fields):
        p = para(tf, first=(i == 0), space_before=0 if i == 0 else 11)
        run(p, label.upper(), size=label_size, bold=True, color=MUTED)
        p2 = para(tf, space_before=2)
        run(p2, text, size=size, color=BODY)
    return tf


def figpage(deck, kicker, head, png, fields, sources, *, head_size=20):
    """One full slide: heading, the figure, and the descriptive fields.

    The figure's aspect ratio picks the layout. A wide figure takes the full content width with
    the fields in a strip beneath it; a squarer one keeps the fields in a right-hand column so it
    can use the full slide height. Whichever renders the figure LARGER wins — a 3x3 panel grid is
    unreadable if it is forced into a letterbox slot.
    """
    if not os.path.exists(png):
        raise SystemExit("MISSING ARTIFACT: %s" % png)
    from PIL import Image
    with Image.open(png) as im:
        iw, ih = im.size

    strip_h = Inches(1.02) if len(fields) <= 4 else Inches(1.16)
    top = Inches(1.40)
    bottom = H - Inches(0.60)

    # candidate A — fields in a strip underneath
    a_w, a_h = CW, (bottom - strip_h) - top - Inches(0.12)
    a_scale = min(a_w / iw, a_h / ih)

    # candidate B — fields in a right-hand column
    side_w = Inches(3.30)
    b_w, b_h = CW - side_w - Inches(0.30), bottom - top
    b_scale = min(b_w / iw, b_h / ih)

    s = deck._next()
    running_head(s, kicker)
    tf = txbox(s, ML, Inches(0.80), CW, Inches(0.52))
    run(para(tf, first=True), head, size=head_size, bold=True, color=INK)
    footer(s, deck.n)

    if b_scale > a_scale * 1.06:      # switch only when it is a real gain
        pic(s, png, ML, top, b_w, b_h)
        sidepanel(s, fields, ML + b_w + Inches(0.30), top, side_w, b_h)
    else:
        pic(s, png, ML, top, a_w, a_h)
        factstrip(s, fields, y=bottom - strip_h, height=strip_h)
    provenance(s, sources)
    return s


def caption(slide, y, text, *, size=10.5):
    """A small muted line under a table — units, sign conventions, significance keys."""
    tf = txbox(slide, ML, y, CW, Inches(0.4))
    p = para(tf, first=True)
    run(p, text, size=size, color=MUTED)
    return y + Inches(0.3)


class Deck:
    """Slide counter + the two openers every deck here uses.

    ``newslide`` returns (slide, y_below_heading); ``divider`` writes a full-bleed navy act
    opener and returns the slide.
    """

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.n = 0

    def _next(self):
        self.n += 1
        return blank(self.prs)

    def newslide(self, kicker=None, head=None):
        s = self._next()
        y = Inches(0.86)
        if kicker:
            running_head(s, kicker)
        if head:
            y = heading(s, head)
        footer(s, self.n)
        return s, y

    def divider(self, eyebrow, title, sub, *, accent=PTO_C, title_size=40, y=Inches(2.5)):
        s = self._next()
        rect(s, 0, 0, W, H, INK)
        tf = txbox(s, ML, y, CW, Inches(2.6))
        p = para(tf, first=True)
        run(p, eyebrow, size=15, bold=True, color=accent)
        p = para(tf, space_before=12)
        run(p, title, size=title_size, bold=True, color=PAPER)
        p = para(tf, space_before=16)
        run(p, sub, size=15, color=SLATE)
        footer(s, self.n)
        return s

    def save(self, path, repo=None):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)
        shown = os.path.relpath(path, repo) if repo else path
        print("wrote %s  (%d slides)" % (shown, len(self.prs.slides._sldIdLst)))
