"""Build the supervisor-meeting deck for Exp3 (PTO vs GRPO) — 2026-07-16 refresh.

Full-rigor progress + results deep-dive + next-steps decision, real EDA figures,
plus a backbench appendix of the heavy stats tables (native pptx tables).

Refresh vs the 2026-07-13 deck:
- NEW HEADLINE RESULT promoted: best-vs-best paired contrast (method_paired_best) —
  PTO@10 beats GRPO even at its iter-8 peak (+0.18 Q1+Q2, dz 0.30, Holm .010) —
  preempts the "you only win because GRPO regressed" objection
- figure paths moved to the 2026-07-16 tier-based EDA reorg (7 families + 0_headline;
  behavior_drift -> miti_detail_grid superset; endpoint figures as *_final/*_best pairs)
- NEW slide: the drill-down reorg (every questionnaire now decomposes, zero oracle cost)
- Q2 reward-composition slide now uses the *_best deltas (credits GRPO at its peak)
- NEW appendix slide: the method_paired_best table
- what's-new / threats / next-steps updated accordingly
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = r"c:\Users\baruc\Desktop\Projects\Thesis_PTO_GRPO\Exp3_PTO_GRPO"
FIG  = os.path.join(ROOT, "eda", "results", "L0", "figures")
TBL  = os.path.join(ROOT, "eda", "results", "L0", "tables")
OUT  = os.path.join(ROOT, "supervisor_meeting_2026-07-16.pptx")

def f(p):  return os.path.join(FIG, *p.split("/"))
def t(p):  return os.path.join(TBL, *p.split("/"))

NAVY  = RGBColor(0x1F,0x3A,0x5F); PTO = RGBColor(0x00,0x72,0xB2)
GRPO  = RGBColor(0xE6,0x9F,0x00); GREY= RGBColor(0x5A,0x5A,0x5A)
LIGHT = RGBColor(0xEF,0xF2,0xF6); GREEN=RGBColor(0x00,0x8A,0x63)
RED   = RGBColor(0xC0,0x4A,0x1A); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x22,0x28,0x30); ROWALT=RGBColor(0xF4,0xF7,0xFA)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_N = 0

def slide(fnum=True):
    global _N
    s = prs.slides.add_slide(BLANK)
    if fnum:
        _N += 1
        tb = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"Exp3 · {_N}"; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xB0,0xB8,0xC0); r.font.name = FONT
    return s

def box(s,l,tp,w,h):
    tb = s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tb.text_frame.word_wrap = True; return tb

def rect(s,l,tp,w,h,color,line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp

def set_runs(para, segs, size, align=PP_ALIGN.LEFT):
    para.alignment = align
    for txt,bold,color in segs:
        r = para.add_run(); r.text = txt; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

def title_bar(s, text, kicker=None):
    rect(s,0,0,13.333,1.15,NAVY); rect(s,0,1.15,13.333,0.06,PTO)
    tb = box(s,0.55,0.18,12.2,0.9); tf = tb.text_frame
    p = tf.paragraphs[0]
    if kicker:
        set_runs(p,[(kicker,True,RGBColor(0x8F,0xC7,0xEC))],13)
        set_runs(tf.add_paragraph(),[(text,True,WHITE)],26)
    else:
        set_runs(p,[(text,True,WHITE)],27)

def bullets(s, items, left=0.6, top=1.5, width=12.1, size=16, gap=8):
    tb = box(s,left,top,width,5.6); tf = tb.text_frame
    first = True
    for segs,level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level = level; p.space_before = Pt(gap); p.space_after = Pt(2)
        mk = "•  " if level==0 else "–  "
        set_runs(p,[("    "*level+mk, False, PTO if level==0 else GREY)]+segs, size)

def fit(img,mw,mh):
    with Image.open(img) as im: iw,ih = im.size
    ar = iw/ih; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    return w,h

def figure(s,img,mw,mh,left,top):
    w,h = fit(img,mw,mh)
    return s.shapes.add_picture(img,Inches(left+(mw-w)/2),Inches(top+(mh-h)/2),Inches(w),Inches(h))

def caption(s,text,left,top,width,color=GREY,size=11.5,align=PP_ALIGN.CENTER):
    p = box(s,left,top,width,0.5).text_frame.paragraphs[0]
    set_runs(p,[(text,False,color)],size,align)

def fig_panel(kicker,title,img,cap,ptitle,pitems,mw=7.8):
    s = slide(); title_bar(s,title,kicker)
    figure(s,img,mw,5.2,0.35,1.5); caption(s,cap,0.35,6.75,mw)
    px = mw+0.7; tb = box(s,px,1.6,13.0-px,5.2); tf = tb.text_frame
    set_runs(tf.paragraphs[0],[(ptitle,True,NAVY)],14.5)
    for segs in pitems:
        p = tf.add_paragraph(); p.space_before = Pt(10)
        set_runs(p,[("•  ",False,PTO)]+segs,12.5)
    return s

# ---- native pptx table from a markdown table file ----
def _num(x):
    try: float(x); return True
    except: return False
def _fmt(col,v):
    if not _num(v): return v
    fv = float(v)
    if col in ("p","p_holm","wilcoxon_p") and fv == 0: return "<.001"
    s = f"{fv:.3f}".rstrip("0").rstrip(".")
    return s if s else "0"

def md_table(s, md_path, left, top, width, height, drop=(), keep=None,
             fontsize=9.0, rename=None):
    raw = [l for l in open(md_path,encoding="utf-8").read().splitlines() if l.strip().startswith("|")]
    rows = [[c.strip() for c in l.strip().strip("|").split("|")] for l in raw]
    header, body = rows[0], rows[2:]
    if keep: body = [r for r in body if keep(dict(zip(header,r)))]
    idx = [i for i,h in enumerate(header) if h not in drop]
    disp = [(rename or {}).get(header[i],header[i]) for i in idx]
    data = [[_fmt(header[i], r[i]) for i in idx] for r in body]
    nr, nc = len(data)+1, len(idx)
    gt = s.shapes.add_table(nr,nc,Inches(left),Inches(top),Inches(width),Inches(height)).table
    # proportional widths
    lens = [max(len(disp[c]), *(len(data[r][c]) for r in range(len(data)))) for c in range(nc)]
    tot = sum(lens)
    for c in range(nc):
        gt.columns[c].width = Inches(width*lens[c]/tot)
    for c,htext in enumerate(disp):
        cell = gt.cell(0,c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_top=Pt(1);cell.margin_bottom=Pt(1);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        r = pr.add_run(); r.text = htext; r.font.size=Pt(fontsize); r.font.bold=True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for ri,row in enumerate(data,1):
        for c,val in enumerate(row):
            cell = gt.cell(ri,c); cell.fill.solid()
            cell.fill.fore_color.rgb = ROWALT if ri%2 else WHITE
            cell.margin_top=Pt(0);cell.margin_bottom=Pt(0);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
            pr = cell.text_frame.paragraphs[0]
            col = DARK
            if disp[c] in ("dz","delta","mean_delta") and val.startswith("-"): col = RED
            if disp[c]=="effect" and val=="large": col = GREEN
            rr = pr.add_run(); rr.text = val; rr.font.size=Pt(fontsize)
            rr.font.color.rgb = col; rr.font.name = FONT
    return gt

def divider(text, sub):
    s = slide(fnum=False)
    rect(s,0,0,13.333,7.5,NAVY); rect(s,0,3.55,13.333,0.06,PTO)
    tb = box(s,0.9,2.7,11.5,1.8); tf = tb.text_frame
    set_runs(tf.paragraphs[0],[(text,True,WHITE)],36)
    set_runs(tf.add_paragraph(),[(sub,False,RGBColor(0xB0,0xC6,0xDA))],16)

# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.75,13.333,0.08,PTO)
tb = box(s,0.9,1.7,11.5,2.4); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Looking Ahead in Goal-Oriented Dialogue",True,WHITE)],34)
set_runs(tf.add_paragraph(),[("Exp3 — Preference-Tree vs Group-Relative optimization of a small MI therapist",False,RGBColor(0xBD,0xD6,0xEA))],18)
tb2 = box(s,0.9,5.0,11.5,1.6); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("Supervisor progress meeting · results deep-dive + measurement validity + next-steps decision",True,WHITE)],15)
set_runs(tf2.add_paragraph(),[("Lior Baruch · Reichman University · 2026-07-16",False,RGBColor(0x9F,0xB4,0xC8))],13)
p = tf2.add_paragraph(); p.space_before = Pt(8)
set_runs(p,[("Llama-3.2-1B therapist (bf16)  ·  gpt-4o-mini patient + oracle  ·  96 personas",False,RGBColor(0x7E,0x93,0xA8))],12)

# =====================================================================
# 2 · WHERE WE ARE
# =====================================================================
s = slide(); title_bar(s,"Where we are",kicker="STATUS")
rows = [("Arm","Iterations","Status",None,True),
    ("PTO  K=0","base + 1–10","COMPLETE (10) — fully scored",GREEN,False),
    ("GRPO K=0","base + 1–10","COMPLETE (10) — fully scored",GREEN,False),
    ("PTO  K=5","base + 1–5 trained","PAUSED — I1–4 scored; iter-5 unscored*",RED,False),
    ("GRPO K=5","base + 1","PAUSED — iter 1 trained + scored",RED,False)]
y=1.5
for i,(a,b,c,col,hdr) in enumerate(rows):
    rc = NAVY if hdr else (LIGHT if i%2 else WHITE)
    rect(s,0.6,y,8.4,0.55,rc,line=RGBColor(0xD0,0xD6,0xDD))
    set_runs(box(s,0.75,y+0.06,2.3,0.45).text_frame.paragraphs[0],[(a,hdr,WHITE if hdr else DARK)],14)
    set_runs(box(s,3.1,y+0.06,2.3,0.45).text_frame.paragraphs[0],[(b,hdr,WHITE if hdr else DARK)],13)
    set_runs(box(s,5.3,y+0.06,3.6,0.45).text_frame.paragraphs[0],[(c,hdr,WHITE if hdr else (col or DARK))],11.5)
    y+=0.55
set_runs(box(s,0.6,y+0.02,8.4,0.35).text_frame.paragraphs[0],
    [("* the iter-5 adapter exists but its eval conversations were never generated — one generate-only pass (no training) buys a 5th K=5 point.",False,GREY)],10.5)
rect(s,9.3,1.5,3.45,3.35,LIGHT)
tb = box(s,9.5,1.62,3.1,3.1); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Three headlines",True,NAVY)],15)
for a,bd,c in [("PTO wins even vs GRPO's PEAK: best-vs-best ","+0.18 (dz 0.30, p .010)","; matched endpoint 4.26 vs 3.75"),
        ("GRPO peaks @ iter 8 then ","regresses"," into sycophancy"),
        ("Halo-rubric gains come ","with"," MI-inconsistency rising ~2.3× (PTO) / ~4× (GRPO)")]:
    p = tf.add_paragraph(); p.space_before = Pt(9)
    set_runs(p,[("•  "+a,False,DARK),(bd,True,PTO),(c,False,DARK)],12.5)
rect(s,0.6,y+0.45,8.4,0.7,RGBColor(0xFB,0xEE,0xE6))
set_runs(box(s,0.78,y+0.54,8.1,0.55).text_frame.paragraphs[0],
    [("Binding constraint:  ",True,RED),("OpenAI spend ≈ $300 → both K=5 arms paused → RQ-i (look-ahead) on hold",False,DARK)],13)

# =====================================================================
# 3 · NEW SINCE LAST MEETING
# =====================================================================
s = slide(); title_bar(s,"New since the last meeting (07-13) — the steelman result + the drill-down EDA",kicker="WHAT'S NEW")
bullets(s,[
 ([("NEW RESULT — best-vs-best paired contrast: ",True,NAVY),("PTO at ITS best (iter 10) vs GRPO at ITS best (iter 8, the peak): ",False,DARK),("PTO still wins Q1+Q2 by +0.18 (dz 0.30, Holm p .010)",True,PTO),(" — and WAI-SR / CSQ-8 / MI-SAT / PCT too (MITI, MICI n.s.). Kills the \"you only win because GRPO regressed\" objection: GRPO is credited at its peak and still loses.",False,DARK)],0),
 ([("EDA reorganized into a tier-based drill-down ",True,NAVY),("(07-16): global scores → inside-every-questionnaire → validity/hacking. Every endpoint artifact now reported at ",False,DARK),("final AND best",True,DARK),(" (one selected checkpoint per arm, own training oracle); a ",False,DARK),("0_headline/",True,PTO),(" folder holds the 7 presentation figures, always in sync. No science change — self-check green, headline means reproduce.",False,DARK)],0),
 ([("Every questionnaire now decomposes ",True,NAVY),("— item-level detail grids for Q1 / WAI-SR / CSQ-8 / MI-SAT are NEW (the per-item scores were already on disk → ",False,DARK),("zero oracle cost",True,GREEN),("); Q2 items, MITI behaviours, MICI and PCT detail existed and moved into the same uniform style, each with \"which items drive the change\" bars at final + best.",False,DARK)],0),
 ([("Judge-reliability pipeline STILL PENDING ",True,RED),("(built 07-12, gated, ≈ $5): oracle ICC + a second judge from a different model family + the contrast-preservation check → the single-oracle circularity threat. Decision carried over to today.",False,DARK)],0),
],size=15.5,gap=12)

# =====================================================================
# 3b · NEW · the tier-based drill-down (one slide, example = Q1 items)
# =====================================================================
fig_panel("EDA · NEW","Every questionnaire now decomposes — the drill-down EDA",
  f("2_questionnaires/q1_detail_grid.png"),
  "Example: Q1 (the satisfaction half of the training reward) split into its 5 items — content/motivation items carry the gain; GRPO's iter-9/10 drop is visible in every item. Zero oracle cost (per-item scores were already stored).",
  "One system, three levels",[
   [("L1 · global scores ",True,PTO),("— the all-metric grid (next slides) + final-vs-best endpoint pairs",False,DARK)],
   [("L2 · inside each rubric ",True,PTO),("— the same small-multiples style for all 8 instruments: items (Q1/Q2/WAI/CSQ/MI-SAT), behaviours (MITI/MICI), patient signals (PCT)",False,DARK)],
   [("L3 · cross-cutting ",True,PTO),("— factor structure, reward-hacking, heterogeneity, reliability",False,DARK)],
   [("Every level answers \"which component moved, and was it still there at the arm's BEST iteration?\"",False,DARK)],
   [("0_headline/ = the 7 meeting figures, auto re-saved — this deck reads straight from it",False,GREY)],
  ],mw=7.6)

# =====================================================================
# 4 · RECAP
# =====================================================================
s = slide(); title_bar(s,"Recap — the setup and the three questions",kicker="FRAMING")
bullets(s,[
 ([("Task: ",True,NAVY),("train a small LLM therapist to do Motivational Interviewing vs a simulated patient; reward = a larger oracle grading validated MI questionnaires.",False,DARK)],0),
 ([("RQ-i  (look-ahead K∈{0,5}): ",True,PTO),("does anticipating K future turns help?  ",False,DARK),("— ON HOLD (K=5 paused for cost)",True,RED)],0),
 ([("RQ-ii (PTO vs GRPO at matched K + MCL): ",True,PTO),("is iterative GRPO competitive with preference-tree PTO?  ",False,DARK),("— the active question",True,GREEN)],0),
 ([("RQ-iii (oracle questionnaire): ",True,PTO),("held for later.",False,DARK)],0),
 ([("The methodological crux — a deliberate gap:",True,NAVY)],0),
 ([("Training reward = ",False,DARK),("Q1+Q2 only",True,NAVY),(" (22-item global-eval rubric), scored on ",False,DARK),("partial",True,RED),(" conversations.",False,DARK)],1),
 ([("Evaluation = ",False,DARK),("all 6 questionnaires + orthogonal axes",True,NAVY),(", scored on ",False,DARK),("full",True,GREEN),(" conversations.",False,DARK)],1),
 ([("That gap is how we catch a model that games the short training reward instead of doing real MI.",False,GREY)],1),
],size=16.5,gap=9)

# =====================================================================
# 5 · DESIGN / pairing
# =====================================================================
s = slide(); title_bar(s,"The design that earns the statistical power",kicker="RIGOR")
bullets(s,[
 ([("96 personas recur across every iteration and both methods",True,NAVY),(" — nothing is independent-sample.",False,DARK)],0),
 ([("Every comparison is ",False,DARK),("paired by persona (n = 96 paired deltas)",True,PTO),(" — controls persona difficulty; far stronger than an independent Mann–Whitney.",False,DARK)],0),
 ([("Test battery (all repeated-measures-aware):",True,NAVY)],0),
 ([("vs-base / PTO-vs-GRPO: ",True,DARK),("Wilcoxon signed-rank + Cohen's dz + 2000-resample bootstrap CI (fixed seed).",False,DARK)],1),
 ([("Trajectory over iterations: ",True,DARK),("Friedman χ² + Kendall's W",False,DARK),("  (Spearman ρ / OLS slope reported descriptive-only — refuses a pseudo-replicated p).",False,GREY)],1),
 ([("Multiplicity: ",True,DARK),("Holm–Bonferroni, family = rubrics within one arm / one iteration (not the whole grid).",False,DARK)],1),
 ([("Honest boundary: ",True,RED),("pairing controls persona difficulty, NOT patient-simulator stochasticity — matched-subjects, not a deterministic re-run.",False,DARK)],0),
],size=16,gap=9)

# =====================================================================
# 6 · HEADLINE trajectory
# =====================================================================
fig_panel("RESULT · RQ-ii","Headline — PTO climbs stably; GRPO peaks then regresses",
  f("1_outcomes/trajectories/trajectory_Q1Q2.png"),
  "Q1+Q2 (training-reward rubric) over iterations, persona-paired. Peak auto-flagged where an arm regresses after.",
  "Read-off",[
   [("PTO ",True,PTO),("3.00 → 4.26",True,DARK),(" @ iter 10 (final = peak); dz 1.43 large",False,GREY)],
   [("GRPO ",True,GRPO),("3.07 → 4.08 @ iter 8",True,DARK),(", falls to 3.75 by iter 10",False,GREY)],
   [("Climb rate (OLS): PTO ",False,DARK),("0.120/iter",True,PTO),(" vs GRPO ",False,DARK),("0.072/iter",True,GRPO)],
   [("Friedman: both change reliably — PTO W=0.45, GRPO W=0.33 (n=96, k=11)",False,DARK)],
   [("⇒ GRPO needs early stopping; PTO sustains gains",True,GREEN)],
  ])

# =====================================================================
# 7 · all-metric grid
# =====================================================================
s = slide(); title_bar(s,"Same pattern across the whole battery",kicker="RESULT")
figure(s,f("1_outcomes/trajectories_all_metrics.png"),12.4,5.2,0.45,1.45)
caption(s,"All eval metrics over iterations. The 5 global-eval (halo) rubrics rise together; MICI (lower=better) rises too — the reward-hack. PTO peak-then-flat, GRPO peak-then-drop.",0.6,6.75,12.1)

# =====================================================================
# 8 · effect forest
# =====================================================================
s = slide(); title_bar(s,"Effect vs base — every rubric, with CIs and dz",kicker="RESULT")
figure(s,f("1_outcomes/effect_vs_base_forest_final.png"),7.9,5.2,0.4,1.5)
caption(s,"Δ-vs-base per arm×rubric at the FINAL iteration, 95% bootstrap CI + Cohen's dz. MICI direction-colored (a positive Δ is bad). The *_best companion (each arm at its selected checkpoint) is in 1_outcomes/.",0.4,6.75,7.9)
tb = box(s,8.5,1.6,4.4,5.2); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("PTO K=0 vs base (final)",True,PTO)],14)
for x in ["Q1+Q2  Δ+1.26  dz 1.43  large","MITI   Δ+1.14  dz 1.35  large","MI-SAT Δ+0.82  dz 0.90  large","every halo rubric: large, Holm p≈0"]:
    p=tf.add_paragraph();p.space_before=Pt(5);set_runs(p,[("•  "+x,False,DARK)],12.5)
p=tf.add_paragraph();p.space_before=Pt(12);set_runs(p,[("GRPO K=0 vs base",True,GRPO)],14)
for x in ["final (iter10) Q1+Q2 dz 0.72 medium","best  (iter8)  Q1+Q2 dz 1.22 large","MICI dz 1.72 large — worst offender"]:
    p=tf.add_paragraph();p.space_before=Pt(5);set_runs(p,[("•  "+x,False,DARK)],12.5)

# =====================================================================
# 9 · PTO vs GRPO ladder
# =====================================================================
s = slide(); title_bar(s,"PTO vs GRPO — paired, and the honest ladder",kicker="RIGOR · RQ-ii")
data = [("Iter","PTO−GRPO Q1+Q2","dz","Holm p","note"),
  ("1–2","≈ 0","~0","1.00","tie"),("3","−0.179","−0.33","0.014","GRPO briefly ahead"),
  ("8","+0.138","+0.31","0.028","GRPO at its peak"),("9","+0.431","+0.77","<0.001","GRPO collapses"),
  ("10","+0.507","+0.73","<0.001","matched endpoint"),
  ("best","+0.177","+0.30","0.010","PTO@10 vs GRPO@8 (best-vs-best) — NEW")]
y=1.5; colx=[0.6,1.7,4.3,5.3,6.4]; colw=[1.1,2.6,1.0,1.1,3.0]
for i,row in enumerate(data):
    hdr=i==0; rc=NAVY if hdr else (LIGHT if i%2 else WHITE)
    rect(s,0.6,y,8.9,0.5,rc,line=RGBColor(0xD0,0xD6,0xDD))
    for cx,cw,val in zip(colx,colw,row):
        col = WHITE if hdr else (PTO if val.startswith("+") else (RED if val.startswith("−") else DARK))
        set_runs(box(s,cx,y+0.04,cw,0.42).text_frame.paragraphs[0],[(val,hdr,col)],12.5)
    y+=0.5
rect(s,9.7,1.5,3.1,4.4,LIGHT)
tb=box(s,9.88,1.62,2.8,4.2); tf=tb.text_frame
set_runs(tf.paragraphs[0],[("What a committee will press",True,NAVY)],13.5)
for segs in [[("The endpoint win is largely GRPO's ",False,DARK),("late collapse",True,RED),(", not PTO pulling away.",False,DARK)],
 [("Answer — now a FORMAL paired contrast (",False,DARK),("method_paired_best",True,GREEN),(", NEW): PTO@its-best vs GRPO@its-best = ",False,DARK),("+0.18, dz 0.30, Holm .010",True,GREEN),(" — PTO also wins WAI-SR / CSQ-8 / MI-SAT / PCT (MITI, MICI n.s.).",False,DARK)],
 [("Persona pairing holds across different iterations — every iteration reshuffles the same 96 personas.",False,GREY)],
 [("iter-9 drop localized: it8→it9 dz −0.41 (Holm 0); it9→it10 n.s. → a step-down, then plateau.",False,GREY)]]:
    p=tf.add_paragraph();p.space_before=Pt(8);set_runs(p,[("•  ",False,PTO)]+segs,11)
caption(s,"Holm family = rubrics within each contrast (not re-pooled across iterations). Full best-vs-best table in the appendix.",0.6,6.35,8.9,align=PP_ALIGN.LEFT)

# =====================================================================
# 10 · reward hack panel
# =====================================================================
fig_panel("MECHANISM","The gains come WITH a measurable reward-hack",
  f("3_validity/reward_hack_panel.png"),
  "Per arm: the reward proxy (Q1+Q2, left) climbs while MI-inconsistency (MICI, right) climbs with it and real change-talk (PCT) barely moves.",
  "“All rubrics up” ≠ multi-skill",[
   [("MI-inconsistency rises: MICI 0.21 → ",False,DARK),("0.49 PTO (~2.3×) / 0.84 GRPO (~4×)",True,DARK),(" @ iter 10",False,DARK)],
   [("Real change-talk (PCT) rises only modestly — more for PTO (0.49→0.63) than GRPO (0.49→0.57)",False,DARK)],
   [("Both kill the early degeneration loops (loop% 0.49 → 0)",False,DARK)],
   [("The halo-rubric gain is partly over-praise/advice in ",False,DARK),("both",True,RED),(" methods — worse in GRPO",False,DARK)],
  ])

# =====================================================================
# 11 · behavior drift
# =====================================================================
fig_panel("MECHANISM","Mechanism — what the therapist actually does",
  f("2_questionnaires/miti_detail_grid.png"),
  "The full MITI drill-down (new grid): 4 global ratings + all 7 behaviour rates per therapist turn + the proficiency ratios (rates, so conv length can't inflate them).",
  "Affirmation drift — in BOTH arms",[
   [("GRPO @ iter 10: ",True,GRPO),("B6-Affirm 0.52→1.98, questions 0.83→0.15/turn, R:Q→1.44",False,DARK)],
   [("PTO @ iter 10: ",True,PTO),("milder & plateaus — B6-Affirm 1.64, 0.55 q/turn",False,DARK)],
   [("Across all 96 iter-10 convs: GRPO ",False,DARK),("~3.5× more praise, ~4× fewer questions",True,RED)],
   [("GRPO's late eval regression IS this over-praise, which the full-conv oracle penalizes",False,GREY)],
  ],mw=8.0)

# =====================================================================
# 12 · NEW · MITI 4.2.1 thresholds
# =====================================================================
fig_panel("MECHANISM","Better than base ≠ competent — official MITI 4.2.1 thresholds",
  f("2_questionnaires/miti_proficiency_thresholds.png"),
  "The 4 official MITI summary scores vs the manual's fair (amber) / good (green) competency lines (Moyers et al., 4.2.1 §H–I).",
  "The absolute anchor",[
   [("Both arms go below-competence → fair-to-good on the ",False,DARK),("global ratings",True,NAVY),("; Relational crosses \"good\"",False,DARK)],
   [("Neither arm reaches \"good\" on the ",False,DARK),("technique ratios",True,RED),(" (R:Q, %CR)",False,DARK)],
   [("GRPO's R:Q 1.43 \"fair\" is the ",False,DARK),("pathological route",True,GRPO),(" — fewer questions, not more reflections",False,DARK)],
   [("Caveats stated with the lines: thresholds are expert opinion (the manual's own words) and defined for ~20-min human audio sessions",False,GREY)],
  ])

# =====================================================================
# 13 · NEW (07-09) · MICI behaviour detail
# =====================================================================
fig_panel("MECHANISM","Which MI-inconsistent moves are rising? (MICI detail)",
  f("2_questionnaires/mici_detail_grid.png"),
  "Each of the 6 MI-inconsistent behaviours + severity, per therapist turn (lower = better).",
  "Decomposing the MICI rise",[
   [("The MICI climb is driven mainly by ",False,DARK),("over-praise / sycophancy",True,RED),(" and unsolicited advice",False,DARK)],
   [("Confront / warn / judge stay low — it's not hostility, it's ",False,DARK),("misplaced warmth",True,NAVY)],
   [("GRPO's curves climb steeper than PTO's on the praise-type behaviours",False,GREY)],
   [("Ties the aggregate MICI number to concrete, codeable moves",False,GREY)],
  ])

# =====================================================================
# 14 · PCT patient detail
# =====================================================================
fig_panel("MECHANISM","Did the patient actually move? (change-talk detail)",
  f("2_questionnaires/pct_detail_grid.png"),
  "PCT globals (Importance / Confidence / Readiness) + change / sustain / neutral proportions of patient utterances.",
  "Real outcome, not just felt-warmth",[
   [("PCT is the ",False,DARK),("patient-side",True,NAVY),(" axis — did the client express motivation, not just rate the therapist highly?",False,DARK)],
   [("Change-talk proportion rises modestly; sustain-talk (lower=better) eases",False,DARK)],
   [("More movement for PTO than GRPO — consistent with the headline",False,DARK)],
   [("Modest effect: the halo-rubric gains outpace real behaviour change",False,RED)],
  ])

# =====================================================================
# 15 · NEW · Q2 item-level reward composition
# =====================================================================
fig_panel("MECHANISM","Which Q2 items drive the reward? — the incentive, decomposed",
  f("2_questionnaires/q2_item_deltas_best.png"),
  "Δ vs base per Q2 item at each arm's BEST iteration (GRPO credited at its iter-8 peak), colored by face-content group (our analytical grouping — NOT a validated subscale).",
  "The reward composition pays for the drift",[
   [("“Revealed his thinking” (self-disclosure) ",True,NAVY),("tops BOTH arms' Δ ranking",True,DARK)],
   [("Q2 items 1/2/3/10 reward therapist ",False,DARK),("self-disclosure",True,RED),(" — behaviour MI does not prescribe",False,DARK)],
   [("So the emotive drift isn't only oracle leniency: it is ",False,DARK),("what the Q1+Q2 training reward pays for",True,NAVY)],
   [("The group-trajectory companion shows the exploited components take off early, not just at the endpoint",False,GREY)],
   [("Direct thesis ammunition for the reward-design discussion (RQ-iii link)",False,GREY)],
  ])

# =====================================================================
# 16 · factor loadings (halo)
# =====================================================================
s = slide(); title_bar(s,"One global-eval halo factor — the orthogonal axes earn their keep",kicker="RIGOR")
figure(s,f("3_validity/factor_loadings.png"),7.7,5.2,0.4,1.5)
caption(s,"PC1/PC2 loadings on the standardized 10-metric factor space (global-eval rubrics + orthogonal axes).",0.4,6.75,7.7)
tb=box(s,8.3,1.6,4.6,5.2); tf=tb.text_frame
set_runs(tf.paragraphs[0],[("The halo is one factor — not skill breadth",True,NAVY)],14.5)
for segs in [[("5 global-eval rubrics alone → ",False,DARK),("PC1 ≈ 91%",True,RED),(" of variance",False,DARK)],
 [("Add PCT / MICI / R:Q / %CR / %MICO → ",False,DARK),("PC1 drops to ≈ 55%",True,GREEN),(" per arm",False,DARK)],
 [("A genuine second factor (technique + MI-inconsistency) appears",False,DARK)],
 [("Honest read: 55% still means the halo is over half the variance",False,GREY)],
 [("Structural fact about the score matrix — not a null-test → hard to attack",True,GREEN)]]:
    p=tf.add_paragraph();p.space_before=Pt(9);set_runs(p,[("•  ",False,PTO)]+segs,12.5)

# =====================================================================
# 17 · rubric correlation
# =====================================================================
fig_panel("RIGOR","The global-eval (halo) rubrics are highly inter-correlated",
  f("3_validity/rubric_correlation.png"),
  "Spearman correlation among all metrics (per-conversation, pooled). Diverging scale; heavy line = halo block boundary.",
  "Why we needed orthogonal axes",[
   [("The 5 global-eval rubrics form a tight ",False,DARK),("high-correlation block",True,RED),(" — they largely move together",False,DARK)],
   [("MICI and the technique ratios sit ",False,DARK),("outside",True,GREEN),(" that block (weak / opposite correlation)",False,DARK)],
   [("Visual companion to the PCA: correlation → the single halo factor",False,GREY)],
   [("PCT co-moves partly with the halo (loads ~0.39 on PC1)",False,GREY)],
  ])

# =====================================================================
# 18 · subscale trajectories
# =====================================================================
s = slide(); title_bar(s,"Subscale detail — WAI-SR (Goal/Task/Bond) & MITI globals",kicker="MECHANISM")
figure(s,f("2_questionnaires/wai_subscales.png"),12.4,5.15,0.45,1.45)
caption(s,"Opening the aggregate rubrics: WAI-SR alliance subscales and the 4 MITI global ratings over iterations — the gains aren't concentrated in a single sub-item.",0.6,6.72,12.1)

# =====================================================================
# 19 · heterogeneity
# =====================================================================
s = slide(); title_bar(s,"GRPO's collapse concentrates on the resistant personas",kicker="HETEROGENEITY")
figure(s,f("4_heterogeneity/subgroup_endpoint_cooperation_level_final.png"),11.8,5.1,0.75,1.5)
caption(s,"Final-iteration score per cooperation level × arm (the *_best companion bars GRPO at iter 8). GRPO's endpoint regression is worst on low-cooperation (Resistant) personas — where sycophancy fails and concrete guidance is demanded.",0.75,6.7,11.8)

# =====================================================================
# 20 · reward faithfulness
# =====================================================================
fig_panel("RELIABILITY","Is the short training reward faithful? — why MCL exists",
  f("5_training/reward_reliability_curve.png"),
  "Rank agreement between the partial-conv training reward and the full-conv eval, vs conversation length.",
  "The proxy is unreliable when short",[
   [("At n_turns=2, agreement ≈ ",False,DARK),("0.66–0.73",True,RED),(" — barely above chance (0.5)",False,DARK)],
   [("Clears 0.8 at ~10 turns, 0.9 at ~30 (monotone)",False,DARK)],
   [("Motivates ",False,DARK),("MIN_CONV_LENGTH",True,NAVY),(" — drop training slices below N utterances; set to 12 in both trainers",False,DARK)],
   [("At MCL=12: GRPO proxy grows MORE faithful (0.86→0.94), PTO LESS (0.87→0.76) — both stay out of the bad regime",False,GREY)],
  ])

# =====================================================================
# 21 · reward distribution by arm
# =====================================================================
fig_panel("TRAINING","Candidate-reward distributions — the training signal",
  f("5_training/reward_distribution_by_arm.png"),
  "Per-candidate oracle reward distribution by arm/iteration (from generations.jsonl).",
  "What the update actually saw",[
   [("The reward the optimizer trained on, ",False,DARK),("before",True,NAVY),(" the eval pass",False,DARK)],
   [("Distributions shift ",False,DARK),("upward",True,GREEN),(" over iterations — the policy really is finding higher-scored candidates",False,DARK)],
   [("Spread stays wide → enough within-group contrast for GRPO advantages / PTO pairs",False,DARK)],
   [("Confirms the signal isn't collapsing to a degenerate spike",False,GREY)],
  ],mw=8.0)

# =====================================================================
# 22 · preference probe
# =====================================================================
fig_panel("PTO INTERNALS","PTO's DPO signal is real (preference probe)",
  f("6_preference/PTO_LA0_pref_word_ranking.png"),
  "Mass-mean-probe word ranking along the mean(chosen − rejected) direction.",
  "The pair signal separates — and drifts",[
   [("wins_correct 0.65 → 0.71",True,PTO),(" over iters (>0.5 = the chosen−rejected direction genuinely separates pairs)",False,DARK)],
   [("Strengthens late — the DPO update has real signal, not noise",False,DARK)],
   [("Its latent target drifts toward affirmation / achievement language — the latent echo of the behaviour drift",False,GREY)],
  ],mw=8.0)

# =====================================================================
# 23 · threats to validity
# =====================================================================
s = slide(); title_bar(s,"Threats to validity — ranked, with the answer",kicker="RIGOR")
bullets(s,[
 ([("1 · Single-oracle circularity (biggest). ",True,RED),("Training reward & eval are the SAME gpt-4o-mini; no human validation. “PTO higher” could be “PTO games this oracle better.”",False,DARK)],0),
 ([("Defence: ",True,GREEN),("the orthogonal axes show the oracle PENALIZES the hack (MICI up, questions down) — but same model scores them. ",False,DARK),("The judge-reliability run (next slide) bounds this empirically for ≈ $5.",True,NAVY)],1),
 ([("2 · Endpoint privileging. ",True,RED),("Why iter 10? ",False,DARK),("Answer: pre-committed matched-budget = 10 iters, AND the NEW formal best-vs-best contrast (+0.18, dz 0.30, Holm .010; wins on 4 further rubrics) credits GRPO at its peak and still favours PTO.",False,DARK)],0),
 ([("3 · One draw per (model, persona). ",True,RED),("No within-cell replication → conversation noise loads into the persona effect; dz slightly overstates precision. The ICC run also quantifies the oracle's share of this noise.",False,DARK)],0),
 ([("4 · Kendall's W modest (0.33–0.45). ",True,RED),("Significant ≠ large — frame as reliable directional change.",False,DARK)],0),
 ([("5 · Absolute scores are Exp3-internal only ",True,RED),("(bf16 vs Exp2's 4-bit) — never cross-compare experiments.",False,DARK)],0),
],size=15,gap=9)

# =====================================================================
# 24 · NEW · judge-reliability proposal
# =====================================================================
s = slide(); title_bar(s,"Proposal — the ≈ $5 measurement-validity run (built, gated, ready)",kicker="PROPOSAL · PENDING")
bullets(s,[
 ([("What: ",True,NAVY),("re-score an anchor subset — base, PTO iter-10, GRPO iter-8 (peak), GRPO iter-10 — on Q1, Q2, MICI. Everything is resume-safe and writes OUTSIDE the real eval scores.",False,DARK)],0),
 ([("Part 1 · Oracle repeatability (ICC): ",True,PTO),("same oracle × 3 reps with per-rep seeds → ICC(2,1) + mean |Δ| — upgrades the “oracle noise ≈ 0.10” folklore into a citable statistic.  (~3,460 gpt-4o-mini calls ≈ $1.6)",False,DARK)],0),
 ([("Part 2 · Second judge: ",True,PTO),("a judge from a DIFFERENT model family (default: Claude Haiku; pluggable) scores the same conversations → per-metric agreement (r/ρ/bias) ",False,DARK),("and the defense-critical check: does the PTO−GRPO endpoint contrast survive a judge that never simulated the patient?",True,NAVY),("  (~1,150 calls ≈ $3.7)",False,DARK)],0),
 ([("Why it matters: ",True,GREEN),("threat #1 (patient = oracle = same model) goes from “acknowledged” to “empirically bounded” — the cheapest defensibility purchase available.",False,DARK)],0),
 ([("Asks: ",True,RED),("approve the ≈ $5 spend · confirm the second-judge model (Haiku economical / Sonnet stronger) · agree the anchor-subset choice.",False,DARK)],0),
],size=15.5,gap=12)
rect(s,0.6,6.15,12.1,0.75,LIGHT)
set_runs(box(s,0.8,6.28,11.7,0.55).text_frame.paragraphs[0],
  [("Interpretation guard rails are pre-registered in the notebook: ICC ≥ 0.75 = good; cross-judge r read against the ICC ceiling √(ICC₁·ICC₂), not against 1.0.",False,GREY)],12)

# =====================================================================
# 25 · DECISION
# =====================================================================
s = slide(); title_bar(s,"Decision needed — RQ-i (look-ahead) vs the budget",kicker="DECISION")
rect(s,0.6,1.45,12.1,1.15,RGBColor(0xFB,0xEE,0xE6))
tb=box(s,0.8,1.55,11.7,1.0); tf=tb.text_frame
set_runs(tf.paragraphs[0],[("The constraint: ",True,RED),("OpenAI ≈ $300 spent; cost ∝ candidate count (prompts×G / branch×M) × iterations. ",False,DARK)],14)
set_runs(tf.add_paragraph(),[("Caching is already maxed (~50% off the oracle prefix) → the ",False,DARK),("only lever is call COUNT",True,NAVY),(". First K=5 point is nearly free: a generate-only pass with the existing PTO iter-5 adapter (96 convs, no training).",False,DARK)],14)
opts=[("A. Resume one K=5 arm, cost-capped","Free iter-5 point first, then M/G 8→4, cap ~5–6 iters, compare at matched iter. Makes RQ-i conclusive. Keep K + gpt-4o-mini oracle FIXED.",GREEN,"Recommended"),
 ("B. Declare RQ-ii the thesis core","Write up PTO>GRPO + reward-hack now; RQ-i stays preliminary / future work (still take the free iter-5 point).",NAVY,""),
 ("C. Spend to full K=5 (10 iters)","Cleanest RQ-i answer but ~another few-hundred $ — hard to justify given curves plateau by iter ~4.",RED,"")]
y=2.85
for name,desc,col,tag in opts:
    rect(s,0.6,y,12.1,1.2,WHITE,line=RGBColor(0xCF,0xD6,0xDE)); rect(s,0.6,y,0.14,1.2,col)
    tb=box(s,0.95,y+0.12,11.5,1.0); tf=tb.text_frame
    seg=[(name,True,col)]
    if tag: seg.append(("    ("+tag+")",True,GREEN))
    set_runs(tf.paragraphs[0],seg,15); set_runs(tf.add_paragraph(),[(desc,False,DARK)],12.5)
    y+=1.32

# =====================================================================
# 26 · next steps
# =====================================================================
s = slide(); title_bar(s,"Proposed next steps & asks",kicker="NEXT")
bullets(s,[
 ([("Decide RQ-i path ",True,NAVY),("(A/B/C above) — my recommendation: ",False,DARK),("A, free iter-5 point now + a cost-capped K=5 resume",True,GREEN)],0),
 ([("Approve the ≈ $5 judge-reliability run ",True,NAVY),("(ICC + second judge + contrast preservation) — answers the circularity threat; results feed LIMITATIONS + the defense.",False,DARK)],0),
 ([("Start the thesis chapter ",True,NAVY),("on RQ-ii: PTO sustains gains; GRPO peaks then reward-hacks into sycophancy; orthogonal axes break the halo; the Q2 composition explains WHY the drift is incentivized.",False,DARK)],0),
 ([("Framing to agree on: ",True,NAVY),("lead the PTO>GRPO claim with the formal best-vs-best contrast (+0.18, dz 0.30 — NEW), use the endpoint (+0.51) as the stability story; report MITI thresholds honestly (neither arm technique-\"good\").",False,DARK)],0),
 ([("Done since last time: ",True,GREY),("EDA reorganized into the tier-based drill-down + 0_headline (07-16) · best-vs-best contrast formalized (method_paired_best) · item-level detail for every questionnaire at zero oracle cost · all committed + pushed.",False,DARK)],0),
],size=16.5,gap=13)
rect(s,0.6,6.05,12.1,0.9,LIGHT)
set_runs(box(s,0.8,6.2,11.7,0.7).text_frame.paragraphs[0],
  [("Bottom line:  ",True,NAVY),("RQ-ii is in hand and defensible; RQ-i needs a small budget decision; $5 buys the measurement-validity answer.",False,DARK)],15)

# =====================================================================
# APPENDIX
# =====================================================================
divider("Appendix — heavy stats tables","Persona-paired, n = 96. Full-conversation eval. All tables reproducible from eda/results/L0/tables/.")

# A1 · main results FINAL
s = slide(); title_bar(s,"Main results — each arm vs base at the FINAL iteration",kicker="APPENDIX · main_results")
md_table(s, t("7_stats/main_results.md"), 0.4, 1.45, 12.55, 5.3,
   drop=("target","target_iter","wilcoxon_p","traj_rho"),
   keep=lambda r: r["target"]=="final",
   rename={"traj_slope":"slope","ci_low":"ci_lo","ci_high":"ci_hi"}, fontsize=9.5)
caption(s,"delta = paired mean(target − base); dz = Cohen's paired d; p_holm = Holm across rubrics within arm; slope = OLS/iter. MICI is lower-is-better.",0.4,6.9,12.5,align=PP_ALIGN.LEFT,size=10)

# A2 · main results BEST
s = slide(); title_bar(s,"Main results — each arm vs base at its BEST iteration",kicker="APPENDIX · main_results")
md_table(s, t("7_stats/main_results.md"), 0.4, 1.45, 12.55, 5.3,
   drop=("target","wilcoxon_p","traj_rho"),
   keep=lambda r: r["target"]=="best",
   rename={"target_iter":"best_it","traj_slope":"slope","ci_low":"ci_lo","ci_high":"ci_hi"}, fontsize=9.5)
caption(s,"GRPO's best is iter 8 (4.08); PTO's best is iter 10 (4.26) — PTO peaks later and higher.",0.4,6.9,12.5,align=PP_ALIGN.LEFT,size=10)

# A3 · vs_base paired (Q1Q2 by iter)
s = slide(); title_bar(s,"Q1+Q2 vs base — every iteration, paired",kicker="APPENDIX · vs_base_paired")
md_table(s, t("7_stats/vs_base_paired.md"), 2.0, 1.45, 9.3, 5.3,
   rename={"iteration":"iter","mean_delta":"Δ vs base","ci_low":"ci_lo","ci_high":"ci_hi"}, fontsize=10)
caption(s,"Both arms climb hard early; PTO keeps climbing to iter 10 (dz 1.43) while GRPO peaks at iter 8 (dz 1.22) then fades.",2.0,6.9,9.3,align=PP_ALIGN.LEFT,size=10)

# A4 · method paired (two cuts)
s = slide(); title_bar(s,"PTO − GRPO paired — Q1+Q2 by iteration & all metrics @ iter 10",kicker="APPENDIX · method_paired_by_K")
md_table(s, t("7_stats/method_paired_by_K.md"), 0.45, 1.5, 5.9, 5.1,
   drop=("K","n"), keep=lambda r: r["metric"]=="Q1Q2",
   rename={"iteration":"iter","mean_delta":"Δ (PTO−GRPO)"}, fontsize=9.5)
caption(s,"Q1+Q2, all iterations",0.45,6.7,5.9,align=PP_ALIGN.CENTER,size=10)
md_table(s, t("7_stats/method_paired_by_K.md"), 6.7, 1.5, 6.2, 5.1,
   drop=("K","n","iteration"), keep=lambda r: r["iteration"]=="10.000",
   rename={"mean_delta":"Δ (PTO−GRPO)"}, fontsize=9.5)
caption(s,"All metrics @ iter 10 (+ = PTO higher; MICI − = PTO less inconsistent = good)",6.7,6.7,6.2,align=PP_ALIGN.CENTER,size=10)

# A4b · NEW · best-vs-best (method_paired_best)
s = slide(); title_bar(s,"PTO@best vs GRPO@best — the model-selection contrast (NEW)",kicker="APPENDIX · method_paired_best")
md_table(s, t("7_stats/method_paired_best.md"), 2.0, 1.5, 9.3, 4.6,
   rename={"iter_a":"PTO iter","iter_b":"GRPO iter","mean_delta":"Δ (PTO−GRPO)"}, fontsize=10)
caption(s,"Each method at its own-oracle BEST iteration (PTO@10 vs GRPO@8), persona-paired across the different iterations (valid — every iteration reshuffles the same 96 personas). + = PTO higher; MICI negative = PTO less MI-inconsistent (good, n.s.). Holm across rubrics within the contrast.",2.0,6.4,9.3,align=PP_ALIGN.LEFT,size=10)

# A5 · Friedman omnibus
s = slide(); title_bar(s,"Friedman omnibus — reliable change over iterations",kicker="APPENDIX · friedman_omnibus")
md_table(s, t("7_stats/friedman_omnibus.md"), 1.7, 1.5, 9.9, 5.1,
   rename={"kendall_w":"Kendall W","k_iters":"k","n_personas":"n"}, fontsize=10)
caption(s,"χ² across the 11 model states (base + 10 iters), n=96 personas. W = Kendall's concordance (effect size). All p<0.001; W moderate (0.06–0.47).",1.7,6.75,9.9,align=PP_ALIGN.LEFT,size=10)

# A6 · slopes
s = slide(); title_bar(s,"Trajectory slopes & peak iterations (descriptive)",kicker="APPENDIX · slope_by_arm")
md_table(s, t("7_stats/slope_by_arm.md"), 1.7, 1.5, 9.9, 5.1,
   rename={"spearman_rho":"Spearman ρ","ols_slope":"OLS slope","peak_iter":"peak","final_iter":"final"}, fontsize=10)
caption(s,"ρ / slope are DESCRIPTIVE (pool repeated personas) — inference is carried by Friedman. Note PTO peaks at iter 10 on every metric; GRPO peaks earlier (8) on the halo rubrics.",1.7,6.75,9.9,align=PP_ALIGN.LEFT,size=10)

# A7 · PCA + iter9
s = slide(); title_bar(s,"Factor structure & GRPO iter-9 anomaly",kicker="APPENDIX · rubric_pca / grpo_iter9")
tb=box(s,0.6,1.35,5.0,0.5); set_runs(tb.text_frame.paragraphs[0],[("PC1 share (10-metric space)",True,NAVY)],13)
md_table(s, t("7_stats/rubric_pca_pc1.md"), 0.6, 1.8, 4.6, 1.3,
   rename={"PC1_pct":"PC1 %"}, fontsize=11)
caption(s,"vs ≈91% with the halo rubrics alone → a real 2nd factor.",0.6,3.25,4.8,align=PP_ALIGN.LEFT,size=10)
tb=box(s,6.0,1.35,6.5,0.5); set_runs(tb.text_frame.paragraphs[0],[("GRPO one-iteration paired dips",True,NAVY)],13)
md_table(s, t("7_stats/grpo_iter9_check.md"), 6.0, 1.8, 6.7, 3.6,
   rename={"mean_delta":"Δ","p_holm":"Holm"}, fontsize=10)
caption(s,"The regression is a step-down at it8→it9 (Q1+Q2 dz −0.41, Holm 0); it9→it10 is n.s. — not a slow decay.",6.0,5.55,6.7,align=PP_ALIGN.LEFT,size=10)

# A8 · MITI threshold verdicts
s = slide(); title_bar(s,"MITI 4.2.1 threshold verdicts — base vs final, per arm",kicker="APPENDIX · miti_thresholds")
md_table(s, t("2_questionnaires/miti_threshold_verdicts.md"), 1.2, 1.5, 10.9, 3.6, fontsize=10.5)
caption(s,"✓good / ✓fair / ✗ = the manual's clinician competency verdicts per summary score. Neither arm reaches \"good\" on the technique ratios (R:Q, %CR); the global ratings fare better.",1.2,6.0,10.9,align=PP_ALIGN.LEFT,size=10)

prs.save(OUT)
print("SAVED", OUT, "slides:", len(prs.slides._sldIdLst))
