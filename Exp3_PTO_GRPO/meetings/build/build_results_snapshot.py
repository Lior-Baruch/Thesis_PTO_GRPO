"""Build the LEAN results-snapshot deck for Exp3 — the attachment for the supervisor
publication-meeting email (2026-07-26).

Deliberately NOT the full supervisor deck (`build_supervisor_deck.py`): this one is
**results only, minimum interpretation** — what was run, what the numbers are, and
exactly how far the K=5 arms got. No recommendations, no next-steps, no threat model;
those live in the email body and the meeting itself.

Captions state what is plotted and the value. They do not draw conclusions.

Exports .pptx, then converts to .pdf via PowerPoint COM (see export_pdf.ps1).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))       # meetings/build/
ROOT = os.path.dirname(os.path.dirname(HERE))           # Exp3_PTO_GRPO/
L0F  = os.path.join(ROOT, "eda", "results", "L0", "figures")
L0T  = os.path.join(ROOT, "eda", "results", "L0", "tables")
L5F  = os.path.join(ROOT, "eda", "results", "L5", "figures")
OUT  = os.path.join(ROOT, "meetings", "2026-07-26", "results_snapshot_2026-07-26.pptx")

# Since 2026-07-28 every grader has its own leaf: results/<view>/figures/<family>/<judge>/<name>.
# The decks show the PRIMARY grader's figures, so the judge segment is injected here rather than
# spelled out at ~40 call sites. Point JUDGE at another grader to rebuild the same deck off it.
JUDGE = "gpt-4o-mini"

def _jp(base, p):
    """``<base>/<family>[/<sub>]/<JUDGE>/<name>`` — judge goes ahead of the FILENAME, not the family."""
    *parts, name = p.split("/")
    return os.path.join(base, *parts, JUDGE, name)

def f0(p): return _jp(L0F, p)
def t0(p): return _jp(L0T, p)
def f5(p): return _jp(L5F, p)

NAVY  = RGBColor(0x1F,0x3A,0x5F); PTO = RGBColor(0x00,0x72,0xB2)
GRPO  = RGBColor(0xE6,0x9F,0x00); GREY= RGBColor(0x5A,0x5A,0x5A)
LIGHT = RGBColor(0xEF,0xF2,0xF6); GREEN=RGBColor(0x00,0x8A,0x63)
RED   = RGBColor(0xC0,0x4A,0x1A); WHITE=RGBColor(0xFF,0xFF,0xFF)
DARK  = RGBColor(0x22,0x28,0x30); ROWALT=RGBColor(0xF4,0xF7,0xFA)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_N = 0

def slide(fnum=True):
    global _N
    s = prs.slides.add_slide(BLANK)
    if fnum:
        _N += 1
        tb = s.shapes.add_textbox(Inches(11.6), Inches(7.05), Inches(1.6), Inches(0.35))
        p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"Exp3 · {_N}"; r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0xB0,0xB8,0xC0); r.font.name = FONT
    return s

def box(s,l,tp,w,h):
    tb = s.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tb.text_frame.word_wrap = True; return tb

def rect(s,l,tp,w,h,color,line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(tp),Inches(w),Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False; return sp

def set_runs(para, segs, size, align=PP_ALIGN.LEFT):
    para.alignment = align
    for txt,bold,color in segs:
        r = para.add_run(); r.text = txt; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color; r.font.name = FONT

def title_bar(s, text, kicker=None):
    rect(s,0,0,13.333,1.15,NAVY); rect(s,0,1.15,13.333,0.06,PTO)
    tb = box(s,0.55,0.18,12.2,0.9); tf = tb.text_frame
    p = tf.paragraphs[0]
    if kicker:
        set_runs(p,[(kicker,True,RGBColor(0x8F,0xC7,0xEC))],13)
        set_runs(tf.add_paragraph(),[(text,True,WHITE)],26)
    else:
        set_runs(p,[(text,True,WHITE)],27)

def bullets(s, items, left=0.6, top=1.5, width=12.1, size=16, gap=8):
    tb = box(s,left,top,width,5.6); tf = tb.text_frame
    first = True
    for segs,level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.level = level; p.space_before = Pt(gap); p.space_after = Pt(2)
        mk = "•  " if level==0 else "–  "
        set_runs(p,[("    "*level+mk, False, PTO if level==0 else GREY)]+segs, size)

def fit(img,mw,mh):
    with Image.open(img) as im: iw,ih = im.size
    ar = iw/ih; w=mw; h=w/ar
    if h>mh: h=mh; w=h*ar
    return w,h

def figure(s,img,mw,mh,left,top,anchor="center"):
    """Place img fitted inside the (mw × mh) band at (left, top). Returns the image's bottom edge."""
    w,h = fit(img,mw,mh)
    dy = 0.0 if anchor=="top" else (mh-h)/2
    s.shapes.add_picture(img,Inches(left+(mw-w)/2),Inches(top+dy),Inches(w),Inches(h))
    return top+dy+h

def caption(s,text,left,top,width,color=GREY,size=11.5,align=PP_ALIGN.CENTER):
    p = box(s,left,top,width,0.5).text_frame.paragraphs[0]
    set_runs(p,[(text,False,color)],size,align)

# ---- native pptx table from a markdown table file ----
def _num(x):
    try: float(x); return True
    except: return False

def _fmt(col,v):
    if not _num(v): return v
    fv = float(v)
    if col in ("p","p_holm","wilcoxon_p") and fv == 0: return "<.001"
    if col in ("iteration","target_iter"): return f"{fv:.0f}"
    s = f"{fv:.3f}".rstrip("0").rstrip(".")
    return s if s else "0"

def md_table(s, md_path, left, top, width, height, drop=(), keep=None,
             fontsize=9.0, rename=None, order=None):
    raw = [l for l in open(md_path,encoding="utf-8").read().splitlines() if l.strip().startswith("|")]
    rows = [[c.strip() for c in l.strip().strip("|").split("|")] for l in raw]
    header, body = rows[0], rows[2:]
    if keep: body = [r for r in body if keep(dict(zip(header,r)))]
    if order: body = sorted(body, key=lambda r: order(dict(zip(header,r))))
    idx = [i for i,h in enumerate(header) if h not in drop]
    disp = [(rename or {}).get(header[i],header[i]) for i in idx]
    data = [[_fmt(header[i], r[i]) for i in idx] for r in body]
    nr, nc = len(data)+1, len(idx)
    gt = s.shapes.add_table(nr,nc,Inches(left),Inches(top),Inches(width),Inches(height)).table
    lens = [max(len(disp[c]), *(len(data[r][c]) for r in range(len(data)))) for c in range(nc)]
    tot = sum(lens)
    for c in range(nc):
        gt.columns[c].width = Inches(width*lens[c]/tot)
    for c,htext in enumerate(disp):
        cell = gt.cell(0,c); cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.margin_top=Pt(1);cell.margin_bottom=Pt(1);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
        pr = cell.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.LEFT
        r = pr.add_run(); r.text = htext; r.font.size=Pt(fontsize); r.font.bold=True
        r.font.color.rgb = WHITE; r.font.name = FONT
    for ri,row in enumerate(data,1):
        for c,val in enumerate(row):
            cell = gt.cell(ri,c); cell.fill.solid()
            cell.fill.fore_color.rgb = ROWALT if ri%2 else WHITE
            cell.margin_top=Pt(0);cell.margin_bottom=Pt(0);cell.margin_left=Pt(4);cell.margin_right=Pt(4)
            pr = cell.text_frame.paragraphs[0]
            col = DARK
            if disp[c] in ("dz","delta","mean_delta") and val.startswith("-"): col = RED
            rr = pr.add_run(); rr.text = val; rr.font.size=Pt(fontsize)
            rr.font.color.rgb = col; rr.font.name = FONT
    return gt

def grid_table(s, rows, left, top, width, colw, fontsize=12.5, rowh=0.42):
    """rows = list of (cells:list[str], color_or_None, is_header:bool)"""
    y = top
    for cells, color, is_hdr in rows:
        x = left
        for i,c in enumerate(cells):
            w = width*colw[i]
            rect(s,x,y,w,rowh, NAVY if is_hdr else (LIGHT if (rows.index((cells,color,is_hdr))%2==0) else WHITE))
            p = box(s,x+0.1,y+0.04,w-0.15,rowh).text_frame.paragraphs[0]
            col = WHITE if is_hdr else (color if (color and i==len(cells)-1) else DARK)
            set_runs(p,[(c,is_hdr or i==0,col)],fontsize)
            x += w
        y += rowh + 0.04
    return y

# =====================================================================
# 1 · TITLE
# =====================================================================
s = slide(fnum=False)
rect(s,0,0,13.333,7.5,NAVY); rect(s,0,4.6,13.333,0.08,PTO)
tb = box(s,0.9,1.6,11.5,2.6); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Looking Ahead in Goal-Oriented Dialogue",True,WHITE)],34)
set_runs(tf.add_paragraph(),[("Preference-Tree (PTO) vs Group-Relative (GRPO) optimization "
                             "of a small Motivational-Interviewing therapist",False,
                             RGBColor(0xBD,0xD6,0xEA))],18)
tb2 = box(s,0.9,4.85,11.5,1.9); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("Results snapshot — Exp3",True,WHITE)],20)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Numbers and figures only; interpretation deliberately left for the meeting.",
             False,RGBColor(0x9F,0xB4,0xC8))],14)
p = tf2.add_paragraph(); p.space_before = Pt(10)
set_runs(p,[("Lior Baruch · Reichman University · 26 July 2026",False,RGBColor(0x9F,0xB4,0xC8))],13)
p = tf2.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Llama-3.2-1B therapist (bf16)  ·  gpt-4o-mini simulated patient + oracle  ·  "
             "96 personas, persona-paired",False,RGBColor(0x7E,0x93,0xA8))],12)

# =====================================================================
# 2 · REMINDER — THE PUBLISHED PAPER (Exp1)
# =====================================================================
s = slide(); title_bar(s,"Reminder — what we published (ICLR 2025)","BACKGROUND · 1")

rect(s,0.6,1.45,12.1,1.12,LIGHT)
tb = box(s,0.85,1.55,11.6,1.0); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Preference Tree Optimization: Enhancing Goal-Oriented Dialogue with "
                           "Look-Ahead Simulations",True,NAVY)],17)
p = tf.add_paragraph(); p.space_before = Pt(4)
set_runs(p,[("Baruch, Butman, Bar, Friedman · ICLR 2025 · attached to this email",False,GREY)],13)

bullets(s,[
  ([("The method. ",True,NAVY),("At each therapist turn, branch several candidate replies; for each "
    "candidate simulate ",False,DARK),("K further turns",True,DARK),(" of the conversation; let the "
    "oracle score the resulting trajectory; keep the best and worst as a preference pair; update "
    "with DPO. Repeat, regenerating the tree from the improved model each round.",False,DARK)],0),
  ([("Why look ahead. ",True,NAVY),("Scoring a reply on its own rewards openings that look good in "
    "isolation. Scoring the reply plus K simulated turns rewards openings that lead somewhere.",
    False,DARK)],0),
  ([("Setup. ",True,NAVY),("Llama-2-7B therapist, GPT-3.5 as both simulated patient and oracle, "
    "96 patient personas, K ∈ {0, 5}, 7 iterations, reward = mean(Q1, Q2).",False,DARK)],0),
  ([("Headline finding. ",True,NAVY),("Every PTO model beat the untrained baseline on session "
    "satisfaction and working alliance, and ",False,DARK),("K = 5 gave higher and more stable "
    "scores than K = 0",True,DARK),(".",False,DARK)],0),
], top=2.85, size=15, gap=15)

# =====================================================================
# 3 · WHAT CHANGED SINCE
# =====================================================================
s = slide(); title_bar(s,"What changed since the paper","BACKGROUND · 2")

grid_table(s,[
  (["","Exp1 — ICLR paper","Exp2","Exp3 — this snapshot"],None,True),
  (["Therapist","Llama-2-7B","Llama-3.2-1B, 4-bit","Llama-3.2-1B, bf16"],None,False),
  (["Patient + oracle","GPT-3.5","gpt-4o-mini","gpt-4o-mini"],None,False),
  (["Patient personas","cooperative","less cooperative","less cooperative"],None,False),
  (["Oracle output","regex-parsed","JSON schema","JSON schema"],None,False),
  (["Evaluation","Q1, Q2","6 questionnaires","8 metrics (see next slide)"],None,False),
  (["Methods","PTO","PTO (4 oracles) + weak GRPO","PTO vs GRPO, matched"],None,False),
  (["Iterations","7","7 per arm","10 per arm (K=0)"],None,False),
],0.6,1.45,12.1,[0.16,0.24,0.28,0.32],fontsize=12,rowh=0.36)

p = box(s,0.6,5.15,12.1,0.4).text_frame.paragraphs[0]
set_runs(p,[("Each experiment is a fresh re-implementation — no data carries over between them.",
             True,NAVY)],13.5)

rect(s,0.6,5.62,12.1,1.25,RGBColor(0xFC,0xF3,0xEC))
rect(s,0.6,5.62,0.06,1.25,RED)
tb = box(s,0.85,5.72,11.7,1.1); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("One caveat worth stating up front — ",True,RED),
                           ("Exp2 and Exp3 absolute scores are not on the same axis.",True,DARK)],13)
p = tf.add_paragraph(); p.space_before = Pt(3)
set_runs(p,[("Same therapist model, but Exp2 generated its conversations in 4-bit and Exp3 in bf16. "
             "4-bit produces far more degenerate, looping therapist turns (≈9.5% of turns vs "
             "≈0.3%), which the oracle floors — so Exp2's base model scores 2.38 and Exp3's "
             "scores 3.00 for the same model. The non-degenerate Exp2 subset scores ≈2.93. "
             "Compare within an experiment, not across.",False,DARK)],12)

# =====================================================================
# 4 · WHAT WAS RUN
# =====================================================================
s = slide(); title_bar(s,"What was run in this experiment","SETUP")

bullets(s,[
  ([("Task. ",True,NAVY),("A 1B therapist model converses with a simulated patient (96 fixed "
    "personas). A larger oracle model then grades the full conversation on validated MI "
    "questionnaires.",False,DARK)],0),
  ([("Training signal. ",True,NAVY),("Oracle score on Q1+Q2 only. All other questionnaires are "
    "held out for evaluation.",False,DARK)],0),
  ([("Methods compared. ",True,NAVY),("PTO",True,PTO),(" (branch each therapist turn → oracle → "
    "preference pairs → DPO update) vs ",False,DARK),("GRPO",True,GRPO),
    (" (group-relative update over the same generations). Matched hyper-parameters, matched "
     "look-ahead K, matched min-conversation-length filter.",False,DARK)],0),
  ([("Loop. ",True,NAVY),("Each iteration regenerates all conversations from the current policy; "
    "those same conversations are the evaluation set.",False,DARK)],0),
  ([("Evaluation. ",True,NAVY),("Full-conversation scores on eight metrics: the five questionnaire "
    "rubrics ",False,DARK),("Q1+Q2, WAI-SR, CSQ-8, MI-SAT, MITI",True,DARK),(", plus ",False,DARK),
    ("PCT",True,DARK),(" (patient change-talk), ",False,DARK),("MICI",True,DARK),
    (" (MI-inconsistent therapist behaviour, lower is better), and the MITI ratios ",False,DARK),
    ("R:Q / %CR / %MICO",True,DARK),(" derived from the behaviour counts.",False,DARK)],0),
], top=1.45, size=15.5, gap=13)

rect(s,0.6,5.35,12.1,0.03,RGBColor(0xD5,0xDD,0xE5))
grid_table(s,[
  (["Arm","Iterations trained","Iterations scored","Status"],None,True),
  (["PTO   K=0","1–10","base + 1–10","complete"],GREEN,False),
  (["GRPO  K=0","1–10","base + 1–10","complete"],GREEN,False),
  (["PTO   K=5","1–5","base + 1–4","paused (iter-5 adapter unscored)"],RED,False),
  (["GRPO  K=5","1","base + 1","paused"],RED,False),
],0.6,5.55,12.1,[0.16,0.19,0.19,0.46],fontsize=12,rowh=0.33)

# =====================================================================
# 3 · MAIN FIGURE — trajectories
# =====================================================================
s = slide(); title_bar(s,"All metrics across 10 iterations (K = 0)","RESULTS · 1")
figure(s,f0("0_headline/trajectories_all_metrics.png"),12.3,5.15,0.5,1.35)
caption(s,"Mean over the 96 personas at each iteration; shaded band = 95% CI. "
          "Iteration 0 = untrained base model. MICI is inverted in meaning (lower = fewer "
          "MI-inconsistent therapist behaviours).",0.5,6.65,12.3)

# =====================================================================
# 4 · ENDPOINT NUMBERS
# =====================================================================
s = slide(); title_bar(s,"Endpoint numbers (K = 0)","RESULTS · 2")
caption(s,"Two endpoints reported for each arm: the matched final iteration (10) and each arm's "
          "own best iteration on its training rubric (Q1+Q2).",
        0.5,1.32,12.3,size=12.5,align=PP_ALIGN.LEFT)
md_table(s,t0("0_headline/leaderboard_scorecard.md"),0.5,1.85,12.3,1.9,
         fontsize=10.5,
         rename={"target":"endpoint","iteration":"iter",
                 "WAI-SR (Working Alliance)":"WAI-SR","CSQ-8 (Client Satisfaction)":"CSQ-8",
                 "MI-SAT (MI Satisfaction)":"MI-SAT","MITI (MI Integrity)":"MITI",
                 "PCT (Patient Change-Talk)":"PCT","MICI (MI-Inconsistency) ↓":"MICI ↓",
                 "Reflection:Question (MITI)":"R:Q","% Complex Reflections (MITI)":"%CR",
                 "% MI-Consistent (MITI)":"%MICO"})

rect(s,0.5,4.1,12.3,0.03,RGBColor(0xD5,0xDD,0xE5))
p = box(s,0.5,4.25,12.3,0.4).text_frame.paragraphs[0]
set_runs(p,[("Paired contrasts on the 96 shared personas (PTO − GRPO), Holm-corrected",True,NAVY)],14)
grid_table(s,[
  (["Contrast","Q1+Q2 Δ","dz","p (Holm)"],None,True),
  (["PTO iter-10  vs  GRPO iter-10  (matched final)","+0.51","0.73","< .001"],None,False),
  (["PTO iter-10  vs  GRPO iter-8  (each at its own best)","+0.18","0.30",".010"],None,False),
],0.5,4.75,8.6,[0.55,0.16,0.13,0.16],fontsize=12,rowh=0.36)

p = box(s,0.5,6.05,12.3,1.0).text_frame.paragraphs[0]
set_runs(p,[("Trajectory shape: ",True,NAVY),
            ("PTO Q1+Q2 rises 3.00 → 4.26 with its maximum at the last iteration (OLS slope "
             "0.120/iter). GRPO rises 3.07 → 4.08 by iteration 8, then falls to 3.75 by "
             "iteration 10 (slope 0.072/iter). Every arm × rubric effect vs base is significant "
             "at Holm p < .001.",False,DARK)],12.5)

# =====================================================================
# 5 · EFFECT VS BASE
# =====================================================================
s = slide(); title_bar(s,"Effect vs the untrained base model (K = 0)","RESULTS · 3")
figure(s,f0("0_headline/effect_vs_base_forest_final.png"),6.05,4.9,0.35,1.45)
figure(s,f0("0_headline/effect_vs_base_forest_best.png"),6.05,4.9,6.75,1.45)
caption(s,"at the matched final iteration (10)",0.35,6.45,6.05,color=NAVY,size=12.5)
caption(s,"at each arm's own best iteration (PTO 10, GRPO 8)",6.75,6.45,6.05,color=NAVY,size=12.5)
caption(s,"Paired Cohen's dz vs base over the 96 personas, with 95% CI. Positive = improvement, "
          "except MICI where positive = more MI-inconsistent behaviour.",0.35,6.85,12.45,size=11.5)

# =====================================================================
# 6 · WHAT MOVED ALONGSIDE THE SCORES
# =====================================================================
s = slide(); title_bar(s,"What else moved as the scores rose (K = 0)","RESULTS · 4")
bot = figure(s,f0("0_headline/reward_hack_panel.png"),7.85,5.15,0.35,1.5,anchor="top")
caption(s,"Global rubric score (left axis) vs MI-inconsistent behaviour and patient change-talk "
          "(right axis), per iteration",0.35,bot+0.1,7.85,size=11.5)

px = 8.35
p = box(s,px,1.5,4.6,0.5).text_frame.paragraphs[0]
set_runs(p,[("Measured alongside the gains",True,NAVY)],14.5)
tb = box(s,px,2.0,4.6,4.6); tf = tb.text_frame
first = True
for segs in [
  [("MICI",True,DARK),(" (MI-inconsistent behaviour) rises 0.21 at base → ",False,DARK),
   ("0.49",True,PTO),(" PTO and ",False,DARK),("0.84",True,GRPO),(" GRPO at iter 10 "
   "(GRPO 0.54 at its iter-8 peak).",False,DARK)],
  [("Reflection-to-question ratio at iter 10: PTO 0.75, GRPO 1.44.",False,DARK)],
  [("The 5 questionnaire rubrics load on a single factor (PC1 ≈ 91%). Adding the other three "
    "metrics drops PC1 to ≈ 55% — but that separation comes from ",False,DARK),
   ("MICI and the MITI ratios",True,DARK),(". PCT correlates with the 5 rubrics (ρ 0.79–0.94), so "
   "it is not measuring something they miss.",False,DARK)],
  [("Both arms drift toward affirmation-heavy language, more so in GRPO's late iterations.",
    False,DARK)],
  [("Full per-item decomposition of all eight scored instruments is on disk (not shown here).",
    False,GREY)],
]:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
    p.space_before = Pt(11)
    set_runs(p,[("•  ",False,PTO)]+segs,12.5)

# =====================================================================
# 7 · MICI DETAIL
# =====================================================================
s = slide(); title_bar(s,"MI-inconsistency, decomposed (K = 0)","RESULTS · 5")
figure(s,f0("0_headline/mici_detail_grid.png"),12.3,5.15,0.5,1.35)
caption(s,"Per-conversation rate of each MI-inconsistent therapist behaviour, by iteration. "
          "Lower is better throughout.",0.5,6.65,12.3)

# =====================================================================
# 8 · K=5 — WHAT EXISTS
# =====================================================================
s = slide(); title_bar(s,"Look-ahead (K = 5): exactly what exists today","RESULTS · 6")

grid_table(s,[
  (["Arm","Trained","Scored","Missing"],None,True),
  (["PTO  K=5","iters 1–5","base + iters 1–4","iter-5 eval conversations never generated"],RED,False),
  (["GRPO K=5","iter 1","base + iter 1","iters 2–10"],RED,False),
],0.5,1.4,12.3,[0.13,0.15,0.20,0.52],fontsize=12,rowh=0.34)

bot5 = figure(s,f5("1_outcomes/trajectories_all_metrics.png"),7.5,3.95,0.35,2.65,anchor="top")
px = 8.2
tb = box(s,px,2.6,4.75,4.2); tf = tb.text_frame
set_runs(tf.paragraphs[0],[("Numbers so far (PTO K=5 only)",True,NAVY)],14)
for segs in [
  [("Q1+Q2 3.00 → 3.89 over 4 iterations (dz 0.88).",False,DARK)],
  [("MICI 0.18 → 0.33 over the same 4 iterations.",False,DARK)],
  [("K=0 vs K=5, paired at matched early iterations (PTO): no significant difference on any "
    "metric (all Holm p > .5).",False,DARK)],
  [("GRPO K=5 has one trained iteration — not comparable.",False,DARK)],
  [("Both K=5 arms were paused on API cost, not on any result.",False,GREY)],
]:
    p = tf.add_paragraph(); p.space_before = Pt(10)
    set_runs(p,[("•  ",False,PTO)]+segs,12)

rect(s,8.2,6.15,4.75,0.9,LIGHT)
tb2 = box(s,8.35,6.22,4.5,0.8); tf2 = tb2.text_frame
set_runs(tf2.paragraphs[0],[("Against the ICLR paper: ",True,NAVY),
                            ("there K=5 beat K=0 on Llama-2-7B over 7 iterations. Here the K "
                             "comparison rests on 4 iterations of one method and shows nothing "
                             "yet.",False,DARK)],11)
caption(s,"K=5 trajectories — 4 iterations for PTO, 1 for GRPO. Read as preliminary.",
        0.35,bot5+0.05,7.5,size=11)

# =====================================================================
# 9 · WHAT'S ON DISK
# =====================================================================
s = slide(); title_bar(s,"What exists behind these slides","REFERENCE")
bullets(s,[
  ([("2,784 scored conversations",True,NAVY),(" across 29 model checkpoints; every conversation "
    "graded on all eight instruments by the oracle.",False,DARK)],0),
  ([("Per-candidate generation records",True,NAVY),(" for every training iteration — all M/G "
    "candidates with their scores and sub-scores.",False,DARK)],0),
  ([("Three analysis views",True,NAVY),(" (K=0 only, K=5 only, combined), each with figures, "
    "tables and a written summary, regenerated from one command.",False,DARK)],0),
  ([("Per-item decomposition",True,NAVY),(" of every instrument — which individual items move, "
    "at both endpoints — already computed, no further oracle calls needed.",False,DARK)],0),
  ([("Persona-level heterogeneity",True,NAVY),(" splits (patient cooperation level, problem "
    "type) for every metric.",False,DARK)],0),
  ([("Second-judge reliability pipeline",True,NAVY),(" is built and ready to run "
    "(inter-rater agreement + a pluggable second grader model); not yet executed.",False,DARK)],0),
], top=1.6, size=15.5, gap=15)

prs.save(OUT)
print("wrote", OUT, f"({_N+1} slides)")
